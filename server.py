
import os
import json
import uuid
import time
import re
import asyncio
import hmac
import html
import ipaddress
import socket
from io import BytesIO
from pathlib import Path
from datetime import datetime, timezone
from urllib.parse import urlparse, quote

from fastapi import FastAPI, UploadFile, File, Form, Header, HTTPException
from fastapi.responses import HTMLResponse, FileResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
import httpx
from bs4 import BeautifulSoup
from search_strategy import (
    _SEARCH_TOPIC_LABELS, _SPAM_DOMAINS,
    _WECHAT_MAX_ARTICLES, _WECHAT_MAX_FULLTEXT,
    _SEARCH_COVERAGE_PATTERNS,
    _assess_search_coverage, _build_baidu_search_plan,
    _build_gap_search_plan, _build_quick_search_plan,
    _build_screening_search_plan, _company_reference_aliases,
    _coverage_text_from_results, _dedupe_wechat_results,
    _filter_relevant_wechat_results, _filter_verified_search_results,
    _govern_search_result_groups, _has_company_reference, _search_result_quality_score,
    _wechat_supplement_queries,
)


# ============ 配置 ============
DEEPSEEK_API_KEY = os.environ.get("DEEPSEEK_API_KEY", "")
DEEPSEEK_BASE_URL = os.environ.get("DEEPSEEK_BASE_URL", "https://api.deepseek.com")
DEEPSEEK_MODEL = os.environ.get("DEEPSEEK_MODEL", "deepseek-v4-pro")

KIMI_API_KEY = os.environ.get("KIMI_API_KEY", "sk-TScI4VoJMeuQcWNp5UIwdxYEnAqwa7YVJMeB1Cg8eRMvFZBH")
KIMI_BASE_URL = "https://api.moonshot.cn/v1"

AVAILABLE_MODELS = {
    "deepseek-v4-pro": {"provider": "deepseek", "label": "DeepSeek V4 Pro（默认）", "base_url": DEEPSEEK_BASE_URL},
    "deepseek-v4-flash": {"provider": "deepseek", "label": "DeepSeek V4 Flash（快速）", "base_url": DEEPSEEK_BASE_URL},
    "kimi-k2.6": {"provider": "kimi", "label": "Kimi K2.6（长上下文）", "base_url": KIMI_BASE_URL},
    "moonshot-v1-128k": {"provider": "kimi", "label": "Moonshot V1 128K（备选）", "base_url": KIMI_BASE_URL},
}

def _get_model_config(model_name: str) -> dict:
    """Resolve model config, falling back to default."""
    return AVAILABLE_MODELS.get(model_name, AVAILABLE_MODELS[DEEPSEEK_MODEL])

# AI 自主搜索工具定义
ANALYSIS_TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "search_web",
            "description": (
                "搜索互联网获取公司相关的公开信息。"
                "用于补充缺失的信息维度（客户、供应商、产品、市场地位、规模、财务数据、新闻等）。"
                "查询词应包含公司名称和具体关注点，用空格分隔关键词。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "搜索查询词，如 'XX公司 大客户 供应商 合作订单'"
                    }
                },
                "required": ["query"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "fetch_page",
            "description": (
                "获取指定URL的完整网页文本内容。"
                "当搜索结果摘要信息不足、需要核实具体数据或深入了解某篇文章时使用。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {
                        "type": "string",
                        "description": "要获取内容的网页URL"
                    }
                },
                "required": ["url"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "search_wechat",
            "description": (
                "搜索微信公众号文章。微信公众号是小微企业最重要的公开信息渠道，"
                "常包含产品介绍、客户案例、产能规模、行业动态等一手经营信息。"
                "查询词应包含公司名称，可搭配行业或具体关注点。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "微信搜索查询词，如 'XX公司 产品 客户'"
                    }
                },
                "required": ["query"]
            }
        }
    }
]



# 垃圾域名黑名单 — Tavily 对中国公司搜索时常返回爬虫/聚合站
_SPAM_DOMAINS = {
    "dayinmao.com",
    "bjx.com", "china-10.com",
}
# 微信公众号搜索参数
_WECHAT_MAX_ARTICLES = 6
_WECHAT_MAX_FULLTEXT = 3

# 搜索覆盖度评估模式
_SEARCH_COVERAGE_PATTERNS = {
    "identity": r"主营|业务|产品|生产|加工|制造|代工|OEM|ODM|公司简介",
    "operations": r"客户|供应商|供货|采购|订单|排产|产量|招标|合作",
    "scale": r"营收|收入|销售额|产能|产量|吨|员工|融资|亿元|万元",
    "ownership": r"股东.{0,30}(?:持股|出资|控股)|实际控制人|实控人",
    "risk": r"行政处罚|处罚|诉讼|失信|被执行|经营异常|监管|召回|通报",
    "market": r"市场|行业|竞争|份额|领先|排名|品牌|渠道|协同",
}

# 搜索主题中文标签
_SEARCH_TOPIC_LABELS = {
    "identity": "公司身份",
    "operations": "客户与经营",
    "scale": "规模与营收",
    "ownership": "股权与实控人",
    "risk": "风险合规",
    "market": "市场与竞争",
    "industry": "行业研究",
    "official": "官方网站",
}

# 企查查 MCP 配置
QCC_MCP_BASE = "https://agent.qcc.com/mcp"
QCC_MCP_KEY = os.environ.get("QCC_MCP_KEY", "Miyyy5X4Ik5nm2p7HdVRy0iahYM30NAE1g0qFbfDsahRRCeC")
_QCC_MCP_HEADERS = {"Authorization": f"Bearer {QCC_MCP_KEY}", "Content-Type": "application/json"}


def _parse_qcc_mcp_response(response_text: str) -> dict:
    """Parse QCC MCP's text/event-stream or JSON-RPC response body."""
    text = str(response_text or "").strip()
    if not text:
        return {}
    if text.startswith("event:") or "\ndata: " in text or text.startswith("data: "):
        payload = "\n".join(line[6:] for line in text.splitlines() if line.startswith("data: "))
    else:
        payload = text
    if not payload.strip():
        return {}
    return json.loads(payload)


def _qcc_content_text(data: dict) -> str:
    result = data.get("result", {}) if isinstance(data, dict) else {}
    if result.get("isError"):
        return ""
    content = result.get("content", [])
    if isinstance(content, str):
        return content
    return "\n".join(
        item.get("text", "") for item in content
        if isinstance(item, dict) and item.get("type") == "text"
    )


def _qcc_first_company_name(text: str, fallback: str) -> str:
    """Pick the exact registered name returned by get_company_by_query when available."""
    try:
        data = json.loads(text)
        company_name = data.get("企业信息", {}).get("企业名称") or data.get("企业名称")
        if company_name:
            return str(company_name).strip()
    except (TypeError, ValueError):
        pass
    for pattern in (
        r"(?:企业名称|公司名称|名称|name)[：:\s]+([^\n，,]+(?:有限公司|有限责任公司|股份有限公司|集团有限公司))",
        r"([\u4e00-\u9fff（）()A-Za-z0-9·]+(?:有限公司|有限责任公司|股份有限公司|集团有限公司))",
    ):
        match = re.search(pattern, text)
        if match:
            return match.group(1).strip()
    return fallback


def _qcc_extract_registration_info(text: str, company_name: str) -> dict:
    info = {"company_name": company_name}
    try:
        data = json.loads(text)
        if isinstance(data, dict):
            mapping = {
                "registered_capital": "注册资本",
                "established_at": "成立日期",
                "legal_representative": "法定代表人",
                "business_scope": "经营范围",
                "credit_code": "统一社会信用代码",
                "registration_status": "登记状态",
                "insured_count": "参保人数",
            }
            info["company_name"] = data.get("企业名称") or company_name
            for target, source in mapping.items():
                if data.get(source) not in {None, ""}:
                    info[target] = str(data[source]).strip()
            return info
    except (TypeError, ValueError):
        pass
    for field, pattern in [
        ("registered_capital", r"注册资本[：:]\s*([^\n]+)"),
        ("established_at", r"成立(?:日期|时间)[：:]\s*([^\n]+)"),
        ("legal_representative", r"法定代表[人]?[：:]\s*([^\n]+)"),
        ("business_scope", r"经营范围[：:]\s*([^\n]+)"),
        ("credit_code", r"统一社会信用代码[：:]\s*([^\n]+)"),
        ("registration_status", r"(?:登记状态|经营状态)[：:]\s*([^\n]+)"),
        ("insured_count", r"参保人数[：:]\s*([^\n]+)"),
    ]:
        match = re.search(pattern, text)
        if match:
            info[field] = match.group(1).strip().rstrip("。；;,")
    return info


async def _qcc_call_tool(client: httpx.AsyncClient, name: str, arguments: dict, request_id: int) -> dict:
    resp = await client.post(
        f"{QCC_MCP_BASE}/company/stream",
        headers=_QCC_MCP_HEADERS,
        json={
            "jsonrpc": "2.0",
            "id": request_id,
            "method": "tools/call",
            "params": {"name": name, "arguments": arguments},
        },
    )
    if resp.status_code != 200:
        print(f"[QCC] {name} failed: {resp.status_code}")
        return {}
    data = _parse_qcc_mcp_response(resp.text)
    if data.get("error"):
        print(f"[QCC] {name} error: {data['error'].get('message', '')}")
        return {}
    return data


async def query_qcc_company_info(company_name: str) -> dict:
    """通过企查查 MCP 查询公司工商信息。返回 {注册资本, 成立时间, 法人, 经营范围, ...}"""
    if not QCC_MCP_KEY:
        return {}
    try:
        async with httpx.AsyncClient(timeout=httpx.Timeout(20.0, connect=8.0)) as client:
            entity_data = await _qcc_call_tool(
                client,
                "get_company_by_query",
                {"searchKey": company_name},
                1,
            )
            entity_text = _qcc_content_text(entity_data)
            exact_name = _qcc_first_company_name(entity_text, company_name)
            registration_data = await _qcc_call_tool(
                client,
                "get_company_registration_info",
                {"searchKey": exact_name},
                2,
            )
            registration_text = _qcc_content_text(registration_data)
            if not registration_text:
                return {}
            return _qcc_extract_registration_info(registration_text, exact_name)
    except Exception as e:
        print(f"[QCC] Error: {e}")
        return {}


def _qcc_registration_summary_text(info: dict) -> str:
    lines = ["## 企查查工商信息"]
    for label, key in [
        ("企业名称", "company_name"),
        ("统一社会信用代码", "credit_code"),
        ("法定代表人", "legal_representative"),
        ("登记状态", "registration_status"),
        ("成立日期", "established_at"),
        ("注册资本", "registered_capital"),
        ("企业类型", "company_type"),
        ("人员规模", "staff_size"),
        ("参保人数", "insured_count"),
        ("国标行业", "industry"),
        ("注册地址", "registered_address"),
        ("经营范围", "business_scope"),
    ]:
        if info.get(key):
            lines.append(f"{label}：{info[key]}")
    return "\n".join(lines)


async def collect_qcc_evidence(company_name: str) -> tuple[str, list[dict]]:
    """Collect structured QCC evidence for the report and AI prompt."""
    if not QCC_MCP_KEY:
        return "", []
    qcc_text_parts = []
    qcc_refs: list[dict] = []
    try:
        async with httpx.AsyncClient(timeout=httpx.Timeout(20.0, connect=8.0)) as client:
            entity_data = await _qcc_call_tool(
                client, "get_company_by_query", {"searchKey": company_name}, 1
            )
            entity_text = _qcc_content_text(entity_data)
            exact_name = _qcc_first_company_name(entity_text, company_name)
            if entity_text:
                qcc_text_parts.append("\n".join([
                    "## 企查查企业实体识别",
                    f"检索词：{company_name}",
                    f"命中主体：{exact_name}",
                ]))
                qcc_refs.append({"provider": "企查查MCP", "title": "企业实体识别", "url": "https://www.qcc.com/"})

            registration_data = await _qcc_call_tool(
                client, "get_company_registration_info", {"searchKey": exact_name}, 2
            )
            registration_text = _qcc_content_text(registration_data)
            registration_info = _qcc_extract_registration_info(registration_text, exact_name)
            if registration_info:
                qcc_text_parts.append(_qcc_registration_summary_text(registration_info))
                qcc_refs.append({"provider": "企查查MCP", "title": "企业工商信息", "url": "https://www.qcc.com/"})

            profile_data = await _qcc_call_tool(
                client, "get_company_profile", {"searchKey": exact_name}, 3
            )
            profile_text = _qcc_content_text(profile_data)
            profile_info = json.loads(profile_text) if profile_text else {}
            if isinstance(profile_info, dict) and profile_info:
                summary_lines = ["## 企查查企业简介"]
                for label, key in [
                    ("企业名称", "企业名称"),
                    ("简介", "简介"),
                    ("企业当前经营状态", "企业当前经营状态"),
                    ("所属行业", "企查查行业"),
                ]:
                    if profile_info.get(key):
                        summary_lines.append(f"{label}：{profile_info[key]}")
                qcc_text_parts.append("\n".join(summary_lines))
                qcc_refs.append({"provider": "企查查MCP", "title": "企业简介", "url": "https://www.qcc.com/"})

            controller_data = await _qcc_call_tool(
                client, "get_actual_controller", {"searchKey": exact_name}, 4
            )
            controller_text = _qcc_content_text(controller_data)
            controller_info = json.loads(controller_text) if controller_text else {}
            if isinstance(controller_info, dict) and controller_info:
                lines = ["## 企查查实际控制人"]
                for item in controller_info.get("实际控制人信息", [])[:3]:
                    if not isinstance(item, dict):
                        continue
                    name = item.get("实际控制人名称", "")
                    direct = item.get("直接持股比例", "")
                    vote = item.get("表决权比例", "")
                    lines.append(f"- {name}（直接持股{direct}，表决权{vote}）")
                qcc_text_parts.append("\n".join(lines))
                qcc_refs.append({"provider": "企查查MCP", "title": "实际控制人", "url": "https://www.qcc.com/"})

            shareholder_data = await _qcc_call_tool(
                client, "get_shareholder_info", {"searchKey": exact_name}, 5
            )
            shareholder_text = _qcc_content_text(shareholder_data)
            shareholder_info = json.loads(shareholder_text) if shareholder_text else {}
            if isinstance(shareholder_info, dict) and shareholder_info:
                lines = ["## 企查查股东信息"]
                for item in shareholder_info.get("股东信息", [])[:5]:
                    if not isinstance(item, dict):
                        continue
                    name = item.get("股东名称", "")
                    ratio = item.get("持股比例", "")
                    lines.append(f"- {name}（持股{ratio}）")
                qcc_text_parts.append("\n".join(lines))
                qcc_refs.append({"provider": "企查查MCP", "title": "股东信息", "url": "https://www.qcc.com/"})

            annual_data = await _qcc_call_tool(
                client, "get_annual_reports", {"searchKey": exact_name}, 6
            )
            annual_text = _qcc_content_text(annual_data)
            annual_info = json.loads(annual_text) if annual_text else {}
            if isinstance(annual_info, dict):
                reports = annual_info.get("企业年报信息", [])
                if reports:
                    for report in [item for item in reports[:3] if isinstance(item, dict)]:
                        basic = report.get("企业基本信息", {}) if isinstance(report, dict) else {}
                        assets = report.get("企业资产状况信息", {}) if isinstance(report, dict) else {}
                        year = report.get("年报年度", "最新") if isinstance(report, dict) else "最新"
                        lines = [f"## 企查查年报（{year}）"]
                        if basic.get("企业主营业务活动"):
                            lines.append(f"主营业务活动：{basic['企业主营业务活动']}")
                        if basic.get("企业经营状态"):
                            lines.append(f"经营状态：{basic['企业经营状态']}")
                        if basic.get("从业人数"):
                            lines.append(f"从业人数：{basic['从业人数']}")
                        if basic.get("企业是否有网站或网店"):
                            lines.append(f"是否有网站/网店：{basic['企业是否有网站或网店']}")
                        if basic.get("企业联系电话"):
                            lines.append(f"联系电话：{basic['企业联系电话']}")
                        if assets.get("营业总收入"):
                            lines.append(f"营业总收入：{assets['营业总收入']}")
                        if assets.get("利润总额"):
                            lines.append(f"利润总额：{assets['利润总额']}")
                        qcc_text_parts.append("\n".join(lines))
                        qcc_refs.append({
                            "provider": "企查查MCP",
                            "title": f"企业年报（{year}）",
                            "url": "https://www.qcc.com/",
                        })
    except Exception as e:
        print(f"[QCC] Evidence collection error: {e}")
    return "\n\n".join(qcc_text_parts), qcc_refs


TAVILY_API_KEY = os.environ.get("TAVILY_API_KEY", "")
SEARCH_PROVIDER = os.environ.get("SEARCH_PROVIDER", "exa").strip().lower()
TAVILY_FALLBACK = os.environ.get("TAVILY_FALLBACK", "true").strip().lower() not in {"0", "false", "no", "off"}
EXA_MCP_URL = os.environ.get("EXA_MCP_URL", "https://mcp.exa.ai/mcp?tools=web_search_exa")
EXA_API_KEY = os.environ.get("EXA_API_KEY", "")
ADMIN_API_TOKEN = os.environ.get("ADMIN_API_TOKEN", "")

BASE_DIR = Path(__file__).parent
UPLOAD_DIR = BASE_DIR / "uploads"
REPORT_DIR = BASE_DIR / "reports"
UPLOAD_DIR.mkdir(exist_ok=True)
REPORT_DIR.mkdir(exist_ok=True)

app = FastAPI(title="PE Deep Screening", version="0.1.0")
app.mount("/static", StaticFiles(directory=str(BASE_DIR / "static")), name="static")

# ============ 任务存储 ============
tasks: dict = {}


def _require_admin_token(provided_token: str) -> None:
    if not ADMIN_API_TOKEN or not hmac.compare_digest(provided_token, ADMIN_API_TOKEN):
        raise HTTPException(403, "管理接口未授权")


def _is_blocked_network_address(address: str) -> bool:
    try:
        ip = ipaddress.ip_address(address)
    except ValueError:
        return True
    return (
        ip.is_private
        or ip.is_loopback
        or ip.is_link_local
        or ip.is_multicast
        or ip.is_reserved
        or ip.is_unspecified
    )


async def _is_safe_public_url(url: str, host_cache: dict[str, bool] | None = None) -> bool:
    """Allow web crawling only for publicly routable HTTP(S) hosts."""
    try:
        parsed = urlparse(url)
        if parsed.scheme not in ("http", "https") or not parsed.hostname:
            return False
        if parsed.username or parsed.password:
            return False
        host = parsed.hostname.rstrip(".").lower()
        if host in {"localhost", "localhost.localdomain"} or host.endswith(".local"):
            return False
        if host_cache is not None and host in host_cache:
            return host_cache[host]
        try:
            host_ip = ipaddress.ip_address(host)
        except ValueError:
            host_ip = None
        if host_ip is not None:
            safe = not _is_blocked_network_address(host)
        else:
            port = parsed.port or (443 if parsed.scheme == "https" else 80)
            infos = await asyncio.to_thread(socket.getaddrinfo, host, port, type=socket.SOCK_STREAM)
            safe = bool(infos) and all(not _is_blocked_network_address(info[4][0]) for info in infos)
        if host_cache is not None:
            host_cache[host] = safe
        return safe
    except (OSError, ValueError):
        return False


async def _protect_browser_context(context) -> None:
    """Block private-network navigation and subrequests, including redirects."""
    host_cache: dict[str, bool] = {}

    async def guard_route(route, request):
        parsed = urlparse(request.url)
        if parsed.scheme in ("http", "https") and not await _is_safe_public_url(request.url, host_cache):
            print(f"[Playwright] Blocked non-public request: {request.url[:120]}")
            await route.abort()
            return
        await route.continue_()

    await context.route("**/*", guard_route)


# ============ 公开网页搜索（Exa 优先，Tavily 可选兜底）============
# ============ Tavily 多 Key 轮换系统 ============
TAVILY_KEYS_FILE = BASE_DIR / "tavily_keys.json"
TAVILY_MONTHLY_LIMIT = 1000


def _load_keys() -> dict:
    if TAVILY_KEYS_FILE.exists():
        try:
            return json.loads(TAVILY_KEYS_FILE.read_text(encoding="utf-8"))
        except (json.JSONDecodeError, KeyError):
            pass
    # Initialize with current env key
    init_key = os.environ.get("TAVILY_API_KEY", "")
    if not init_key:
        init_key = TAVILY_API_KEY
    return {
        "keys": [{"key": init_key, "credits_used": 0, "label": "默认Key"}],
        "current_index": 0,
    }


def _save_keys(data: dict):
    TAVILY_KEYS_FILE.write_text(json.dumps(data, ensure_ascii=False), encoding="utf-8")


def _get_current_key() -> str:
    d = _load_keys()
    idx = d.get("current_index", 0)
    return d["keys"][idx]["key"]


def _track_credits(credits: int):
    d = _load_keys()
    idx = d.get("current_index", 0)
    d["keys"][idx]["credits_used"] = d["keys"][idx].get("credits_used", 0) + credits
    _save_keys(d)


def _mark_current_key_exhausted() -> dict:
    d = _load_keys()
    idx = d.get("current_index", 0)
    current = d["keys"][idx]
    limit = int(current.get("credits_limit") or TAVILY_MONTHLY_LIMIT)
    current["credits_used"] = max(int(current.get("credits_used", 0)), limit)
    current["credits_limit"] = limit
    _save_keys(d)
    return current


async def _refresh_tavily_usage() -> dict:
    """Read official Tavily plan usage; local response tracking is not reliable across API versions."""
    d = _load_keys()
    async with httpx.AsyncClient(timeout=httpx.Timeout(20.0, connect=8.0)) as client:
        for entry in d.get("keys", []):
            key = str(entry.get("key", "")).strip()
            if not key:
                continue
            try:
                resp = await client.get(
                    "https://api.tavily.com/usage",
                    headers={"Authorization": f"Bearer {key}"},
                )
                if resp.status_code != 200:
                    entry["usage_status"] = f"http_{resp.status_code}"
                    continue
                usage_data = resp.json()
                key_data = usage_data.get("key", {})
                account_data = usage_data.get("account", {})
                entry["credits_used"] = int(key_data.get("usage", entry.get("credits_used", 0)) or 0)
                entry["credits_limit"] = int(
                    key_data.get("limit") or account_data.get("plan_limit") or TAVILY_MONTHLY_LIMIT
                )
                entry["usage_status"] = "ok"
            except Exception as e:
                entry["usage_status"] = f"error:{type(e).__name__}"
    _save_keys(d)
    return d


def _rotate_key_if_needed(reserve: int = 0) -> bool:
    """Rotate to the next key that still has usable credits."""
    d = _load_keys()
    idx = d.get("current_index", 0)
    current = d["keys"][idx]
    current_limit = int(current.get("credits_limit") or TAVILY_MONTHLY_LIMIT)
    threshold = max(current_limit - reserve, 0)
    if current.get("credits_used", 0) >= threshold:
        for offset in range(1, len(d["keys"])):
            next_idx = (idx + offset) % len(d["keys"])
            next_key = d["keys"][next_idx]
            next_limit = int(next_key.get("credits_limit") or TAVILY_MONTHLY_LIMIT)
            if next_key.get("credits_used", 0) >= max(next_limit - reserve, 0):
                continue
            d["current_index"] = next_idx
            _save_keys(d)
            print(f"[Tavily] Key rotated: {idx} -> {next_idx} ({d['keys'][next_idx].get('label', '')})")
            return True
    return False


async def _send_credit_warning_email(key_label: str, credits_used: int):
    smtp_host = os.environ.get("SMTP_HOST", "")
    smtp_port = int(os.environ.get("SMTP_PORT", "587"))
    smtp_user = os.environ.get("SMTP_USER", "")
    smtp_pass = os.environ.get("SMTP_PASS", "")
    alert_email = os.environ.get("ALERT_EMAIL", "jefftang@mtxpartners.com")
    if not smtp_host:
        return
    import smtplib
    from email.mime.text import MIMEText
    pct = credits_used / TAVILY_MONTHLY_LIMIT * 100
    msg = MIMEText(f"Tavily Key「{key_label}」已使用 {credits_used}/{TAVILY_MONTHLY_LIMIT}（{pct:.0f}%）。系统将自动轮换到下一个Key。", "plain", "utf-8")
    msg["Subject"] = f"[PE Deep Screening] Tavily Key 额度将尽 — {key_label}"
    msg["From"] = smtp_user
    msg["To"] = alert_email
    try:
        if smtp_port == 465:
            server = smtplib.SMTP_SSL(smtp_host, smtp_port, timeout=15)
        else:
            server = smtplib.SMTP(smtp_host, smtp_port, timeout=15)
            server.starttls()
        server.login(smtp_user, smtp_pass)
        server.sendmail(smtp_user, [alert_email], msg.as_string())
        server.quit()
    except Exception as e:
        print(f"[Tavily] Email failed: {e}")


async def tavily_search(
    query: str, max_results: int = 10, search_depth: str = "advanced", _retry_on_key_error: bool = True
) -> list:
    """调用 Tavily 搜索 API"""
    # 提前切换接近耗尽的 Key，防止任务中途失败
    _rotate_key_if_needed(reserve=50)
    async with httpx.AsyncClient(timeout=httpx.Timeout(90.0, connect=15.0)) as client:
        resp = await client.post(
            "https://api.tavily.com/search",
            json={
                "api_key": _get_current_key(),
                "query": query,
                "max_results": max_results,
                "search_depth": search_depth,
                "include_answer": True,
                "include_raw_content": False,
            }
        )
        if resp.status_code == 432 and _retry_on_key_error:
            unavailable = _mark_current_key_exhausted()
            print(f"[Tavily] Key unavailable (HTTP 432): {unavailable.get('label', '')}")
            if _rotate_key_if_needed():
                return await tavily_search(query, max_results, search_depth, _retry_on_key_error=False)
        resp.raise_for_status()
        data = resp.json()

        # 多Key轮换追踪
        credits_used = data.get("usage", {}).get("credits", 0)
        if credits_used > 0:
            _track_credits(credits_used)
            d = _load_keys()
            idx = d["current_index"]
            current = d["keys"][idx]
            used = current.get("credits_used", 0)
            print(f"[Tavily] +{credits_used} credits, key[{idx}] '{current.get('label','')}': {used}/{TAVILY_MONTHLY_LIMIT}")
            if used >= int(current.get("credits_limit") or TAVILY_MONTHLY_LIMIT):
                _rotate_key_if_needed()
                asyncio.create_task(_send_credit_warning_email(current.get("label", ""), used))

        results = []
        # Tavily answer
        if data.get("answer"):
            results.append({"type": "answer", "content": data["answer"]})
        # Search results
        for r in data.get("results", []):
            results.append({
                "type": "result",
                "provider": "Tavily",
                "title": r.get("title", ""),
                "url": r.get("url", ""),
                "content": r.get("content", ""),
            })
        return results


def _parse_exa_search_text(text: str) -> list[dict]:
    """Convert Exa MCP text output into the same result shape used by the report."""
    results = []
    for block in re.split(r"\n\s*---\s*\n", text or ""):
        title_match = re.search(r"(?m)^Title:\s*(.+?)\s*$", block)
        url_match = re.search(r"(?m)^URL:\s*(https?://\S+)\s*$", block)
        if not title_match or not url_match:
            continue
        content = re.sub(r"(?m)^(Title|URL|Published|Author):\s*.*\n?", "", block).strip()
        content = re.sub(r"(?m)^Highlights:\s*", "", content).strip()
        results.append({
            "type": "result",
            "provider": "Exa",
            "title": title_match.group(1).strip(),
            "url": url_match.group(1).strip(),
            "content": content,
        })
    return results


async def exa_search(query: str, max_results: int = 10, _attempt: int = 0) -> list:
    """Search via Exa's hosted MCP endpoint. Retries on 429 with backoff."""
    headers = {"Accept": "application/json, text/event-stream", "Content-Type": "application/json"}
    if EXA_API_KEY:
        headers["x-api-key"] = EXA_API_KEY
    payload = {
        "jsonrpc": "2.0",
        "id": 1,
        "method": "tools/call",
        "params": {"name": "web_search_exa", "arguments": {"query": query, "numResults": max_results}},
    }
    async with httpx.AsyncClient(timeout=httpx.Timeout(45.0, connect=15.0)) as client:
        resp = await client.post(EXA_MCP_URL, headers=headers, json=payload)
        if resp.status_code == 429:
            print(f"[Exa] 429 rate limited, falling back immediately")
            resp.raise_for_status()
        resp.raise_for_status()
    response_text = resp.text
    if "text/event-stream" in resp.headers.get("content-type", ""):
        data_lines = [line[6:] for line in response_text.splitlines() if line.startswith("data: ")]
        response_text = "\n".join(data_lines)
    data = json.loads(response_text)
    result = data.get("result", {})
    if result.get("isError"):
        raise RuntimeError(result.get("content", [{}])[0].get("text", "Exa MCP search failed"))
    content = "\n".join(
        item.get("text", "") for item in result.get("content", []) if item.get("type") == "text"
    )
    return _parse_exa_search_text(content)


def _has_relevant_company_result(results: list[dict], company_name: str) -> bool:
    """Require Exa results to mention the target entity before trusting them for SME reports."""
    if not company_name:
        return any(item.get("type") == "result" for item in results)
    needles = _company_reference_aliases(company_name)
    for item in results:
        if item.get("type") != "result":
            continue
        haystack = re.sub(r"\s+", "", f"{item.get('title', '')} {item.get('content', '')}")
        if any(needle and needle in haystack for needle in needles):
            return True
    return False


async def public_web_search(
    query: str,
    company_name: str = "",
    max_results: int = 10,
    search_depth: str = "advanced",
    comprehensive: bool = False,
) -> list:
    """Search the open web: Exa (with auto-retry on 429), Tavily as fallback."""
    if SEARCH_PROVIDER in {"exa", "hybrid"}:
        try:
            exa_results = await exa_search(query, max_results=max_results)
            if _has_relevant_company_result(exa_results, company_name):
                return exa_results
            print(f"[Exa] No relevant company match for '{query}', trying Tavily fallback")
        except Exception as e:
            print(f"[Exa] Query '{query}' failed ({type(e).__name__}), trying Tavily fallback")

    if TAVILY_FALLBACK:
        try:
            return await tavily_search(query, max_results=max_results, search_depth=search_depth)
        except Exception as e:
            print(f"[Tavily] Fallback query '{query}' failed: {e}")

    return [{"type": "error", "content": "所有搜索渠道均失败"}]


async def multi_search(
    queries: list[str], company_name: str = "", comprehensive: bool = False
) -> dict[str, list]:
    """并行执行多个搜索查询"""
    semaphore = asyncio.Semaphore(3)  # 最多3个并发
    
    async def safe_search(q):
        async with semaphore:
            try:
                return await public_web_search(
                    q, company_name=company_name, comprehensive=comprehensive
                )
            except Exception as e:
                print(f"[Search] Query '{q}' failed: {e}")
                return [{"type": "error", "content": str(e)}]

    coros = {q: safe_search(q) for q in queries}
    # 并行执行
    gathered = await asyncio.gather(*coros.values())
    return dict(zip(coros.keys(), gathered))


# ============ 百度搜索（补充通用网页搜索对中文内容的覆盖不足）============
async def search_baidu(query: str, max_results: int = 8) -> list[dict]:
    """百度搜索，返回 [{title, url, content}, ...]"""
    if not query or not query.strip():
        return []
    encoded = quote(query.strip())
    url = f"https://www.baidu.com/s?wd={encoded}&rn={max_results}"
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
        "Accept": "text/html,application/xhtml+xml;q=0.9,*/*;q=0.8",
        "Accept-Language": "zh-CN,zh;q=0.9",
    }
    results = []
    try:
        async with httpx.AsyncClient(timeout=httpx.Timeout(15.0, connect=8.0)) as client:
            resp = await client.get(url, headers=headers, follow_redirects=True)
            if resp.status_code != 200:
                return []
            html = resp.text
        # Parse Baidu results - typical structure: <div class="result c-container"> or <h3><a>
        items = re.findall(r'<div[^>]*class="[^"]*result[^"]*c-container[^"]*"[^>]*>(.*?)</div>\s*(?:</div>)?', html, re.DOTALL)
        if not items:
            # Fallback: simpler pattern
            items = re.findall(r'<h3[^>]*>\s*<a[^>]*href="(https?://[^"]+)"[^>]*>(.*?)</a>', html, re.DOTALL)
            for url_match, title_html in items[:max_results]:
                title = re.sub(r'<[^>]+>', '', title_html).strip()
                if title:
                    results.append({"title": title, "url": url_match, "content": title})
            return results
        for item in items[:max_results]:
            title_match = re.search(r'<a[^>]*href="(https?://[^"]+)"[^>]*>(.*?)</a>', item, re.DOTALL)
            if not title_match:
                continue
            real_url = title_match.group(1)
            title = re.sub(r'<[^>]+>', '', title_match.group(2)).strip()
            # Extract snippet
            snippet = ""
            sn_match = re.search(r'<span[^>]*class="[^"]*content-right_[^"]*"[^>]*>(.*?)</span>', item, re.DOTALL)
            if not sn_match:
                sn_match = re.search(r'<div[^>]*class="[^"]*c-abstract[^"]*"[^>]*>(.*?)</div>', item, re.DOTALL)
            if sn_match:
                snippet = re.sub(r'<[^>]+>', '', sn_match.group(1)).strip()[:300]
            if title:
                results.append({"title": title, "url": real_url, "content": snippet or title})
        print(f"[Baidu] {query}: {len(results)} results")
    except Exception as e:
        print(f"[Baidu] Search failed for {query}: {type(e).__name__}")
    return results


# ============ 体系化检索规划与来源治理 ============















# ============ 微信文章搜索（次幂数据 CimiData 付费API）============
CIMIDATA_APP_ID = os.environ.get("CIMIDATA_APP_ID", "")
CIMIDATA_APP_SECRET = os.environ.get("CIMIDATA_APP_SECRET", "")
CIMIDATA_HOST = "https://www.cimidata.com"
_cimidata_token = None
_cimidata_token_time = 0


async def _get_cimidata_token() -> str:
    global _cimidata_token, _cimidata_token_time
    if not CIMIDATA_APP_ID or not CIMIDATA_APP_SECRET:
        return ""
    if _cimidata_token and time.time() - _cimidata_token_time < 3600:
        return _cimidata_token
    async with httpx.AsyncClient(timeout=httpx.Timeout(15.0)) as client:
        resp = await client.post(
            f"{CIMIDATA_HOST}/api/token",
            json={"app_id": CIMIDATA_APP_ID, "app_secret": CIMIDATA_APP_SECRET},
        )
        data = resp.json()
        token_data = data.get("data") if isinstance(data, dict) else None
        _cimidata_token = (
            token_data.get("access_token", "") if isinstance(token_data, dict) else ""
        ) or (data.get("access_token", "") if isinstance(data, dict) else "")
        if not _cimidata_token:
            code = data.get("code", "unknown") if isinstance(data, dict) else "invalid_response"
            message = data.get("message", data.get("msg", "")) if isinstance(data, dict) else ""
            print(f"[CimiData] Token error: code={code}, message={message}")
            return ""
        _cimidata_token_time = time.time()
        return _cimidata_token


async def search_wechat_articles(query: str, days: int = 30, max_results: int = 10) -> list[dict]:
    """搜索微信公众号文章，返回 [{title, url, digest, account, time}, ...]"""
    if not query or not query.strip():
        return []
    token = await _get_cimidata_token()
    if not token:
        return []
    results = []
    try:
        async with httpx.AsyncClient(timeout=httpx.Timeout(30.0)) as client:
            resp = await client.post(
                f"{CIMIDATA_HOST}/api/v3/articles/search",
                params={"access_token": token},
                json={"keyword": query.strip(), "page": 1},
            )
            data = resp.json()
            articles = data.get("data", {}).get("items", []) if isinstance(data.get("data"), dict) else []
            for art in articles[:max_results]:
                results.append({
                    "title": art.get("title", ""),
                    "url": art.get("content_url", art.get("url", "")),
                    "digest": art.get("digest", art.get("summary", ""))[:200],
                    "account": art.get("nickname", art.get("account_name", art.get("author", ""))),
                    "time": art.get("published_at", art.get("create_time", "")),
                })
        print(f"[WeChat] CimiData found {len(results)} articles for: {query}")
    except Exception as e:
        print(f"[WeChat] CimiData search failed: {type(e).__name__}: {e}")
    return results


async def fetch_wechat_article_content(url: str) -> dict:
    """获取微信文章正文"""
    token = await _get_cimidata_token()
    if not token:
        return {}
    try:
        async with httpx.AsyncClient(timeout=httpx.Timeout(30.0)) as client:
            resp = await client.post(
                f"{CIMIDATA_HOST}/api/v2/articles/detail",
                params={"access_token": token},
                json={"url": url},
            )
            data = resp.json()
            content = data.get("data", {}).get("content", "") if isinstance(data.get("data"), dict) else ""
            title = data.get("data", {}).get("title", "") if isinstance(data.get("data"), dict) else ""
            return {"title": title, "content": content[:3000], "url": url}
    except Exception as e:
        print(f"[WeChat] CimiData detail failed: {type(e).__name__}")
        return {}




















# ============ 微信文章搜索（旧搜狗实现，已弃用保留参考）============
async def search_wechat_articles_old(query: str, days: int = 30, max_results: int = 10) -> list[dict]:
    """搜索微信公众号文章。底层搜狗，不可靠时静默返回空。"""
    if not query or not query.strip():
        return []

    time_to = int(time.time())
    time_from = time_to - days * 86400
    encoded = quote(query.strip())

    user_agents = [
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
        "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
    ]
    ua = user_agents[int(time.time()) % len(user_agents)]

    results = []
    for page in range(1, 3):  # 尝试翻2页
        url = f"https://weixin.sogou.com/weixin?type=2&query={encoded}&page={page}&ie=utf8&ft={time_from}&et={time_to}"
        headers = {
            "User-Agent": ua,
            "Accept": "text/html,application/xhtml+xml;q=0.9,*/*;q=0.8",
            "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.5",
            "Referer": "https://weixin.sogou.com/",
            "DNT": "1",
        }
        try:
            async with httpx.AsyncClient(timeout=httpx.Timeout(30.0, connect=10.0)) as client:
                resp = await client.get(url, headers=headers, follow_redirects=True)
                if resp.status_code != 200:
                    break
                html = resp.text

            if "请输入验证码" in html or "antispider" in html.lower():
                if page == 1:
                    print(f"[WeChat] Captcha triggered for: {query}")
                break

            li_pattern = r'<li[^>]*id="sogou_vr_\d+_box_\d+"[^>]*>(.*?)</li>'
            items = re.findall(li_pattern, html, re.DOTALL)
            if not items:
                break

            for item in items:
                if len(results) >= max_results:
                    break
                result = {}
                link_match = re.search(r'href="(/link\?url=[^"]+)"', item)
                if link_match:
                    result["sogou_link"] = "https://weixin.sogou.com" + link_match.group(1).replace("&amp;", "&")

                title_match = re.search(r'<h3>\s*<a[^>]*>(.*?)</a>\s*</h3>', item, re.DOTALL)
                if title_match:
                    result["title"] = re.sub(r'<[^>]+>', '', title_match.group(1)).strip()

                digest_match = re.search(r'<p[^>]*class="txt-info"[^>]*>(.*?)</p>', item, re.DOTALL)
                if digest_match:
                    result["digest"] = re.sub(r'<[^>]+>', '', digest_match.group(1)).strip()[:200]

                account_match = re.search(r'<span[^>]*class="all-time-y2"[^>]*>([^<]+)</span>', item)
                if account_match:
                    result["account"] = account_match.group(1).strip()

                time_match = re.search(r"timeConvert\('(\d+)'\)", item)
                if time_match:
                    ts = int(time_match.group(1))
                    result["time"] = datetime.fromtimestamp(ts).strftime("%Y-%m-%d")

                if result.get("title"):
                    results.append(result)

            if len(results) >= max_results:
                break
            await asyncio.sleep(0.5)  # 翻页间隔，降低触发验证码概率

        except Exception as e:
            print(f"[WeChat] Search failed for {query} p{page}: {type(e).__name__}")
            break

    if results:
        print(f"[WeChat] {query}: {len(results)} articles")

    return results


async def fetch_wechat_article_content_old(sogou_link: str) -> dict:
    """通过搜狗跳转链接获取微信文章真实URL和原文内容"""
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
            "Referer": "https://weixin.sogou.com/",
        }
        async with httpx.AsyncClient(timeout=httpx.Timeout(30.0, connect=10.0)) as client:
            resp = await client.get(sogou_link, headers=headers, follow_redirects=False)
            wx_url = ""
            if resp.status_code == 302:
                wx_url = resp.headers.get("Location", "")
            else:
                html = resp.text
                url_parts = re.findall(r"url \+= '([^']+)'", html)
                if url_parts:
                    wx_url = "".join(url_parts).replace("@", "")
                else:
                    wx_match = re.search(r'var\s+url\s*=\s*["\']([^"\']+)["\']', html)
                    if wx_match:
                        wx_url = wx_match.group(1)

            if not wx_url:
                return {"error": "Cannot extract article URL"}

            article_resp = await client.get(wx_url, headers=headers)
            html = article_resp.text

        result = {"wx_url": wx_url}

        title_match = re.search(r'<h1[^>]*class="rich_media_title"[^>]*>(.*?)</h1>', html, re.DOTALL)
        if not title_match:
            title_match = re.search(r'<title>([^<]+)</title>', html)
        if title_match:
            result["title"] = re.sub(r'<[^>]+>', '', title_match.group(1)).strip()

        content_match = re.search(r'<div[^>]*id="js_content"[^>]*>(.*?)</div>\s*<script', html, re.DOTALL)
        if content_match:
            content = content_match.group(1)
            content = re.sub(r'<script[^>]*>.*?</script>', '', content, flags=re.DOTALL)
            content = re.sub(r'<style[^>]*>.*?</style>', '', content, flags=re.DOTALL)
            content = re.sub(r'<[^>]+>', '\n', content)
            content = re.sub(r'\n\s*\n', '\n', content)
            content = re.sub(r'&nbsp;|&lt;|&gt;', ' ', content)
            result["content"] = content.strip()[:3000]

        return result

    except Exception as e:
        return {"error": f"{type(e).__name__}: {e}"}


# ============ Playwright 网页抓取 ============

# 不需要抓取的域名（企查查等后面会用MCP直接拿结构化数据，抓页面反而浪费）
SKIP_DOMAINS = {
    "qcc.com", "tianyancha.com", "qixin.com", "aiqicha.baidu.com",
    "sbj.cnipa.gov.cn", "gsxt.gov.cn",  # 企查查/天眼查/启信宝/爱企查/商标/工商
    "login", "register", "captcha",  # 登录/注册/验证码页面
}

# 值得深度抓取的域名后缀/关键词（高质量行业信息源）
PRIORITY_DOMAINS = {
    "gov.cn", "org.cn",  # 政府/组织
    "36kr.com", "jiemian.com", "yicai.com", "caixin.com",  # 财经媒体
    "sohu.com", "163.com", "sina.com.cn", "qq.com",  # 门户网站
}


def _is_url_worth_crawling(url: str, company_website: str = "") -> bool:
    """判断一个URL是否值得用Playwright抓取全文
    
    策略：
    - 公司官网：最高优先级（如果提供了）
    - 企查查/天眼查等：跳过（后面用MCP拿结构化数据）
    - 纯搜索结果页/API页：跳过
    - 其他正常网页：值得抓取
    """
    try:
        parsed = urlparse(url)
        domain = parsed.netloc.lower()
        path = parsed.path.lower()
        
        # 跳过无意义的页面
        skip_patterns = ["/login", "/register", "/captcha", "/api/", ".pdf", ".jpg", ".png", ".gif"]
        if any(p in path for p in skip_patterns):
            return False
        
        # 跳过企查查等（后面用MCP）
        for skip in SKIP_DOMAINS:
            if skip in domain:
                return False
        
        # 公司官网最高优先
        if company_website:
            try:
                cw_parsed = urlparse(company_website if "://" in company_website else f"https://{company_website}")
                if cw_parsed.netloc.lower() in domain or domain in cw_parsed.netloc.lower():
                    return True  # 官网，必抓
            except:
                pass
        
        # 其他正常页面，值得抓取
        return True
    except:
        return False


def _extract_main_content(html: str, url: str = "") -> str:
    """从HTML中提取正文文本，去掉导航/广告/页脚等噪音
    
    策略：
    1. 用BeautifulSoup解析
    2. 移除 script/style/nav/footer/header/aside 等噪音标签
    3. 提取主要文本内容
    4. 清理多余空白
    """
    try:
        soup = BeautifulSoup(html, "lxml")
    except:
        soup = BeautifulSoup(html, "html.parser")
    
    # 移除噪音标签
    for tag in soup.find_all(["script", "style", "nav", "footer", "header", "aside",
                               "noscript", "iframe", "form", "button"]):
        tag.decompose()
    
    # 移除广告和侧边栏（常见class/id）
    for pattern in ["ad", "ads", "advertisement", "sidebar", "widget", "comment",
                     "share", "social", "breadcrumb", "pagination", "cookie",
                     "popup", "modal", "overlay", "banner"]:
        for el in soup.find_all(class_=re.compile(pattern, re.I)):
            el.decompose()
        for el in soup.find_all(id=re.compile(pattern, re.I)):
            el.decompose()
    
    # 尝试找 article 或 main 标签
    main_content = soup.find("article") or soup.find("main") or soup.find(class_=re.compile(r"content|article|post|entry", re.I))
    
    if main_content:
        text = main_content.get_text(separator="\n", strip=True)
    else:
        text = soup.get_text(separator="\n", strip=True)
    
    # 清理多余空行
    lines = [l.strip() for l in text.split("\n") if l.strip()]
    text = "\n".join(lines)
    
    return text


async def crawl_urls_with_playwright(urls: list[str], company_website: str = "", max_concurrent: int = 3, timeout_per_page: int = 10) -> dict[str, dict]:
    """用 Playwright 抓取URL列表的完整页面内容
    
    共享一个浏览器实例，用多个page顺序抓取，避免反复启停浏览器。
    
    Args:
        urls: 要抓取的URL列表
        company_website: 公司官网URL（优先级最高）
        max_concurrent: 最大并发浏览器tab数（当前为顺序抓取，此参数保留）
        timeout_per_page: 每页超时秒数
    
    Returns:
        dict: {url: {"title": str, "content": str, "length": int}}
    """
    if not urls:
        return {}
    
    # 过滤值得抓取的URL
    worth_crawling = []
    domain_counts = {}  # domain -> count
    company_domain = ""
    
    if company_website:
        try:
            cw = company_website if "://" in company_website else f"https://{company_website}"
            company_domain = urlparse(cw).netloc.lower().replace("www.", "")
        except:
            pass
    
    for url in urls:
        if not _is_url_worth_crawling(url, company_website):
            continue
        
        # 去重：同域名下最多抓2个页面（官网不限）
        try:
            domain = urlparse(url).netloc.lower().replace("www.", "")
        except:
            domain = url
        
        count = domain_counts.get(domain, 0)
        if count >= 2 and domain != company_domain:
            continue  # 非官网域名最多2个页面
        
        domain_counts[domain] = count + 1
        worth_crawling.append(url)
    
    # 最多抓10个页面（避免耗时过长）
    worth_crawling = worth_crawling[:10]
    
    if not worth_crawling:
        return {}
    
    print(f"[Playwright] Crawling {len(worth_crawling)} URLs: {[u[:60] for u in worth_crawling]}")
    
    results = {}
    
    try:
        from playwright.async_api import async_playwright
        
        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=True)
            context = await browser.new_context(
                user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/125.0.0.0 Safari/537.36",
                viewport={"width": 1280, "height": 800},
            )
            await _protect_browser_context(context)
            page = await context.new_page()
            
            for url in worth_crawling:
                try:
                    await asyncio.wait_for(
                        page.goto(url, wait_until="domcontentloaded", timeout=timeout_per_page * 1000),
                        timeout=timeout_per_page + 5,
                    )
                    # 等待主要内容渲染
                    await asyncio.sleep(1)

                    title = await page.title()
                    html = await page.content()

                    # 提取正文
                    content = _extract_main_content(html, url)
                    
                    results[url] = {
                        "title": title,
                        "content": content,
                        "length": len(content),
                    }
                    print(f"[Playwright] Crawled {url[:60]}: title='{title[:40]}', content_len={len(content)}")
                except (asyncio.TimeoutError, Exception) as e:
                    print(f"[Playwright] Failed to crawl {url[:60]}: {type(e).__name__}")
                    results[url] = {"title": "", "content": f"[抓取失败: {type(e).__name__}]", "length": 0}
            
            await browser.close()
    
    except Exception as e:
        print(f"[Playwright] Browser error: {e}")
    

    return results


async def crawl_company_website(website: str, max_pages: int = 5) -> dict[str, dict]:
    """专门抓取公司官网的关键页面
    
    策略：先抓首页，从首页提取内部链接，再抓产品/关于我们/客户案例等关键页面
    """
    if not website:
        return {}
    
    base_url = website if "://" in website else f"https://{website}"
    if not await _is_safe_public_url(base_url):
        print(f"[Playwright] Blocked invalid or non-public website: {base_url[:120]}")
        return {}
    parsed = urlparse(base_url)
    domain = parsed.netloc.lower()
    
    # 首先抓取首页
    key_paths = [
        "/",           # 首页
        "/about",      # 关于我们
        "/aboutus",
        "/product",    # 产品
        "/products",
        "/service",    # 服务
        "/case",       # 案例
        "/cases",
        "/client",     # 客户
        "/clients",
        "/contact",    # 联系
    ]
    
    urls_to_crawl = []
    for path in key_paths:
        url = f"{parsed.scheme}://{domain}{path}"
        urls_to_crawl.append(url)
    
    # 限制最多抓取max_pages个页面
    urls_to_crawl = urls_to_crawl[:max_pages]
    
    print(f"[Playwright] Crawling company website {base_url}: {len(urls_to_crawl)} pages")
    
    results = {}
    
    try:
        from playwright.async_api import async_playwright
        
        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=True)
            context = await browser.new_context(
                user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/125.0.0.0 Safari/537.36",
                viewport={"width": 1280, "height": 800},
            )
            await _protect_browser_context(context)
            page = await context.new_page()
            
            # 先抓首页
            try:
                await page.goto(base_url, wait_until="domcontentloaded", timeout=15000)
                await asyncio.sleep(3)
                
                title = await page.title()
                html = await page.content()
                content = _extract_main_content(html, base_url)
                
                results[base_url] = {"title": title, "content": content, "length": len(content)}
                print(f"[Playwright] Homepage: title='{title[:40]}', content_len={len(content)}")
                
                # 从首页提取内部链接，找产品/关于/案例等页面
                links = await page.eval_on_selector_all("a[href]", "els => els.map(e => e.href)")
                internal_links = set()
                for link in links:
                    try:
                        lp = urlparse(link)
                        if lp.netloc.lower() == domain and lp.path != "/":
                            # 过滤关键词匹配的路径
                            path_lower = lp.path.lower()
                            if any(kw in path_lower for kw in ["product", "about", "case", "client", "service", "solution", "partner"]):
                                internal_links.add(link)
                    except:
                        pass
                
                # 抓取发现的内部链接
                pages_crawled = 1
                for link in list(internal_links)[:max_pages - 1]:
                    try:
                        await page.goto(link, wait_until="domcontentloaded", timeout=10000)
                        await asyncio.sleep(2)
                        
                        t = await page.title()
                        h = await page.content()
                        c = _extract_main_content(h, link)
                        
                        results[link] = {"title": t, "content": c, "length": len(c)}
                        pages_crawled += 1
                        print(f"[Playwright] Sub-page: {link[:60]}, content_len={len(c)}")
                    except Exception as e:
                        print(f"[Playwright] Sub-page failed {link[:60]}: {e}")
                
            except Exception as e:
                print(f"[Playwright] Homepage crawl failed: {type(e).__name__}: {e}")
            finally:
                await browser.close()
    
    except Exception as e:
        print(f"[Playwright] Browser error during website crawl: {e}")
    
    return results


# ============ LLM API (DeepSeek + Kimi, OpenAI-compatible) ============
async def call_deepseek(messages: list[dict], retry_count: int = 1, max_tokens: int = 16384, model: str = "") -> str:
    """调用 LLM API（兼容 OpenAI 格式），返回完整回复文本。

    model 参数可选，默认使用 DEEPSEEK_MODEL。支持 DeepSeek 和 Kimi (Moonshot)。
    """
    cfg = _get_model_config(model) if model else _get_model_config(DEEPSEEK_MODEL)
    effective_model = model or DEEPSEEK_MODEL
    api_key = DEEPSEEK_API_KEY if cfg["provider"] == "deepseek" else KIMI_API_KEY
    base_url = cfg["base_url"]
    provider = cfg["provider"]

    for attempt in range(retry_count + 1):
        print(f"[LLM-{provider}] Calling {effective_model}, messages={len(messages)}, attempt={attempt+1}/{retry_count+1}")

        try:
            timeout_sec = 600.0 if provider == "kimi" else 300.0
            request_body = {
                "model": effective_model,
                "messages": messages,
                "stream": False,
                "temperature": 0,
            }
            if provider == "kimi":
                request_body["max_completion_tokens"] = max_tokens
                request_body["thinking"] = {"type": "disabled"}
                request_body["prompt_cache_key"] = "pe-screening"
            else:
                request_body["max_tokens"] = max_tokens
            async with httpx.AsyncClient(timeout=httpx.Timeout(timeout_sec, connect=30.0)) as client:
                resp = await client.post(
                    f"{base_url}/chat/completions",
                    headers={
                        "Authorization": f"Bearer {api_key}",
                        "Content-Type": "application/json",
                    },
                    json=request_body,
                )
                print(f"[LLM-{provider}] Response: {resp.status_code}")

                if resp.status_code != 200:
                    print(f"[LLM-{provider}] Error: {resp.text[:500]}")
                    resp.raise_for_status()

                data = resp.json()
                usage = data.get("usage", {})
                if usage:
                    print(f"[LLM-{provider}] Tokens - prompt: {usage.get('prompt_tokens','?')}, completion: {usage.get('completion_tokens','?')}")
                choice = data["choices"][0]
                msg = choice["message"]
                finish_reason = choice.get("finish_reason", "unknown")
                content = msg.get("content", "") or ""
                reasoning = msg.get("reasoning_content", "") or ""

                print(f"[LLM-{provider}] Finish: {finish_reason}, content: {len(content)} chars, reasoning: {len(reasoning)} chars")
                print(f"[LLM-{provider}] Preview: {content[:200] if content else '(empty)'}")

                # 如果 content 为空但有 reasoning_content，说明思维链太长耗尽了 token
                if not content.strip():
                    if reasoning.strip():
                        if attempt < retry_count:
                            print(f"[LLM-{provider}] Content empty (reasoning took all tokens), retrying...")
                            continue
                        return f"[思维链输出耗尽了回复配额，以下是思维过程摘要]\n{reasoning[:2000]}"
                    if attempt < retry_count:
                        print(f"[LLM-{provider}] Both content and reasoning empty, retrying...")
                        continue
                    return ""
                return content
        except httpx.TimeoutException as e:
            print(f"[LLM-{provider}] Timeout: {e}")
            if attempt < retry_count:
                continue
            raise
        except Exception as e:
            print(f"[LLM-{provider}] Exception: {type(e).__name__}: {e}")
            if attempt < retry_count:
                continue
            raise

    return ""


async def _call_deepseek_raw(messages: list[dict], tools: list[dict] | None = None, model: str = "",
                             stream_callback=None) -> dict:
    """Raw LLM API call with optional tool support. Supports DeepSeek and Kimi models.

    If ``stream_callback`` is provided, uses SSE streaming and calls the
    callback with each text delta as it arrives.  The return value still
    includes the fully accumulated content.

    Returns {"content": str, "tool_calls": list | None}
    """
    cfg = _get_model_config(model) if model else _get_model_config(DEEPSEEK_MODEL)
    effective_model = model or DEEPSEEK_MODEL
    api_key = DEEPSEEK_API_KEY if cfg["provider"] == "deepseek" else KIMI_API_KEY
    base_url = cfg["base_url"]
    provider = cfg["provider"]
    use_stream = stream_callback is not None and provider == "deepseek"  # Kimi SSE unreliable

    print(f"[LLM-Raw-{provider}] Calling {effective_model}, messages={len(messages)}, tools={bool(tools)}, stream={use_stream}")

    json_body = {
        "model": effective_model,
        "messages": messages,
        "stream": use_stream,
        "temperature": 0,
    }
    if provider == "kimi":
        json_body["max_completion_tokens"] = 12000
        json_body["thinking"] = {"type": "disabled"}
        json_body["prompt_cache_key"] = "pe-screening"
    if tools:
        json_body["tools"] = tools
        if provider != "kimi":
            json_body["tool_choice"] = "auto"

    raw_timeout = 600.0 if provider == "kimi" else 300.0
    async with httpx.AsyncClient(timeout=httpx.Timeout(raw_timeout, connect=30.0)) as client:
        if use_stream:
            # ---- SSE streaming ----
            content_parts = []
            async with client.stream("POST", f"{base_url}/chat/completions",
                                     headers={"Authorization": f"Bearer {api_key}",
                                              "Content-Type": "application/json"},
                                     json=json_body) as stream_resp:
                if stream_resp.status_code != 200:
                    err = await stream_resp.aread()
                    print(f"[LLM-Raw-{provider}] Stream error: {stream_resp.status_code} {err[:300]}")
                    stream_resp.raise_for_status()
                async for line in stream_resp.aiter_lines():
                    if not line or not line.startswith("data: "):
                        continue
                    payload = line[6:]
                    if payload == "[DONE]":
                        break
                    try:
                        chunk = json.loads(payload)
                        delta = (chunk.get("choices", [{}])[0].get("delta") or {})
                        delta_content = delta.get("content", "")
                        if delta_content:
                            content_parts.append(delta_content)
                            try:
                                stream_callback(delta_content)
                            except Exception:
                                pass
                    except json.JSONDecodeError:
                        continue
            content = "".join(content_parts)
            print(f"[LLM-Raw-{provider}] Stream complete: {len(content)} chars")
            return {"content": content, "tool_calls": None, "reasoning": ""}

        # ---- Non-streaming ----
        resp = await client.post(
            f"{base_url}/chat/completions",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json=json_body,
        )

    if resp.status_code != 200:
        print(f"[LLM-Raw-{provider}] Error: {resp.status_code} {resp.text[:300]}")
        resp.raise_for_status()

    data = resp.json()
    usage = data.get("usage", {})
    if usage:
        print(f"[LLM-Raw-{provider}] Tokens - prompt: {usage.get('prompt_tokens','?')}, completion: {usage.get('completion_tokens','?')}")

    choice = data["choices"][0]
    msg = choice["message"]
    finish_reason = choice.get("finish_reason", "unknown")
    content = msg.get("content", "") or ""
    reasoning = msg.get("reasoning_content", "") or ""
    tool_calls = msg.get("tool_calls") or None

    print(f"[LLM-Raw-{provider}] Finish: {finish_reason}, content: {len(content)} chars, tool_calls: {len(tool_calls) if tool_calls else 0}")

    return {"content": content, "tool_calls": tool_calls, "reasoning": reasoning}


def _format_search_results_for_ai(
    results: list[dict], start_ref_num: int = 1, seen_urls: set | None = None
) -> tuple[str, list[dict]]:
    """Format search results for AI with reference numbers. Returns (text, ref_items).

    Filters out spam domains, duplicate URLs, and empty results.
    """
    seen = seen_urls if seen_urls is not None else set()
    lines = []
    ref_items = []
    for r in results:
        if r.get("type") != "result":
            continue
        url = (r.get("url", "") or "").strip().rstrip("/").lower()
        if not url:
            continue
        # Block known spam/scraper domains
        if any(d in url for d in _SPAM_DOMAINS):
            continue
        # Dedup
        if url in seen:
            continue
        seen.add(url)
        ref_num = start_ref_num + len(ref_items)
        title = r.get("title", "") or ""
        snippet = (r.get("content", "") or "")[:500]
        lines.append(f"[{ref_num}] {title}\n    URL: {url}\n    {snippet}")
        ref_items.append({
            "provider": r.get("provider", "Tavily"),
            "title": title,
            "url": url,
        })
    text = "\n\n".join(lines) if lines else "（未找到相关搜索结果）"
    return text, ref_items


async def _execute_search_tool(
    tool_name: str,
    arguments: dict,
    ref_counter: int = 1,
    seen_urls: set | None = None,
    company_name: str = "",
) -> tuple[str, list[dict]]:
    """Execute a search tool. Returns (formatted_text_for_ai, new_ref_items)."""
    seen = seen_urls if seen_urls is not None else set()
    if tool_name == "search_web":
        query = arguments.get("query", "")
        if not query:
            return "错误：缺少搜索查询词", []
        print(f"[Tool] search_web: {query[:100]}")
        try:
            results = await public_web_search(
                query, company_name=company_name, max_results=8, comprehensive=False
            )
            if company_name:
                results = _filter_verified_search_results(results, company_name, seen_urls=seen)
                text, refs = _format_search_results_for_ai(results, start_ref_num=ref_counter)
            else:
                text, refs = _format_search_results_for_ai(results, start_ref_num=ref_counter, seen_urls=seen)
            return text, refs
        except Exception as e:
            print(f"[Tool] search_web error: {e}")
            return f"搜索出错：{str(e)[:200]}", []

    elif tool_name == "search_wechat":
        query = arguments.get("query", "")
        if not query:
            return "错误：缺少搜索查询词", []
        print(f"[Tool] search_wechat: {query[:100]}")
        try:
            raw = await search_wechat_articles(query, max_results=5)
            deduped = []
            seen_wx = seen.copy()
            for art in raw:
                u = str(art.get("url", "")).strip().lower()
                if u and u not in seen_wx:
                    seen_wx.add(u)
                    deduped.append(art)
            if company_name:
                deduped = _filter_relevant_wechat_results(
                    _dedupe_wechat_results(deduped), company_name
                )
            # Format with reference numbers for AI citation
            lines = []
            refs = []
            for i, art in enumerate(deduped[:6]):
                title = art.get('title', '')
                url = art.get('url', '')
                lines.append(
                    f"[公众号] {title}\\n"
                    f"    账号: {art.get('account','')} | 时间: {art.get('time','')}\\n"
                    f"    {art.get('digest','')[:300]}"
                )
                if url:
                    refs.append({
                        "provider": "微信补充",
                        "title": f"{title} - {art.get('account', '')}",
                        "url": url,
                    })
            text = "\\n\\n".join(lines) if lines else "（未找到相关微信公众号文章）"
            return text, refs
        except Exception as e:
            print(f"[Tool] search_wechat error: {e}")
            return f"微信搜索出错：{str(e)[:200]}", []

    elif tool_name == "fetch_page":
        url = arguments.get("url", "")
        if not url:
            return "错误：缺少 URL", []
        print(f"[Tool] fetch_page: {url[:100]}")
        try:
            async with httpx.AsyncClient(timeout=httpx.Timeout(30.0, connect=10.0)) as client:
                resp = await client.get(
                    url,
                    headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"},
                    follow_redirects=True,
                )
            if resp.status_code != 200:
                return f"获取页面失败：HTTP {resp.status_code}", []
            text = resp.text
            text = __import__("re").sub(r"<[^>]+>", " ", text)
            text = __import__("re").sub(r"\s+", " ", text)
            return f"页面内容（已提取文本，截取前5000字符）：\n{text[:5000]}", []
        except Exception as e:
            print(f"[Tool] fetch_page error: {e}")
            return f"获取页面出错：{str(e)[:200]}", []

    return f"未知工具：{tool_name}", []


def _dedupe_planned_queries(queries: list, company_name: str, limit: int) -> list[str]:
    selected = []
    seen = set()
    for raw in queries if isinstance(queries, list) else []:
        query = re.sub(r"\s+", " ", str(raw or "")).strip()
        if not query:
            continue
        if not _has_company_reference(query, company_name):
            query = f"{company_name} {query}"
        key = query.lower()
        if key in seen:
            continue
        seen.add(key)
        selected.append(query[:120])
        if len(selected) >= limit:
            break
    return selected


async def _plan_model_supplemental_searches(
    company_name: str,
    current_search_text: str,
    model: str = "",
) -> dict[str, list[str]]:
    """Ask the selected model for company-specific follow-up search queries only."""
    compact_evidence = re.sub(r"\s+", " ", str(current_search_text or ""))[:8000]
    prompt = f"""你只负责为目标公司生成补搜查询词，不做分析。

目标公司：{company_name}

当前已收集证据摘要：
{compact_evidence}

请找出公司个体证据缺口，生成补搜查询词。只允许围绕目标公司的产品、客户、产能、招聘、官网/公众号、资质、处罚、诉讼、股权、仓储/工厂/地址等个体信息。
不要搜索泛行业市场规模、行业趋势、赛道背景。
查询词必须包含公司全称或核心简称。

只输出JSON：
{{"web_queries":["..."],"wechat_queries":["..."]}}
"""
    try:
        response = await call_deepseek(
            [{"role": "user", "content": prompt}],
            retry_count=0,
            max_tokens=1200,
            model=model,
        )
        match = re.search(r"\{.*\}", response or "", re.DOTALL)
        data = json.loads(match.group(0) if match else response)
    except Exception as exc:
        print(f"[ModelSearchPlan] skipped: {type(exc).__name__}: {exc}")
        return {"web_queries": [], "wechat_queries": []}
    return {
        "web_queries": _dedupe_planned_queries(data.get("web_queries", []), company_name, 5),
        "wechat_queries": _dedupe_planned_queries(data.get("wechat_queries", []), company_name, 3),
    }


async def _run_model_planned_supplemental_searches(
    company_name: str,
    current_search_text: str,
    ref_counter: int,
    seen_urls: set[str],
    provider_counts: dict[str, int],
    model: str = "",
) -> tuple[str, str, int]:
    """Execute model-planned searches with the same deterministic source governance as base search."""
    plan = await _plan_model_supplemental_searches(company_name, current_search_text, model=model)
    additions = []
    refs_html = ""

    web_queries = plan.get("web_queries", [])
    if web_queries:
        additions.append("\n\n## 模型补搜结果（公司个体证据）\n")
    for query in web_queries:
        try:
            raw_results = []
            try:
                raw_results += await exa_search(query, max_results=5)
            except Exception as exc:
                print(f"[ModelSearchPlan] exa failed for {query}: {type(exc).__name__}: {exc}")
            if TAVILY_FALLBACK:
                try:
                    raw_results += await tavily_search(query, max_results=5, search_depth="advanced")
                except Exception as exc:
                    print(f"[ModelSearchPlan] tavily failed for {query}: {type(exc).__name__}: {exc}")
            results = raw_results or await public_web_search(
                query, company_name=company_name, max_results=8, comprehensive=True
            )
            results = _filter_verified_search_results(results, company_name, seen_urls=seen_urls)
        except Exception as exc:
            print(f"[ModelSearchPlan] web failed for {query}: {type(exc).__name__}: {exc}")
            continue
        if not results:
            continue
        additions.append(f"\n### 网页补搜：{query}\n")
        for item in results:
            if item.get("type") != "result":
                continue
            provider = item.get("provider", "网页搜索")
            if not _can_add_report_reference(provider, ref_counter, provider_counts):
                continue
            additions.append(
                f"[{ref_counter}] {item.get('title', '')}\n"
                f"URL: {item.get('url', '')}\n"
                f"{str(item.get('content', ''))[:700]}\n\n"
            )
            refs_html += _reference_item_html(ref_counter, provider, item.get("title", ""), item.get("url", ""))
            _mark_report_reference(provider, provider_counts)
            ref_counter += 1

    wechat_queries = plan.get("wechat_queries", [])
    if wechat_queries:
        additions.append("\n\n## 模型补搜微信公众号结果（公司经营事实）\n")
    for query in wechat_queries:
        try:
            raw = await search_wechat_articles(query, max_results=5)
            articles = _filter_relevant_wechat_results(_dedupe_wechat_results(raw), company_name)
        except Exception as exc:
            print(f"[ModelSearchPlan] wechat failed for {query}: {type(exc).__name__}: {exc}")
            continue
        if not articles:
            continue
        additions.append(f"\n### 微信补搜：{query}\n")
        for article in articles[:3]:
            url = str(article.get("url", "")).strip().rstrip("/").lower()
            if url and url in seen_urls:
                continue
            if not _can_add_report_reference("微信补充", ref_counter, provider_counts):
                continue
            if url:
                seen_urls.add(url)
            additions.append(
                f"[{ref_counter}] {article.get('title', '')}\n"
                f"URL: {article.get('url', '')}\n"
                f"公众号: {article.get('account', '')} | {article.get('time', '')}\n"
                f"{str(article.get('digest', ''))[:700]}\n\n"
            )
            refs_html += _reference_item_html(
                ref_counter,
                "微信补充",
                f'{article.get("title", "")} - {article.get("account", "")}',
                article.get("url", ""),
            )
            _mark_report_reference("微信补充", provider_counts)
            ref_counter += 1

    return "".join(additions), refs_html, ref_counter


async def call_deepseek_with_tools(
    messages: list[dict],
    tools: list[dict],
    max_rounds: int = 5,
    ref_counter_start: int = 1,
    company_name: str = "",
    model: str = "",
    progress_callback=None,
    stream_callback=None,
) -> tuple[str, str, int]:
    """Multi-turn LLM call with tool support. AI can call search tools autonomously.

    If ``stream_callback`` is provided, streaming is enabled for the final
    analysis round (not intermediate tool-call rounds).

    Returns (analysis_md, ai_refs_html, final_ref_count).
    """
    ref_counter = ref_counter_start
    all_ai_refs = []
    seen_urls = set()  # dedup across all AI search rounds

    for round_num in range(max_rounds):
        print(f"[Tool-Loop] Round {round_num+1}/{max_rounds}, next ref #{ref_counter}")
        if progress_callback:
            progress_callback(f"AI 正在分析中（第{round_num+1}/{max_rounds}轮）...")

        result = await _call_deepseek_raw(messages, tools=tools, model=model)

        tool_calls = result.get("tool_calls")
        content = result.get("content", "")

        if tool_calls:
            messages.append({
                "role": "assistant",
                "content": content or "",
                "tool_calls": tool_calls,
            })

            for tc in tool_calls:
                tc_id = tc.get("id", "")
                func = tc.get("function", {})
                func_name = func.get("name", "")
                try:
                    args = json.loads(func.get("arguments", "{}"))
                except json.JSONDecodeError:
                    args = {}

                if progress_callback:
                    query_preview = args.get("query", args.get("url", ""))[:40]
                    progress_callback(f"AI 正在搜索：{query_preview}...")

                tool_text, new_refs = await _execute_search_tool(
                    func_name,
                    args,
                    ref_counter=ref_counter,
                    seen_urls=seen_urls,
                    company_name=company_name,
                )

                for ref_item in new_refs:
                    ref_item["ref_num"] = ref_counter
                    ref_counter += 1
                all_ai_refs.extend(new_refs)

                messages.append({
                    "role": "tool",
                    "tool_call_id": tc_id,
                    "content": tool_text,
                })
                print(f"[Tool-Loop] Executed {func_name}, {len(new_refs)} refs, next ref #{ref_counter}")

            continue

        if content.strip():
            print(f"[Tool-Loop] AI completed analysis in round {round_num+1}")
            break

        print(f"[Tool-Loop] Empty response, round {round_num+1}")

    if not content.strip():
        messages.append({
            "role": "user",
            "content": "请基于目前收集到的所有信息，输出完整的深筛分析（不要再调用搜索工具）。"
        })
        final = await _call_deepseek_raw(messages, tools=None, model=model, stream_callback=stream_callback)
        content = final.get("content", "") or ""

    ai_refs_html = ""
    for ref in all_ai_refs:
        ai_refs_html += _reference_item_html(
            ref["ref_num"],
            ref.get("provider", "AI搜索"),
            ref.get("title", ""),
            ref.get("url", ""),
        )

    return content, ai_refs_html, ref_counter



# ============ BP 文件提取 ============
async def extract_bp_text(file_path: Path) -> str:
    """提取 BP 文件文本内容"""
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        # 使用 Read 工具类似的逻辑 - 用 PyMuPDF 或简单文本提取
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(str(file_path))
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text[:50000]  # 限制长度
        except ImportError:
            # 回退：尝试用 pdfminer
            try:
                from pdfminer.high_level import extract_text
                text = extract_text(str(file_path))
                return text[:50000]
            except ImportError:
                return "[PDF 提取失败：未安装 PyMuPDF 或 pdfminer]"

    elif suffix in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text += para.text + "\n"
            return text[:50000]
        except ImportError:
            return "[PPTX 提取失败：未安装 python-pptx]"

    elif suffix in (".txt", ".md", ".html", ".htm"):
        return file_path.read_text(encoding="utf-8")[:50000]

    else:
        return f"[不支持的文件格式: {suffix}]"


# ============ 报告生成 ============
_SAFE_RICH_TAGS = {
    "p", "br", "strong", "b", "em", "i", "ul", "ol", "li", "table", "thead",
    "tbody", "tr", "th", "td", "div", "span", "sup", "a", "hr", "h3", "h4",
}
_SAFE_RICH_CLASSES = {
    "ev-a", "ev-b", "ev-c", "ev-d", "note-box", "step-item", "q", "why",
    "ref-num", "ref-content", "cite",
}


def _escape_text(value) -> str:
    return html.escape(str(value if value is not None else ""), quote=True)


def _safe_link(value: str) -> str:
    try:
        parsed = urlparse(str(value).strip())
        if parsed.scheme in ("http", "https") and parsed.netloc:
            return str(value).strip()
    except ValueError:
        pass
    return "#"


_MAX_REPORT_REFERENCES = 30
_MAX_TAVILY_REPORT_REFERENCES = 12


def _can_add_report_reference(provider: str, ref_counter: int, provider_counts: dict[str, int]) -> bool:
    """Keep the final reference list compact and prevent Tavily from dominating it."""
    if ref_counter > _MAX_REPORT_REFERENCES:
        return False
    normalized = str(provider or "")
    if normalized == "Tavily" and provider_counts.get(normalized, 0) >= _MAX_TAVILY_REPORT_REFERENCES:
        return False
    return True


def _mark_report_reference(provider: str, provider_counts: dict[str, int]) -> None:
    normalized = str(provider or "")
    provider_counts[normalized] = provider_counts.get(normalized, 0) + 1


def _reference_item_html(number: int, source: str, title: str, url: str) -> str:
    """Render an actual numbered source rather than accepting model-created source entries."""
    safe_href = html.escape(_safe_link(url), quote=True)
    label = _escape_text(f"[{source}] {title}".strip())
    return (
        f'<li id="ref-{number}"><span class="ref-num">{number}</span>'
        f'<span class="ref-content"><a href="{safe_href}" target="_blank" rel="noopener">'
        f'{label}</a></span></li>\n'
    )


def _sanitize_html_fragment(fragment) -> str:
    """Keep report formatting while removing executable markup from AI output."""
    if isinstance(fragment, (list, tuple)):
        fragment = "".join(str(item) for item in fragment)
    soup = BeautifulSoup(str(fragment or ""), "html.parser")
    for tag in list(soup.find_all(True)):
        if tag.name in {"script", "style", "iframe", "object", "embed", "svg", "math", "form", "input"}:
            tag.decompose()
            continue
        if tag.name not in _SAFE_RICH_TAGS:
            tag.unwrap()
            continue
        attrs = {}
        if tag.name == "a":
            attrs = {"href": _safe_link(tag.get("href", "#")), "target": "_blank", "rel": "noopener"}
        elif tag.name == "li":
            identifier = str(tag.get("id", ""))
            if re.fullmatch(r"ref-\d+", identifier):
                attrs = {"id": identifier}
        safe_classes = [item for item in tag.get("class", []) if item in _SAFE_RICH_CLASSES]
        if safe_classes:
            attrs["class"] = safe_classes
        tag.attrs = attrs
    return "".join(str(node) for node in soup.contents)


def _safe_numeric_values(values, length: int, default: float) -> list[str]:
    result = []
    raw_values = values if isinstance(values, list) else []
    for index in range(length):
        try:
            value = float(raw_values[index])
            value = max(0, min(5, value))
            result.append(f"{value:g}")
        except (IndexError, TypeError, ValueError):
            result.append(f"{default:g}")
    return result


_SCORE_CAPS = {
    1: {1: 2},        # 一次性项目为主 → RR ≤2 (index unchanged: 1)
    3: {2: 2, 3: 2},  # 纯贸易低壁垒 → 客户粘性(now index 2) ≤2, 商业模式(now index 3) ≤2
    6: {2: 2},         # 依赖单一大客户 → 客户粘性(now index 2) ≤2
    10: {0: 2},        # 行业边界模糊 → 行业匹配度(index 0) ≤2
    13: {2: 2},        # 客户质量不足 → 客户粘性(now index 2) ≤2
}
_CRITICAL_RED_FLAGS = {9, 12}


def _bounded_score(value, default: int = 2) -> int:
    try:
        return max(1, min(5, int(round(float(value)))))
    except (TypeError, ValueError):
        return default


def _normalize_flag_marker(value) -> str:
    text = str(value or "")
    if "🔴" in text or text.strip().lower() in {"red", "high", "严重"}:
        return "🔴"
    if "🟡" in text or text.strip().lower() in {"yellow", "medium", "中等"}:
        return "🟡"
    if "🟢" in text or text.strip().lower() in {"green", "low", "无", "未发现"}:
        return "🟢"
    return "🟡"


def _is_routine_sme_information_gap(text: str) -> bool:
    """Check if flag #12 is triggered by routine SME info opacity, not substantive risk."""
    text = str(text or "")
    information_gap_signals = (
        "未提供任何财务数据", "未提供财务数据", "财务数据不公开", "未公开财务",
        "未披露财务", "未披露审计", "缺乏财务数据", "财务数据完全不公开",
        "股权结构.{0,6}(?:不透明|极不透明|不明|无法|不详|未知)",
        "无法查清.{0,8}(?:股东|实际控制人|实控人)",
        "股东.{0,4}(?:不公开|不明|不详|无法查|未公开|未披露)",
        "实际控制人.{0,4}(?:不公开|不明|不详|无法查|未公开|未披露)",
        "实控人.{0,4}(?:不公开|不明|不详|无法查|未公开|未披露)",
        "未透露.{0,8}(?:股东|实控人|实际控制人|集团名称|所属集团)",
        "无法确认.{0,8}(?:实际控制人|实控人|股东|所属)",
        "股东信息.{0,4}(?:不公开|缺乏|极少|有限)",
        "工商信息.{0,4}(?:不完整|有限|缺乏)",
    )
    _GAP_RE = re.compile("|".join(information_gap_signals))
    cleaned_text = re.sub(
        r"未(?:发现|见|有)\s*(?:实质)?(?:行政)?(?:处罚|违规|异常|失信|冻结|被执行)",
        "",
        text,
    )
    substantive_risk_signals = (
        "主体不清", "主体混淆", "多个主体",
        "股权代持", "关联交易", "造假", "虚假", "欠税", "冻结", "失信",
        "被执行", "吊销", "注销", "行政处罚", "审计保留", "审计意见",
        "壳公司", "名义股东", "代持", "挪用", "侵占",
    )
    return (
        bool(_GAP_RE.search(text))
        and not any(signal in cleaned_text for signal in substantive_risk_signals)
    )


def _has_premium_customer_evidence(text: str) -> bool:
    """Check if flag #13 evidence actually lists premium/quality customers."""
    text = str(text or "")
    _PREMIUM_CUSTOMERS = re.compile(
        r"盒马|叮咚|美团|山姆|开市客|Costco|奥乐齐"
        r"|医院|学校|央企|国企|世界500强|机关|部队"
        r"|百果园|天天果园|永辉|大润发|沃尔玛|麦德龙"
        r"|全家|罗森|7-Eleven|便利蜂"
        r"|海底捞|喜茶|奈雪|古茗|瑞幸|星巴克"
    )
    return bool(_PREMIUM_CUSTOMERS.search(text))


def _describes_multiple_customers(text: str) -> bool:
    """Check if flag #6 evidence describes multiple/diverse customers, not single dependency."""
    text = str(text or "")
    _MULTI_CUSTOMER = re.compile(
        r"多[个家种类项]|等[平台渠道客户]|、.*、"
        r"|数个|数个|若干|批量|很多|广泛"
    )
    return bool(_MULTI_CUSTOMER.search(text))


def _has_founder_dependency_mitigation(text: str) -> bool:
    """Check if flag #7 evidence shows mitigating factors against founder dependency."""
    text = str(text or "")
    _MITIGATION = re.compile(
        r"部分存在"
        r"|开始建立|正在建立|已建立|逐步建立"
        r"|职业经理人|招聘高管|聘请.{0,4}(?:经理|总监|负责人)"
        r"|降低了|正在降低|逐步降低|有所降低"
        r"|制度化|体系化|规范化"
        r"|销售团队|业务团队|管理团队"
        r"|交接|过渡|传承|接班人"
        r"|不止.{0,3}老板|非老板.{0,3}单独"
    )
    return bool(_MITIGATION.search(text))


def _has_compliance_green_evidence(text: str) -> bool:
    text = str(text or "")
    green_patterns = (
        r"无.{0,8}(?:诉讼|处罚|行政处罚|被执行|失信|经营异常|股权冻结|重大合规)",
        r"未(?:发现|见|有).{0,10}(?:诉讼|处罚|行政处罚|被执行|失信|经营异常|股权冻结|重大合规)",
        r"(?:诉讼|处罚|行政处罚|被执行|失信|经营异常).{0,8}(?:无|未见|未发现)",
    )
    green = bool(re.search("|".join(green_patterns), text))
    risk_text = re.sub("|".join(green_patterns), "", text)
    substantive = bool(re.search(
        r"被执行|失信|经营异常|行政处罚|重大诉讼|股权冻结|吊销|注销|食品安全事故|监管处罚|处罚决定",
        risk_text,
    ))
    return green and not substantive


def _extract_red_flags(result: dict) -> dict[str, str]:
    flags = {}
    raw_flags = result.get("red_flags")
    if isinstance(raw_flags, dict):
        for key, marker in raw_flags.items():
            try:
                number = int(key)
            except (TypeError, ValueError):
                continue
            if 1 <= number <= 14:
                flags[str(number)] = _normalize_flag_marker(marker)
    html_flags = str(result.get("red_flags_html", ""))
    for marker, number in re.findall(r"([🟢🟡🔴])\s*(\d{1,2})\s*[.。：:]", html_flags):
        if 1 <= int(number) <= 14:
            flags.setdefault(number, marker)

    sec6_verdict = str(result.get("sec7_verdict", "")).strip().upper()
    platform_score = result.get("platform_score")
    addon_score = result.get("addon_score")
    evidence_confidence = str(result.get("evidence_confidence", "")).strip().upper()
    plain_flags_html = re.sub(r"<[^>]+>", " ", html_flags)
    flag_12_match = re.search(
        r"[🟢🟡🔴]\s*12\s*[.。：:](.*?)(?=[🟢🟡🔴]\s*\d{1,2}\s*[.。：:]|$)",
        plain_flags_html,
        flags=re.DOTALL,
    )
    flag_12_text = flag_12_match.group(1) if flag_12_match else ""
    flag_9_match = re.search(
        r"[🟢🟡🔴]\s*9\s*[.。：:](.*?)(?=[🟢🟡🔴]\s*\d{1,2}\s*[.。：:]|$)",
        plain_flags_html,
        flags=re.DOTALL,
    )
    flag_9_text = flag_9_match.group(1) if flag_9_match else ""

    # SME-friendly normalization:
    # - Flag 11 should not be severe when the target is clearly an Add-on candidate.
    # - Flag 14 is an information-quality warning and should only be severe when evidence is truly too sparse.
    # - Flag 12 remains critical for substantive entity/financial irregularities, not routine non-disclosure.
    if flags.get("11") == "🔴":
        if "ADD-ON" in sec6_verdict or (
            isinstance(platform_score, (int, float))
            and isinstance(addon_score, (int, float))
            and addon_score >= platform_score
        ):
            flags["11"] = "🟡"
    if flags.get("14") == "🔴" and evidence_confidence != "D":
        flags["14"] = "🟡"
    if flags.get("9") == "🔴" and _has_compliance_green_evidence(flag_9_text):
        flags["9"] = "🟡"
    if flags.get("12") == "🔴" and _is_routine_sme_information_gap(flag_12_text):
        flags["12"] = "🟡"

    flag_13_match = re.search(
        r"[🟢🟡🔴]\s*13\s*[.。：:](.*?)(?=[🟢🟡🔴]\s*\d{1,2}\s*[.。：:]|$)",
        plain_flags_html, flags=re.DOTALL,
    )
    flag_13_text = flag_13_match.group(1) if flag_13_match else ""
    if flags.get("13") == "🔴" and _has_premium_customer_evidence(flag_13_text):
        flags["13"] = "🟡"

    flag_6_match = re.search(
        r"[🟢🟡🔴]\s*6\s*[.。：:](.*?)(?=[🟢🟡🔴]\s*\d{1,2}\s*[.。：:]|$)",
        plain_flags_html, flags=re.DOTALL,
    )
    flag_6_text = flag_6_match.group(1) if flag_6_match else ""
    if flags.get("6") == "🔴" and _describes_multiple_customers(flag_6_text):
        flags["6"] = "🟡"

    flag_7_match = re.search(
        r"[🟢🟡🔴]\s*7\s*[.。：:](.*?)(?=[🟢🟡🔴]\s*\d{1,2}\s*[.。：:]|$)",
        plain_flags_html, flags=re.DOTALL,
    )
    flag_7_text = flag_7_match.group(1) if flag_7_match else ""
    if flags.get("7") == "🔴" and _has_founder_dependency_mitigation(flag_7_text):
        flags["7"] = "🟡"

    return {str(number): flags.get(str(number), "🟡") for number in range(1, 15)}
def _dimension_evidence(result: dict, quick: bool) -> tuple[list[int], float, str]:
    raw = result.get("evidence_levels")
    values = [min(4, _bounded_score(value, 2)) for value in raw] if isinstance(raw, list) else []
    # New 5-dimension system with 7 evidence levels (removed 平台价值 and 可交易性)
    if quick:
        evidence = (values + [2] * 5)[:5]
    elif len(values) >= 7:
        # values: [行业归属, 商业模式, RR, 客户质量, 切换成本, 市场地位, 规模推断]
        evidence = [
            values[0],                           # 行业匹配度
            values[2],                           # Recurring Revenue
            min(values[3], values[4]),            # 客户粘性 (combined)
            values[1],                           # 商业模式
            values[5],                           # 市场地位
        ]
    else:
        evidence = (values + [2] * 5)[:5]
    confidence_basis = (
        [values[0], values[1], values[2], values[3], values[4], values[5], values[6]]
        if len(values) >= 7
        else evidence
    )
    direct_count = sum(level >= 4 for level in confidence_basis)
    strong_count = sum(level >= 3 for level in confidence_basis)
    average = round(sum(confidence_basis) / len(confidence_basis), 1)
    if direct_count >= 5:
        label = "A"
    elif (direct_count >= 3 and strong_count >= 5) or strong_count >= 6:
        label = "B"
    elif strong_count >= 3:
        label = "C"
    else:
        label = "D"
    return evidence, average, label


def _evidence_confidence_note(label: str) -> str:
    notes = {
        "A": "主要结论由直接证据或强公开证据支撑，适合继续推进。",
        "B": "大部分关键判断有较强公开证据，少量内容仍需推断。",
        "C": "关键判断里推断成分明显，适合观察，不宜强结论。",
        "D": "公开证据太少，只能做初步参考。",
    }
    return notes.get(label, "证据口径待确认")


def _parse_revenue_bounds(text: str) -> tuple[float | None, float | None]:
    """Parse rough RMB revenue bounds from a free-form revenue string."""
    if not text:
        return None, None
    s = str(text).replace(",", "").replace(" ", "")

    def _to_rmb(num: str, unit: str) -> float:
        value = float(num)
        return value * (1e8 if unit == "亿" else 1e4)

    m = re.search(r"(\d+(?:\.\d+)?)\s*-\s*(\d+(?:\.\d+)?)\s*(亿|万)", s)
    if m:
        a, b, unit = m.groups()
        return _to_rmb(a, unit), _to_rmb(b, unit)

    m = re.search(r"(\d+(?:\.\d+)?)\s*(亿|万)\s*-\s*(\d+(?:\.\d+)?)\s*(亿|万)", s)
    if m:
        a, u1, b, u2 = m.groups()
        return _to_rmb(a, u1), _to_rmb(b, u2)

    m = re.search(r"(\d+(?:\.\d+)?)\s*(亿|万)(\+)?", s)
    if m:
        a, unit, lower_bound = m.groups()
        value = _to_rmb(a, unit)
        return (value, None) if lower_bound else (value, value)

    nums = re.findall(r"(\d+(?:\.\d+)?)\s*(亿|万)", s)
    if nums:
        values = sorted(_to_rmb(num, unit) for num, unit in nums)
        return values[0], values[-1]
    return None, None


def _assess_scale_adaptation(result: dict) -> tuple[int, str, str]:
    """Return (category, label, reason) for scale fit vs fund preferences."""
    industry = f"{result.get('industry', '')} {result.get('business_type', '')}".strip()
    revenue_text = str(
        result.get("revenue_range") or result.get("revenue_estimate") or result.get("revenue_est") or ""
    ).strip()
    low, high = _parse_revenue_bounds(revenue_text)
    scale_evidence_level = int(result.get("scale_evidence_level", 2) or 2)
    strong_evidence = scale_evidence_level >= 3 or str(result.get("revenue_evidence", "")).upper() in {"A", "B"}

    # Prefer a simple, stable band: food/OEM around 1-3亿; cleaning/disinfection can be smaller but should not be tiny.
    if any(k in industry for k in ("清洁", "消毒")):
        target_min, target_max = 2e7, 3e8
    elif any(k in industry for k in ("食品", "预制", "调味", "饮料", "OEM", "代工", "加工", "肉制", "餐饮")):
        target_min, target_max = 1e8, 3e8
    else:
        target_min, target_max = 5e7, 3e8

    def _fmt(v: float | None) -> str:
        if v is None:
            return "未知"
        if v >= 1e8:
            x = v / 1e8
            return f"{x:.1f}亿".rstrip("0").rstrip(".") if x % 1 else f"{int(x)}亿"
        return f"{v/1e4:.0f}万"

    # Too large
    if low is not None and low > target_max:
        if strong_evidence:
            return 1, "明确证据证明规模过大", f"营收约{revenue_text}，明显高于基金偏好的{_fmt(target_min)}-{_fmt(target_max)}区间。"
        return 2, "推断规模大概率过大", f"营收约{revenue_text}，大概率高于基金偏好的{_fmt(target_min)}-{_fmt(target_max)}区间。"

    # Too small
    if high is not None and high < target_min:
        if strong_evidence:
            return 3, "明确证据证明规模过小", f"营收约{revenue_text}，明显低于基金偏好的{_fmt(target_min)}-{_fmt(target_max)}区间。"
        return 4, "推断规模大概率过小", f"营收约{revenue_text}，大概率低于基金偏好的{_fmt(target_min)}-{_fmt(target_max)}区间。"

    return 5, "未见规模不合适证据", (
        f"营收约{revenue_text or '未明确'}，区间未整体落在基金偏好范围之外，"
        "暂不据此认定规模不适配。"
    )


def _advance_recommendation(grade: str, scale_category: int, scale_reason: str) -> tuple[str, str]:
    # Category 1/2 = explicit scale mismatch → reject regardless of grade
    if scale_category in {1, 2}:
        return "不建议推进", f"{scale_reason} 规模不适配，不建议优先推进。"

    # Category 3 = DEFINITE evidence of mismatch (too small or too large)
    if scale_category == 3:
        if grade == "A":
            return "不建议推进", f"{scale_reason} 明确证据证明规模不适配，即使企业质量强也不足以弥补。"
        if grade == "B":
            return "不建议推进", f"{scale_reason} 明确证据证明规模不适配。"
        return "不建议推进", f"{scale_reason} 明确证据证明规模不适配且质量弱。"

    # Category 4 = INFERRED mismatch (too small or too large) → 观察 for A/B
    if scale_category == 4:
        if grade == "A":
            return "观察", f"{scale_reason} 虽然企业质量强，但推断规模不适配，建议观察补证。"
        if grade == "B":
            return "观察", f"{scale_reason} 企业质量尚可，但推断规模不适配，先观察。"
        return "不建议推进", f"{scale_reason} 企业质量不足以抵消规模不适配。"

    # Category 5 = no evidence of mismatch → 推进 for A/B
    if grade == "A":
        return "推进", f"{scale_reason} 且企业质量强，建议推进。"
    if grade == "B":
        return "推进", f"{scale_reason} 未见规模不适配证据，企业质量尚可，可推进。"
    return "不建议推进", f"{scale_reason} 但企业质量未达到推进门槛。"


def _assess_transaction_obstacles(
    result: dict, evidence: list[int], red_numbers: set[int]
) -> tuple[str, str, bool, bool]:
    """Return transaction obstacle label, reason, hard-block flag and verification flag.

    Now uses qualitative transaction_feasibility instead of numeric scores[7].
    """
    feasibility = str(result.get("transaction_feasibility", "")).strip()
    trade_evidence = evidence[6] if len(evidence) >= 7 else 2  # now index 6 = 规模推断
    supplied_reason = str(result.get("transaction_feasibility_reason", "")).strip()
    ai_says_no_obstacle = bool(re.search(
        r"未[见发].{0,6}(?:明显|明确|实质).{0,6}(?:交易障碍|执行障碍|审批障碍|出售障碍)"
        r"|不存在.{0,6}(?:明显|明确).{0,6}(?:交易障碍|执行障碍)"
        r"|交易障碍.{0,4}(?:低|极低|无|未见)"
        r"|不构成.{0,8}(?:交易障碍|执行障碍|出售障碍)"
        r"|股权.{0,8}(?:清晰|简单|集中)"
        r"|无.{0,8}(?:诉讼|处罚|行政处罚|国资|外资|复杂因素)"
        r"|未公开出售意愿.{0,12}(?:正常|属正常)",
        supplied_reason,
    ))
    obstacle_check_text = re.sub(
        r"无.{0,12}(?:国资|国企|外资|上市公司|复杂因素|诉讼|处罚|股权冻结|出质)"
        r"|未[见发有].{0,12}(?:国资|国企|外资|上市公司|复杂因素|诉讼|处罚|股权冻结|出质)"
        r"|不构成.{0,8}(?:交易障碍|执行障碍|出售障碍)"
        r"|未公开出售意愿.{0,12}(?:正常|属正常)",
        "",
        supplied_reason,
    )
    explicit_obstacle = bool(re.search(
        r"上市公司|国有|国企|外资|集团体系|母公司.{0,12}(?:不出售|难拆分)|"
        r"明确.{0,12}(?:不出售|拒绝出售|无出售意愿)|股权.{0,10}(?:复杂|冻结|出质)|"
        r"审批障碍|拆分困难",
        obstacle_check_text,
    ))
    if red_numbers & _CRITICAL_RED_FLAGS:
        return "存在重大交易障碍", "存在监管合规或主体清晰度关键红旗。", True, False
    if 7 in red_numbers:
        return "存在明显交易障碍", "存在过度依赖老板个人关系的红旗，交易交接风险高。", True, False
    ai_obstacle = result.get("has_transaction_obstacle")
    if isinstance(ai_obstacle, bool) and ai_obstacle:
        return "存在明显交易障碍", supplied_reason, True, False
    if feasibility == "evidence_no":
        return "明确交易不可行", supplied_reason or "有明确证据表明交易不可推进。", True, False
    if feasibility == "inferred_no":
        # Inferred obstacle: warn but don't hard-block
        reason = supplied_reason or "推断存在交易障碍，需尽调验证。"
        return "交易存在推断障碍", reason, False, True
    if feasibility == "evidence_yes":
        return "交易可行性已验证", supplied_reason or "有明确证据表明交易可推进。", False, False
    if feasibility == "inferred_yes":
        return "未见明显交易障碍", supplied_reason or "公开信息中未见明确交易执行障碍。", False, False
    # Fallback: if no feasibility field, use old heuristic
    if explicit_obstacle and trade_evidence >= 3 and not ai_says_no_obstacle:
        return "存在明显交易障碍", supplied_reason, True, False
    if not explicit_obstacle:
        reason = (
            "公开信息中未见明确交易执行障碍；非上市中小企业通常不会公开出售意愿，"
            "股权细节与接洽意愿作为推进后的核实事项。"
        )
        return "未见明显交易障碍（接触时核实）", reason, False, False
    return "交易障碍线索待核实", supplied_reason, False, True


def _normalize_rr_conclusion(result: dict, rr_score: int, quick: bool) -> None:
    """Keep the report wording consistent with the programmatic RR score."""
    if rr_score >= 5:
        label = "高"
        explanation = "已有合同、重复收入占比或多年连续经营数据等强证据支持持续性收入判断。"
    elif rr_score == 4:
        label = "高"
        explanation = "已有多个公开经营事实支持持续采购或稳定供货；对中小企业不以披露合同期限、续约率或排他条款为前提。"
    elif rr_score == 3:
        label = "中"
        explanation = "产品或客户场景具备复购特征，但尚缺少已发生的持续供货、重复订单或连续生产佐证。"
    else:
        label = "低"
        explanation = "当前公开证据不足以支持稳定重复收入，或业务更偏一次性交易。"
    result["rr_score"] = rr_score
    result["rr"] = label
    result["recurring_revenue"] = label
    result["rr_reason"] = explanation
    if quick:
        return
    result["sec3_verdict"] = f"{label}（{rr_score}分）"
    result["sec3_summary"] = f"{label}（{rr_score}分）｜{explanation}"
    original_content = str(result.get("sec3_content", "")).strip()
    original_content = re.sub(
        r"(判断\s*[：:]\s*)(?:较高|高|中|较低|低|暂无法确认)",
        rf"\g<1>{label}",
        original_content,
        count=1,
    )
    lead = f"<p><strong>程序评分说明：</strong>RR得分为{rr_score}分，判断为{label}。{explanation}</p>"
    result["sec3_content"] = lead + original_content


def _enforce_supported_rr_conclusion(
    result: dict, scores: list[int], red_numbers: set[int], quick: bool
) -> None:
    """Resolve a model output conflict when its documented RR analysis clearly supports 4 points."""
    if quick or scores[1] >= 4 or 1 in red_numbers:
        return
    text = " ".join(
        str(result.get(key, ""))
        for key in ("sec3_verdict", "sec3_summary", "sec3_content")
    )
    if not re.search(r"(?:评分|应得|得分|给予)\s*(?:为|[:：])?\s*4\s*分|评分\s*4\s*分", text):
        return
    support_groups = (
        r"持续供货|稳定供货|连续多年|反复采购|重复订单",
        r"提前排产|敲定.{0,12}订单|周期性订单|节日期间.{0,12}订单",
        r"具名.{0,8}客户|30余家|盒马|全家|联华|小杨生煎|宝立食品",
        r"每月.{0,12}(?:吨|产量)|月产|产量.{0,12}(?:吨|增加)",
    )
    supports = sum(bool(re.search(pattern, text)) for pattern in support_groups)
    has_source_support = bool(
        re.search(r"<sup|\[\d+\]|政府.{0,8}报道|招股书|证据等级\s*[：:]?\s*[AB]", text)
    )
    if supports >= 2 and has_source_support:
        scores[1] = 4
        result["rr_score_note"] = "RR正文明确给出4分且包含持续经营公开证据，程序已同步校正结构化分数。"


def _enforce_core_industry_fit(
    result: dict, scores: list[int], evidence: list[int], red_numbers: set[int], quick: bool
) -> None:
    """Apply a stable score floor when public evidence clearly identifies a core target segment."""
    if 10 in red_numbers or not evidence or evidence[0] < 3:
        return
    text = " ".join(
        str(result.get(key, ""))
        for key in (
            "industry", "business_type", "sec1_summary", "sec1_content",
            "sec2_summary", "sec2_content", "analysis",
        )
    )
    core_reason = ""
    if (
        re.search(r"速冻|冷冻|冻品", text)
        and re.search(r"OEM|ODM|代工", text, re.IGNORECASE)
        and re.search(r"食品|肉制|调理|预制|猪排", text)
        and not re.search(r"鲜切|鲜果|果切|即食鲜|生鲜配送|冷藏鲜|冷鲜", text)
    ):
        core_reason = "公开证据明确显示为速冻/冷冻食品的B端OEM/ODM业务，归入冻品OEM核心赛道。"
    elif (
        re.search(r"饮料|饮品", text)
        and re.search(r"OEM|ODM|代工", text, re.IGNORECASE)
        and not re.search(r"鲜切|鲜果|果切|水果|冻品|速冻|冷冻", text)
        and not re.search(r"(?:不属于|非(?!常)|不是|不包括|不涉及).{0,30}(?:饮料|饮品)", text)
    ):
        core_reason = "公开证据明确显示为饮料OEM/ODM业务，归入核心赛道。"
    elif re.search(r"调味|酱料", text) and re.search(r"OEM|ODM|代工", text, re.IGNORECASE):
        core_reason = "公开证据明确显示为调味品OEM/ODM业务，归入核心赛道。"
    elif re.search(r"商用|工业|餐饮|机构", text) and re.search(r"清洗|清洁|消毒", text):
        core_reason = "公开证据明确显示为商用清洗消毒业务，归入核心赛道。"
    if not core_reason:
        return
    # 只在AI评分偏低时拉高，不覆盖AI正确判断的4分
    if scores[0] >= 4:
        return
    scores[0] = 5
    result["industry_score_note"] = core_reason
    if quick:
        return
    original_summary = str(result.get("sec1_summary", "")).strip()
    original_content = str(result.get("sec1_content", "")).strip()
    result["sec1_summary"] = f"核心赛道（5分）｜{original_summary}" if original_summary else "核心赛道（5分）"
    result["sec1_content"] = (
        f"<p><strong>程序评分说明：</strong>行业匹配度为5分。{core_reason}</p>{original_content}"
    )





def _build_exec_summary(result: dict, verdict_text: str) -> str:
    industry = str(result.get("industry", "该公司")).strip() or "该公司"
    business_type = str(result.get("business_type", "")).strip()
    grade = str(result.get("grade", "C")).strip()
    total = result.get("total_score", result.get("total", "待确认"))
    total_text = f"{total:.1f}" if isinstance(total, (int, float)) else str(total)
    quality_text = str(result.get("business_quality", "中")).strip()
    advance_text = str(result.get("advance_recommendation", "不建议推进")).strip()
    rr = str(result.get("recurring_revenue", "中")).strip()
    revenue = str(result.get("revenue_range", result.get("revenue_estimate", "待推算"))).strip()
    revenue_evidence = str(result.get("revenue_evidence", "C")).strip()
    add_on = str(result.get("sec7_verdict", "待判断")).strip()
    red_count = result.get("red_flags_count", result.get("red_count", 0))
    red_severity = str(result.get("red_flags_severity", result.get("red_severity", "无"))).strip()
    decision_note = str(result.get("decision_note", "")).strip().rstrip("。；;")
    if not decision_note:
        decision_note = "程序复算后生成最终结论"

    parts = [
        f"{industry}{('（' + business_type + '）') if business_type else ''}当前企业质量分{total_text}，ABCD评级{grade}，推进建议为{advance_text}。",
        f"企业质量判断为{quality_text}，规模适配为{result.get('scale_fit_label', '未评估')}，交易障碍为{result.get('transaction_obstacle_label', '未评估')}。",
        f"Recurring Revenue判断为{rr}，规模推断约{revenue}（{revenue_evidence}）。",
        f"角色定位判断偏{add_on}，红旗{red_count}项，严重度{red_severity}。",
        f"综合来看，{decision_note}。",
    ]
    return " ".join(parts)


def _self_critique(analysis: dict, company_name: str) -> dict:
    """Lightweight consistency check after scoring — catches internal contradictions.

    Checks:  A-grade vs small-scale, high-RR vs low-evidence, verdict mismatch.
    If a contradiction is found, adds a ``critique_note`` field and may downgrade
    the advance_recommendation from 推进 to 观察.
    """
    grade = str(analysis.get("grade", "")).upper()
    scale_label = str(analysis.get("scale_fit_label", "") or "")
    scale_reason = str(analysis.get("scale_fit_reason", "") or "")
    evidence = str(analysis.get("evidence_confidence", "C")).upper()
    advance = str(analysis.get("advance_recommendation", "") or "")
    rr_label = str(analysis.get("rr", "") or "")
    total = float(analysis.get("total_score", 0) or 0)
    notes = []

    # Contradiction 1: A-grade + clear small-scale evidence → should not be 推进
    small_scale_kws = ["过小", "不适配", "低于", "不达", "明显低于"]
    is_clearly_small = any(kw in scale_label or kw in scale_reason for kw in small_scale_kws)
    if grade == "A" and advance == "推进" and is_clearly_small:
        analysis["advance_recommendation"] = "观察"
        analysis["advance_recommendation_reason"] = (
            f"[自动审查修正] 企业质量A级但规模明确不适配，推进→观察。{analysis.get('advance_recommendation_reason', '')}"
        )
        notes.append("A级+规模明确过小仍建议推进，自动降为观察")

    # Contradiction 2: high RR (≥4) but evidence is D → suspect confidence
    if rr_label in ("高", "中高") and evidence in ("D",):
        notes.append("RR判断偏高但证据可信度为D，可能存在乐观偏差")

    # Contradiction 3: total >= 4.0 but verdict is REJECT → check if scale is the reason
    if total >= 4.0 and str(analysis.get("verdict_type", "")) == "reject":
        if not is_clearly_small:
            notes.append("企业质量分≥4.0却判定REJECT，请核实排据是否充分")

    if notes:
        analysis["critique_notes"] = " | ".join(notes)
        print(f"[Critique] {company_name}: {analysis['critique_notes']}")

    return analysis


def _apply_screening_score(result: dict, quick: bool = False) -> dict:
    """Recalculate decision fields from AI observations instead of trusting its verdict."""
    score_key = "scores" if quick else "radar_scores"
    raw_scores = result.get(score_key)
    scores = [_bounded_score(value) for value in raw_scores] if isinstance(raw_scores, list) else []
    scores = (scores + [2] * 5)[:5]
    raw_evidence = result.get("evidence_levels")
    revenue_evidence = str(result.get("revenue_evidence", "")).upper()
    scale_evidence_level = {"A": 4, "B": 3, "C": 2, "D": 1}.get(revenue_evidence)
    if scale_evidence_level is None:
        scale_evidence_level = (
            _bounded_score(raw_evidence[6], 2)
            if not quick and isinstance(raw_evidence, list) and len(raw_evidence) >= 7
            else 2
        )
    red_flags = _extract_red_flags(result)
    red_numbers = {int(key) for key, marker in red_flags.items() if marker == "🔴"}
    yellow_count = sum(marker == "🟡" for marker in red_flags.values())
    for red_number in red_numbers:
        for dimension_index, cap in _SCORE_CAPS.get(red_number, {}).items():
            scores[dimension_index] = min(scores[dimension_index], cap)

    evidence, evidence_average, confidence = _dimension_evidence(result, quick)
    # _enforce_core_industry_fit removed — regex-guessing industry from free text is unreliable
    _enforce_supported_rr_conclusion(result, scores, red_numbers, quick)
    # Compounder lens: RR + market position drive moat, stickiness/bizmodel secondary, industry as gate.
    # Weights: RR 30%, Market Position 25%, Business Model 20%, Stickiness 15%, Industry 10%
    quality_weighted = (
        scores[0] * 0.10 + scores[1] * 0.30 + scores[2] * 0.15
        + scores[3] * 0.20 + scores[4] * 0.25
    )
    quality_score = round(quality_weighted, 1)
    full_score = quality_score
    critical_flags = sorted(red_numbers & _CRITICAL_RED_FLAGS)
    if critical_flags or len(red_numbers) >= 3:
        severity = "严重"
    elif red_numbers:
        severity = "中等"
    elif yellow_count:
        severity = "轻微"
    else:
        severity = "无"

    result[score_key] = scores
    result["full_score"] = full_score
    result["red_flags"] = red_flags
    result["red_flags_count" if not quick else "red_count"] = len(red_numbers)
    result["red_flags_severity" if not quick else "red_severity"] = severity
    result["evidence_levels_scoring"] = evidence
    result["evidence_score"] = evidence_average
    result["evidence_confidence"] = confidence
    result["evidence_confidence_note"] = _evidence_confidence_note(confidence)
    result["scale_evidence_level"] = scale_evidence_level
    result["score_method"] = "企业质量5维加权（不含角色定位/交易可行性）"
    scale_category, scale_label, scale_reason = _assess_scale_adaptation(result)
    obstacle_label, obstacle_reason, obstacle_block, obstacle_verify = _assess_transaction_obstacles(
        result, evidence, red_numbers
    )
    grade = "A" if quality_score >= 3.8 else ("B" if quality_score >= 3.0 else ("C" if quality_score >= 2.0 else "D"))
    advance_recommendation, advance_reason = _advance_recommendation(grade, scale_category, scale_reason)
    feasibility = str(result.get("transaction_feasibility", "")).strip()
    feasibility_reason = str(result.get("transaction_feasibility_reason", "")).strip()
    if severity == "严重":
        advance_recommendation = "不建议推进"
        advance_reason = f"存在严重红旗风险，当前不建议推进。{scale_reason}"
    elif scores[0] <= 2:
        advance_recommendation = "不建议推进"
        advance_reason = f"行业匹配度不足，不符合当前筛选方向。{scale_reason}"
    elif feasibility == "evidence_no":
        advance_recommendation = "不建议推进"
        advance_reason = f"交易不可行（有明确证据）：{feasibility_reason or '存在明确交易障碍。'}"
    elif obstacle_block:
        advance_recommendation = "不建议推进"
        advance_reason = f"{obstacle_label}：{obstacle_reason} 因此当前不建议推进。"
    elif advance_recommendation == "推进" and feasibility == "inferred_no":
        # Don't downgrade, but add warning
        advance_reason = (
            f"{scale_reason} 企业质量达到推进门槛，建议推进接触；"
            f"⚠ 交易可行性存在推断障碍，需尽调验证。{feasibility_reason}"
        )
    elif advance_recommendation == "推进" and obstacle_verify:
        advance_reason = (
            f"{scale_reason} 企业质量达到推进门槛，建议推进接触；"
            f"同时核查{obstacle_label}：{obstacle_reason}"
        )
    elif advance_recommendation == "推进" and confidence not in {"A", "B"}:
        advance_recommendation = "观察"
        advance_reason = (
            f"{scale_reason} 企业质量达到推进门槛，但证据可信度不足（{confidence}），"
            "需补充关键证据后再决定是否推进。"
        )
    elif advance_recommendation == "推进" and feasibility == "evidence_yes":
        advance_reason = f"{advance_reason} + 交易可行性已验证。"
    if quality_score >= 3.8:
        business_quality = "高"
    elif quality_score >= 3.0:
        business_quality = "中"
    else:
        business_quality = "低"
    result["business_quality_score"] = quality_score
    result["business_quality"] = business_quality
    result["scale_fit_category"] = scale_category
    result["scale_fit_label"] = scale_label
    result["scale_fit_reason"] = scale_reason
    result["transaction_obstacle_label"] = obstacle_label
    result["transaction_obstacle_reason"] = obstacle_reason
    _normalize_rr_conclusion(result, scores[1], quick)
    result["advance_recommendation"] = advance_recommendation
    result["advance_recommendation_reason"] = advance_reason
    result["dealability"] = {"evidence_yes": "高", "inferred_yes": "中", "inferred_no": "低", "evidence_no": "低"}.get(result.get("transaction_feasibility", "inferred_yes"), "不明")
    if advance_recommendation == "推进":
        verdict = "PASS"
    elif advance_recommendation == "观察":
        verdict = "WATCH"
    else:
        verdict = "REJECT"

    result["total_score"] = quality_score
    result["total"] = quality_score
    result["quality_score"] = quality_score
    result["grade"] = grade
    result["decision_note"] = advance_reason
    verdict_text = {
        "PASS": "推进",
        "WATCH": "观察",
        "REJECT": "不建议推进",
    }[verdict]
    result["exec_summary"] = _build_exec_summary(result, verdict_text)
    if quick:
        result["verdict"] = verdict
    else:
        result["verdict_type"] = "pass" if verdict == "PASS" else ("watch" if verdict == "WATCH" else "reject")
        result["sec9_verdict"] = advance_recommendation
        result["sec9_summary"] = f"ABCD评级：{grade} / 推进建议：{advance_recommendation}（程序复算）"
        result["sec9_reason"] = advance_reason
    return result


def _fix_red_flags_html(fragment: str) -> str:
    if not fragment or not fragment.strip():
        return ""

    def _rewrite_citations(text: str) -> str:
        text = re.sub(
            r'(?<!>)\[(\d+)\s*,\s*(\d+)\]',
            lambda ref: (
                f'<a class="cite external" href="#ref-{ref.group(1)}">[{ref.group(1)}]</a>'
                f'<a class="cite external" href="#ref-{ref.group(2)}">[{ref.group(2)}]</a>'
            ),
            text,
        )
        text = re.sub(
            r'(?<!>)\[(\d+)\]',
            lambda ref: f'<a class="cite external" href="#ref-{ref.group(1)}">[{ref.group(1)}]</a>',
            text,
        )
        return text

    # 14项PE红旗 → 对应分析章节
    _flag_sec_map = {
        1: "sec3",   # 一次性项目收入为主 → Recurring Revenue
        2: "sec4",   # 必须持续开发新客户 → 客户结构
        3: "sec5",   # 纯贸易、低壁垒 → 竞争位置
        4: "sec7",   # 重资产/重库存/重CAPEX → 规模成熟度
        5: "sec5",   # 原材料价格高度暴露 → 竞争位置
        6: "sec4",   # 依赖单一大客户 → 客户结构
        7: "sec4",   # 过度依赖老板个人关系 → 客户结构
        8: "sec5",   # 技术替代风险高 → 竞争位置
        9: "sec8",   # 监管合规风险高 → 红旗核查
        10: "sec5",  # 行业边界模糊 → 竞争位置
        11: "sec6",  # 平台能力不足 → Platform/Add-on
        12: "sec1",  # 财务规范性/主体清晰度可疑 → 公司身份核验
        13: "sec4",  # 客户质量不足 → 客户结构
        14: "sec1",  # 公开信息极度稀缺 → 公司身份核验
    }
    _flag_sec_label = {
        "sec1": "身份",
        "sec3": "RR",
        "sec4": "客户",
        "sec5": "竞争",
        "sec6": "平台",
        "sec7": "交易",
        "sec8": "证据",
    }

    def _make_flag_card(m):
        emoji = m.group(1)
        num = int(m.group(2))
        content = m.group(3).strip()
        content = _rewrite_citations(content)
        # 常见冲突保护：如果"纯贸易、低壁垒"这一条同时出现专利/研发/专业服务/解决方案等强正面信号，
        # 不直接写成"存在"，而改为"待核实"，避免与正文分析打架。
        if num == 3 and emoji == "🔴":
            positive_signals = ("专利", "研发", "技术", "专业服务", "解决方案", "创新", "系统", "先进")
            if ("纯贸易" in content or "低壁垒" in content) and any(token in content for token in positive_signals):
                emoji = "🟡"
                content = re.sub(r'纯贸易、低壁垒：存在。?', '纯贸易、低壁垒：待核实。', content)
                content = re.sub(r'纯贸易、低壁垒：存在', '纯贸易、低壁垒：待核实', content)
                content = re.sub(r'存在。证据：', '待核实。证据：', content, count=1)
        if num == 12 and emoji == "🔴" and _is_routine_sme_information_gap(content):
            emoji = "🟡"
            content = re.sub(r"：存在。?", "：待核实。", content, count=1)
        if num == 13 and emoji == "🔴" and _has_premium_customer_evidence(content):
            emoji = "🟡"
            content = re.sub(r"：存在。?", "：待核实。", content, count=1)
        if num == 6 and emoji == "🔴" and _describes_multiple_customers(content):
            emoji = "🟡"
            content = re.sub(r"：存在。?", "：待核实。", content, count=1)
        if num == 7 and emoji == "🔴" and _has_founder_dependency_mitigation(content):
            emoji = "🟡"
            content = re.sub(r"：存在。?", "：部分存在，已见缓解信号。", content, count=1)
        if emoji == "🟢":
            cls = "flag-clear"
            severity = "-"
            onclick = ""
            link = ""
        elif emoji == "🟡":
            cls = "flag-suspect"
            severity = "中"
            sec = _flag_sec_map.get(num, "sec8")
            onclick = f' onclick="scrollToSection(\'{sec}\')"'
            link = f'<span class="flag-link">&#8594; &#167;{num} {_flag_sec_label.get(sec, "分析")}</span>'
        else:
            cls = "flag-found"
            severity = "高"
            sec = _flag_sec_map.get(num, "sec8")
            onclick = f' onclick="scrollToSection(\'{sec}\')"'
            link = f'<span class="flag-link">&#8594; &#167;{num} {_flag_sec_label.get(sec, "分析")}</span>'
        return (
            f'<div class="flag-item {cls}{" clickable" if onclick else ""}"{onclick}>'
            f'<span class="flag-icon">{emoji}</span>'
            f'<span>{num}. {content}</span>'
            f'<span class="flag-severity">{severity}</span>'
            f'{link}'
            f'</div>'
        )

    pattern = r'<div>\s*([🟢🟡🔴])\s*(\d+)[.、：:](.*?)</div>'
    fragment = re.sub(pattern, _make_flag_card, fragment, flags=re.DOTALL)
    # Models sometimes return a mix of ready-made cards and raw divs. Rewrite
    # citations after card normalization so raw items containing [N] still match.
    return _rewrite_citations(fragment)


def _fix_references_html(fragment: str, search_text: str = "") -> str:
    if not fragment or 'href="#"' not in fragment:
        return fragment
    url_map = {}
    if search_text:
        for m in re.finditer(r'\[(\d+)\]\s*[^\n]+\nURL:\s*(https?://[^\s\n]+)', search_text):
            num = int(m.group(1))
            url = m.group(2).strip()
            if num not in url_map:
                url_map[num] = url
    if not url_map:
        return html

    def _replace_href(m):
        before = fragment[:m.start()]
        ref_matches = list(re.finditer(r'<li id="ref-(\d+)"', before))
        if ref_matches:
            ref_num = int(ref_matches[-1].group(1))
            real_url = html.escape(_safe_link(url_map.get(ref_num, "#")), quote=True)
            return f'<a href="{real_url}" target="_blank" rel="noopener"'
        return m.group(0)

    return re.sub(r'<a href="#"', _replace_href, fragment)


def fill_report_template(
    company_name: str, analysis_json: dict, search_text: str = "", verified_refs_html: str | None = None
) -> str:
    """用分析结果填充报告模板"""
    template_path = BASE_DIR / "report_template.html"
    template = template_path.read_text(encoding="utf-8")

    # 从分析 JSON 中提取各字段，填充模板占位符
    grade = analysis_json.get("grade", "C")
    if grade not in {"A", "B", "C", "D"}:
        grade = "C"
    grade_class = {"A": "a", "B": "b", "C": "c", "D": "d"}.get(grade, "c")

    # 预处理红旗HTML：将AI生成的裸<div>格式转换为可点击的flag-item卡片
    red_flags_raw = _sanitize_html_fragment(analysis_json.get("red_flags_html", ""))
    red_flags_fixed = _fix_red_flags_html(red_flags_raw)
    verdict_type = analysis_json.get("verdict_type", "watch")
    if verdict_type not in {"pass", "watch", "reject"}:
        verdict_type = "watch"
    references_html = (
        analysis_json.get("references_html", "") if verified_refs_html is None else verified_refs_html
    )

    replacements = {
        "{{公司名称}}": _escape_text(company_name),
        "{{公司名称JSON}}": json.dumps(company_name, ensure_ascii=True).replace("</", "<\\/"),
        "{{行业归属}}": _escape_text(analysis_json.get("industry", "待确认")),
        "{{业务类型}}": _escape_text(analysis_json.get("business_type", "待确认")),
        "{{报告日期}}": datetime.now().strftime("%Y年%m月%d日"),
        "grade-{{评级字母}}": f"grade-{grade_class}",
        "{{评级字母}}": grade,
        "{{执行摘要内容，150-300字。包含：是否值得优先接触、Platform/Add-on判断、recurring revenue大致判断、规模推断（含区间和证据等级）、最关键吸引点、最关键风险、建议。}}": _escape_text(analysis_json.get("exec_summary", "")),
        "{{总评分}}": _escape_text(analysis_json.get("total_score", 3.0)),
        "{{规模推断}}": _escape_text(analysis_json.get("revenue_estimate", "待推算")),
        "{{证据等级}}": _escape_text(analysis_json.get("revenue_evidence", "C")),
        "{{revenue判断}}": _escape_text(analysis_json.get("recurring_revenue", "待确认")),
        "{{红旗数量}}": _escape_text(analysis_json.get("red_flags_count", 0)),
        "{{严重度}}": _escape_text(analysis_json.get("red_flags_severity", "无")),
        "{{评分可信度}}": _escape_text(analysis_json.get("evidence_confidence", "C")),
        "{{可信度说明}}": _escape_text(analysis_json.get("evidence_confidence_note", "")),
        "{{评分方法}}": _escape_text(analysis_json.get("score_method", "程序加权复算")),
        "{{决策规则说明}}": _escape_text(analysis_json.get("decision_note", "")),
        "{{雷达与评分说明}}": _escape_text(
            analysis_json.get(
                "radar_score_note",
                "雷达图展示五个企业质量原始维度分数；角色定位与交易可行性为定性标签，独立呈现，不参与评分。",
            )
        ),
        "{{红旗清单HTML}}": red_flags_fixed,
        "{{一句话摘要：行业归属+成立时间+实控人+规模概况}}": _escape_text(analysis_json.get("sec1_summary", "")),
        "{{行业归属判断、主体性质、核心业务、主体清晰度、混淆风险、企查查基础信息摘要表}}": _sanitize_html_fragment(analysis_json.get("sec1_content", "")),
        "{{核心业务+品牌+B端/C端}}": _escape_text(analysis_json.get("sec2_summary", "")),
        "{{核心产品/服务、次要业务、主要收入模式、高复购部分、一次性业务风险}}": _sanitize_html_fragment(analysis_json.get("sec2_content", "")),
        "{{判断：高/中/低+一句话理由}}": _escape_text(analysis_json.get("sec3_summary", "")),
        "{{高/中/低/暂无法确认}}": _escape_text(analysis_json.get("sec3_verdict", "暂无法确认")),
        "{{重复性来源、真重复还是伪重复、客户是否持续采购、增长驱动力、哪些是公开资料支持哪些是推断}}": _sanitize_html_fragment(analysis_json.get("sec3_content", "")),
        "{{客户质量+切换成本+关键风险}}": _escape_text(analysis_json.get("sec4_summary", "")),
        "{{客户类型、需求本质、需求频率、切换成本判断、客户质量与留存线索}}": _sanitize_html_fragment(analysis_json.get("sec4_content", "")),
        "{{市场地位判断+一句话描述}}": _escape_text(analysis_json.get("sec5_summary", "")),
        "{{市场地位详细分析HTML}}": _sanitize_html_fragment(analysis_json.get("sec5_content", "")),
        "{{市场定位+护城河判断}}": _escape_text(analysis_json.get("sec6_summary", "")),
        "{{市场定位、差异化来源、局部护城河判断、公开资料支持强度}}": _sanitize_html_fragment(analysis_json.get("sec6_content", "")),
        "{{判断结果+一句话理由}}": _escape_text(analysis_json.get("sec7_summary", "")),
        "{{Platform/Add-on/Borderline/非目标}}": _escape_text(_role_display_label(analysis_json.get("role_recommendation", ""), analysis_json.get("sec7_verdict", "待判断"))),
        "{{判断理由}}": _sanitize_html_fragment(analysis_json.get("sec7_content", "")),
        "{{营收区间+规模适配+交易障碍}}": _escape_text(analysis_json.get("sec8_summary", "")),
        "{{X-Y 亿}}": _escape_text(analysis_json.get("revenue_range", "待推算")),
        "{{推算方法}}": _escape_text(analysis_json.get("revenue_method", "多维度交叉推算")),
        "{{营收推断过程HTML}}": _sanitize_html_fragment(analysis_json.get("revenue_process_html", "")),
        "{{推算明细表行}}": _sanitize_html_fragment(analysis_json.get("sec8_table", "")),
        "{{判断}}": _escape_text(analysis_json.get("maturity", "待判断")),
        "{{规模适配}}": _escape_text(
            f"{analysis_json.get('scale_fit_category', '')}｜{analysis_json.get('scale_fit_label', '')}".strip("｜")
        ),
        "{{规模适配理由}}": _escape_text(analysis_json.get("scale_fit_reason", "")),
        "{{交易障碍}}": _escape_text(analysis_json.get("transaction_obstacle_label", "未评估")),
        "{{交易障碍理由}}": _escape_text(analysis_json.get("transaction_obstacle_reason", "")),
        "{{推进建议}}": _escape_text(analysis_json.get("advance_recommendation", "不建议推进")),
        "{{推进建议理由}}": _escape_text(analysis_json.get("advance_recommendation_reason", "")),
        "{{各级证据数量统计}}": _escape_text(analysis_json.get("sec9_summary", "")),
        "{{证据等级表行}}": _sanitize_html_fragment(analysis_json.get("sec9_table", "")),
        "{{评级档位+一句话结论}}": _escape_text(analysis_json.get("sec10_summary", "")),
        "{{A：优先深挖 / B：值得继续看 / C：保留观察 / D：暂不优先}}": _escape_text(analysis_json.get("sec10_verdict", "")),
        "{{pass/watch/reject}}": verdict_type,
        "{{企业质量判断}}": _escape_text(analysis_json.get("business_quality", "中")),
        "{{最终建议理由}}": _escape_text(analysis_json.get("sec10_reason", "")),
        "{{关键问题}}": _escape_text(analysis_json.get("key_question", "")),
        "{{最佳切入点}}": _escape_text(analysis_json.get("approach", "")),
        "{{问题清单}}": _sanitize_html_fragment(analysis_json.get("next_steps_html", "")),

        # ---- Card-based section placeholders (maps to template cards) ----
        # Card 1: 行业匹配度 (identity + product merged)
        "{{card1_industry_summary}}": _escape_text(
            f"{analysis_json.get('industry', '')} | {analysis_json.get('business_type', '')}"
        ),
        "{{card1_industry_content}}": (
            _sanitize_html_fragment(analysis_json.get("sec1_content", ""))
            + "<hr>" + _sanitize_html_fragment(analysis_json.get("sec2_content", ""))
        ),

        # Card 2: Recurring Revenue
        "{{card2_rr_summary}}": _escape_text(analysis_json.get("sec3_summary", "")),
        "{{card2_rr_verdict}}": _escape_text(analysis_json.get("sec3_verdict", "暂无法确认")),
        "{{card2_rr_content}}": _sanitize_html_fragment(analysis_json.get("sec3_content", "")),

        # Card 3: 市场地位 (market position + competition merged)
        "{{card3_market_summary}}": _escape_text(analysis_json.get("sec5_summary", "")),
        "{{card3_market_content}}": (
            '<div class="sub-section"><h4>市场地位分析</h4>'
            + _sanitize_html_fragment(analysis_json.get("sec5_content", ""))
            + '</div><hr><div class="sub-section"><h4>竞争位置与差异化</h4>'
            + _sanitize_html_fragment(analysis_json.get("sec6_content", ""))
            + '</div>'
        ),

        # Card 4: 客户粘性
        "{{card4_stickiness_summary}}": _escape_text(analysis_json.get("sec4_summary", "")),
        "{{card4_stickiness_content}}": _sanitize_html_fragment(analysis_json.get("sec4_content", "")),

        # Card 5: 商业模式
        "{{card5_bizmodel_summary}}": _escape_text(analysis_json.get("sec2_summary", "")),
        "{{card5_bizmodel_content}}": _sanitize_html_fragment(analysis_json.get("sec2_content", "")),

        # Card 6: Platform / Add-on
        "{{card6_platform_summary}}": _escape_text(analysis_json.get("sec7_summary", "")),
        "{{card6_platform_verdict}}": _escape_text(analysis_json.get("sec7_verdict", "待判断")),
        "{{card6_platform_content}}": _sanitize_html_fragment(analysis_json.get("sec7_content", "")),

        # Card 7: 交易可行性 & 规模适配
        "{{card7_trade_summary}}": _escape_text(analysis_json.get("sec8_summary", "")),
        "{{card7_revenue}}": _escape_text(
            f"{analysis_json.get('revenue_range', '待推算')} ({analysis_json.get('revenue_evidence', 'C')})"
        ),
        "{{card7_scale}}": _escape_text(analysis_json.get("scale_fit_label", "")),
        "{{card7_obstacle}}": _escape_text(analysis_json.get("transaction_obstacle_label", "")),
        "{{card7_trade_content}}": _sanitize_html_fragment(
            str(analysis_json.get("revenue_process_html", "") or "")
            + str(analysis_json.get("sec8_table", "") or "")
        ),

        # Card 8: 红旗核查
        "{{card8_flags_summary}}": _escape_text(analysis_json.get("sec9_summary", "")),
        "{{card8_flags_content}}": _sanitize_html_fragment(
            str(analysis_json.get("sec9_table", "") or "")
        ),

        # Card 9: 最终建议
        "{{card9_final_summary}}": _escape_text(analysis_json.get("sec10_summary", "")),
        "{{card9_final_content}}": _sanitize_html_fragment(
            _sanitize_html_fragment(analysis_json.get("next_steps_html", ""))
        ),

"{{参考来源清单}}": _sanitize_html_fragment(references_html),
    }

    # 雷达图数据
    scores = _safe_numeric_values(analysis_json.get("radar_scores"), 5, 3)
    replacements["{{5个维度评分，逗号分隔}}"] = ",".join(scores)

    # 证据等级条形图数据
    evidence_levels = _safe_numeric_values(analysis_json.get("evidence_levels"), 7, 2)
    replacements["{{8个结论项的证据等级数字，A=4 B=3 C=2 D=1}}"] = ",".join(evidence_levels)
    colors = []
    for e in evidence_levels:
        value = float(e)
        if value >= 4: colors.append("'#2E7D32'")
        elif value >= 3: colors.append("'#1565C0'")
        elif value >= 2: colors.append("'#F57F17'")
        else: colors.append("'#6B7280'")
    replacements["{{根据等级生成颜色数组}}"] = "[" + ",".join(colors) + "]"

    for old, new in replacements.items():
        template = template.replace(old, new)

    # 后处理：修复红旗区域嵌套grid和hr白块问题
    _flag_start = template.find('class="red-flags"')
    if _flag_start == -1:
        _flag_start = template.find("class='red-flags'")
    if _flag_start != -1:
        _flag_end = template.find('<!-- COLLAPSIBLE', _flag_start)
        if _flag_end != -1:
            _flag_area = template[_flag_start:_flag_end]
            _fixed = _flag_area.replace(
                '<hr>',
                '<hr style="grid-column:1/-1;border:none;border-top:1px solid #e0e0e0;margin:4px 0;">'
            )
            template = template[:_flag_start] + _fixed + template[_flag_end:]

    # 后处理：将AI输出的[1][2,3]等纯文本角标转换为可点击上标
    # 只在正文区域替换（sections到references之间），避免误改参考来源
    sec_start = template.find('id="sec1"')
    ref_start_marker = template.find('class="references"')
    if sec_start != -1 and ref_start_marker > sec_start:
        body = template[sec_start:ref_start_marker]
        body = re.sub(r'\[(\d+)\]', r'<sup class="cite" onclick="scrollToRef(\1)">\1</sup>', body)
        body = re.sub(r'\[(\d+)\s*,\s*(\d+)\]', r'<sup class="cite" onclick="scrollToRef(\1)">\1</sup><sup class="cite" onclick="scrollToRef(\2)">\2</sup>', body)
        template = template[:sec_start] + body + template[ref_start_marker:]

    # 后处理：修复参考来源中的空链接（href="#" → 真实URL）
    ref_start = template.find('class="references-list"')
    if ref_start == -1:
        ref_start = template.find("class='references-list'")
    if ref_start != -1:
        ref_end = template.find('</ol>', ref_start)
        if ref_end != -1:
            ref_end += 5
            ref_area = template[ref_start:ref_end]
            fixed_ref = _fix_references_html(ref_area, search_text)
            if fixed_ref != ref_area:
                template = template[:ref_start] + fixed_ref + template[ref_end:]

    # 清理未替换的占位符
    template = re.sub(r'\{\{.*?\}\}', '', template)

    return template


def _report_filename(company_name: str, task_id: str) -> str:
    safe_name = re.sub(r'[\\/:*?"<>|]', '_', company_name)
    return f"{safe_name}-PE深筛报告-{task_id}.html"


# ============ 核心分析流程 ============
# 两步分析策略：DeepSeek V4 是推理模型，思维链会消耗大量 token。
# 一次性输出完整 JSON 会导致思维链耗尽 max_tokens，content 为空。
# 因此拆分为两步：先分析输出 Markdown，再转为 JSON。

SYSTEM_PROMPT_ANALYSIS = """你是PE并购深筛分析师，专注于中小型非上市公司的公开信息深筛。
严格遵循反幻觉规则：不编造数据，未披露写"未见公开依据"，推断标注为[推断]。

基金偏好：recurring revenue、compounder+acquisition策略、食品饮料OEM(营收1-3亿净利率>10%)或商用清洗消毒(营收2000万+净利率>10%)。

一致性规则（必须严格遵守）：
- 同一家企业、同一类证据、没有新增高等级证据时，维度打分应保持稳定，不要因为措辞不同或搜索顺序不同而改变核心判断。
- 先定"稳定事实"，再基于稳定事实打分。业务定位优先级：公司官网/公司官方公众号自述 > 政府或正式披露 > 权威媒体 > 普通第三方 > 推断。股权、处罚、诉讼、财务与交易障碍优先级：上市公司公告/招股书、政府/司法/工商/监管正式记录 > 权威媒体 > 公司自述 > 普通第三方 > 推断。
- 商业模式只看收入形态、资产强度、壁垒类型；不要把交易障碍、规模、股权复杂度混进商业模式分数。
- 客户粘性只看锁定机制、合同/续约线索、切换成本、客户集中度；不要把规模大或品牌强误当作客户粘性高。
- Platform/Add-on 只看业务角色与协同，不看规模适配或交易障碍；交易障碍单独放到最后一维。
- 不能为了让结论"看起来更顺"而改写已经成立的事实口径。

按以下结构深入分析，其中5个定量维度与雷达图5轴一一对应，2个定性维度独立呈现。每个维度用子节区分子主题。

1.行业匹配度：子节一【公司身份核验】— 工商信息、成立时间、品牌、实控人。子节二【产品与服务结构】— 核心产品、收入模式、OEM vs自有品牌比例、复购特征。企查查/天眼查登记的经营范围不等于实际主业。
2.Recurring Revenue判断(核心)：重复性收入来源、持续供货/重复订单/连续排产等公开经营线索。对非上市中小企业，不把合同期限/续约率作为4分的必要条件。
3.市场地位：子节一【市场地位分析】— 行业排名、媒体冠名、KA覆盖、全国/区域领导地位。子节二【竞争位置与差异化】— 护城河、壁垒类型、可持续性。地位看"赢了多少"，竞争看"怎么赢"。
4.客户粘性：只看护城河——客户为什么不会轻易更换供应商？切换成本（技术认证/配方绑定/合规审核/渠道独占）、客户集中度。
5.商业模式：资产轻重、毛利水平、可复制性、壁垒类型。不把交易障碍或规模混入此项。
角色定位（定性，不评分）：评估目标公司更适合作为平台型收购标的还是补强型并购标的。输出 role_recommendation 字段，取值：platform（适合作为行业整合平台）、addon（适合作为补强型并购标的）、both（两者皆可，说明触发条件）、neither（皆不适合）。同时输出 role_recommendation_reason 简述依据。此判断不参与企业质量分计算，仅作为后续配对/筛选的标签。
交易可行性（定性，不评分）：评估规模以外的明确交易障碍。输出 transaction_feasibility 字段，取值：evidence_yes（有明确证据表明可推进，如已接触、有出售意向、股权清晰、估值预期合理）、inferred_yes（推断可推进，无明确证据但无显著障碍）、inferred_no（推断存在障碍，如股权过于集中且未表态、监管敏感行业）、evidence_no（有明确证据不可推进，如已被收购、明确拒绝、严重合规问题）。同时输出 transaction_feasibility_reason 简述依据。非上市中小企业不公开出售意愿属正常，不得据此将 feasibility 降为 evidence_no 或 inferred_no。此判断不参与任何评分，只在生成最终推进建议时作为否决/支持条件被引用。

8.红旗核查（逐项排查，必须按以下顺序和名称逐一标注🟢/🟡/🔴）：
1.一次性项目收入为主 2.必须持续开发新客户 3.纯贸易低壁垒 4.重资产/重CAPEX 5.原材料价格暴露
6.依赖单一大客户 7.过度依赖老板关系 8.技术替代风险高 9.监管合规风险高 10.行业边界模糊
11.平台能力不足 12.财务规范性/主体清晰度可疑 13.客户质量不足 14.公开信息极度稀缺
（注意：12 不公开财务数据属正常，只有主体不清/造假/处罚/冻结/失信才可标红；14 搜索结果少不等于公司有问题）
另输出执行摘要（必须在最前面输出）：150-300字，包含是否值得优先接触、Platform/Add-on判断、recurring revenue大致判断、规模推断（含区间和证据等级）、最关键吸引点、最关键风险、建议

7维度打分标准（必须严格按照以下rubric打分，不得凭感觉）：
1.行业匹配度：仅看业务类型，不看规模/利润/商业模式（那些属于其他维度）。
  5分（核心赛道）：冻品OEM、饮料OEM、调味料OEM、商用清洗消毒——基金明确在看的方向
  4分（相关赛道）：其他食品加工制造（烘焙、酱料馅料、预制菜、零食代工、食品配料等）
  3分（弱相关但非OEM）：食品供应链（配送、包装、原料贸易）、食品设备
  2分（不相关但有可迁移特征）：其他消费品OEM、化工、环保
  1分：完全不相关
  判定流程：先判断公司是否为食品制造/OEM，是则至少4分；再根据细分归属调整。凡公开资料明确为速冻/冷冻调理食品或冷冻半成品的B端OEM/ODM/代工业务，必须归入"冻品OEM"给5分。注意：鲜切水果/果切/鲜果加工/果蔬分拣/净菜加工属于"其他食品加工制造"给4分，不是冻品OEM，不得因经营范围内含"食品生产"而误判为5分。
2.Recurring Revenue：
  5分=存在公开可核验的强证据（如长约/订阅、重复收入占比、多年稳定复购数据或明确锁定机制）
  4分=存在较强公开经营证据证明业务会持续发生，不要求公开合同期限/续约率/收入占比/排他条款；例如多个具名B端客户+持续供货报道、提前排产/周期性订单、连续月产量或反复采购证据
  3分=仅能从产品特性或客户场景推断会复购，尚未看到已发生的持续供货、重复订单或连续经营事实
  2分=以项目制/一次性收入为主，但存在附加服务收入
  1分=纯一次性交易，无任何重复收入线索
  输出一致性：若正文结论写明"评分4分/应得4分"，radar_scores 中第2项必须同步输出4，不得出现正文与结构化分数矛盾。
3.客户粘性：
  5分=客户切换成本极高（技术认证、独家供应协议、共同研发）+客户集中度合理
  4分=切换成本中等（配方绑定、渠道绑定、质量认证），有明显粘性证据
  3分=推测有粘性（长期合作关系、定制化产品）但无直接证据
  2分=客户关系松散，产品标准化，更换供应商成本低
  1分=纯贸易/批发，客户无任何粘性
4.商业模式：
  5分=轻资产、高毛利、可规模化复制、技术/品牌壁垒明确
  4分=中等资产投入、毛利较好、有一定复制能力
  3分=资产较重或毛利偏低，但行业地位稳固
  2分=重资产或低毛利，增长依赖持续CAPEX投入
  1分=无壁垒的纯贸易/批发/倒卖模式

红旗→维度软联动（🔴触发时相关维度强制上限，必须遵守）：
- #1 一次性项目为主 🔴 → 维度2(Recurring Revenue) ≤2
- #3 纯贸易低壁垒 🔴 → 维度4(客户粘性)≤2 且 维度5(商业模式)≤2；但如果正文同时出现专利、研发、技术、专业服务、解决方案、系统、创新等强正面信号，必须改判为🟡待核实，不得直接写成"存在"
- #6 依赖单一大客户 🔴 → 维度4(客户粘性)≤2
- #7 过度依赖老板关系 🔴 → transaction_feasibility 不能为 evidence_yes（最多为 inferred_yes）
- #8 技术替代风险高 🔴 → role_recommendation 不能为 platform（最多 addon 或 neither）
- #10 行业边界模糊 🔴 → 维度1(行业匹配度)≤2
- #11 平台能力不足 🔴 → role_recommendation 不能为 platform（只能是 addon 或 neither）
- #12 主体清晰度可疑 🔴 → transaction_feasibility 不能为 evidence_yes（最多为 inferred_yes）
- #13 客户质量不足 🔴 → 维度4(客户粘性)≤2
#14 信息极度稀缺 🔴 不影响评分（信息可信度已不作为评分维度，仅作为红旗标注参考）
（上限为2=该维度基本被否决，上限为3=拉低但不致命）
例外：如果搜索结果中有具体证据证明"该红旗虽存在但不影响该维度"，可突破上限，但必须在分析中明确引用证据来源并标注[推断+反证]。

评级与结论规则（由程序复算，分析中不得自行放宽）：
- ABCD评级只衡量企业质量，不纳入角色定位或交易可行性。先应用红旗→维度软联动上限；企业质量分权重（Compounder 透镜）：行业匹配度 10%、Recurring Revenue 30%、客户粘性 15%、商业模式 20%、市场地位 25%。RR 和市场地位是核心（合占 55%），反映 Indutrade/Berkshire 对"复购可预测性 + 利基龙头地位"的优先排序。评级阈值：≥3.8→A，3.0-3.7→B，2.0-2.9→C，＜2.0→D。
- 角色定位（Platform/Add-on）与交易可行性为定性标签，独立呈现，不参与企业质量分计算。
- 雷达图展示5个原始企业质量维度，不直接代表企业质量分；企业质量分由程序复算。
- 推进建议结合ABCD评级、规模适配、角色定位与交易可行性综合判断：规模明确或大概率过大→不建议推进；规模明确或大概率过小、或未见规模不适配证据时，A→推进，B→观察，C/D→不建议推进。交易可行性 evidence_no（有明确证据不可推进）→强制不建议推进；inferred_no（推断存在障碍）→推进建议旁标注⚠警告但不改变结论；evidence_yes→追加"交易可行性已验证"标记。
- 监管合规风险高(#9)、主体清晰度可疑(#12)为🔴，或整体红旗严重度为严重，均覆盖为不建议推进；行业匹配度≤2同样不建议推进。
- 当企业质量结论本身的关键证据可信度为C/D时，程序将推进建议降为观察，提示补证；出售意愿未公开不得作为证据可信度不足的理由。

关键要求：
- 只使用通过公司相关性过滤且已编号的来源；不因行业、地区或同名泛词相近而将不明确属于目标公司的材料纳入事实判断
- 本系统只分析公司个体。不要展开泛行业市场规模、行业趋势或赛道普及性；行业匹配、市场地位、竞争位置必须基于目标公司的产品、客户、渠道、产能、区域覆盖、资质或公开经营事实。
- 按检索覆盖度行事：已见线索的主题可以分析；标记为仍待核实的主题必须写"未见充分公开依据"，不得用邻近信息补造结论
- 逐项给出结论的证据等级：A=明确直接证据，B=较强公开证据，C=合理推断，D=未见公开依据；证据不足不得写成已验证结论
- 证据可信度要简单表达：A=主要由直接/强公开证据支撑，B=大部分有较强公开证据但有少量推断，C=推断成分明显，D=公开证据太少；不要把它写成复杂统计术语
- 营收推算要有逻辑链条（如：员工数→人均产出→营收区间）
- 营收推断过程必须结构化展示到 revenue_process_html，至少 3 步，按"原始线索→中间假设→换算/校正→最终区间"的顺序写，便于人在报告里复核
- RR中小企业口径：媒体或公司公开资料能够共同证明多个具名客户、持续供货、提前排产、周期性订单或稳定连续产量时，可以给4分；不得因未披露合同期限、续约率、收入占比或排他协议而机械压到3分。只有合同/占比/多年连续数据等更强证据才给5分。
- 红旗项要逐一排查，没有发现的也要明确说明"经搜索未发现XX类红旗"
- 引用搜索结果时标注来源编号[1][2]等
- 信息源冲突处理：当工商登记/第三方平台说"销售配送"但公司官网或公司官方公众号自我介绍说"加工/OEM/定制化生产"时，可优先采用公司自述判断业务类型并标注差异。涉及股权、处罚、诉讼、财务与交易障碍时必须优先正式披露或政府/司法/工商/监管记录。
- 微信公众号是经营事实的补充与交叉验证来源，主要用于业务定位、产品、客户案例、产量或订单等经营线索；普通转载文章不得视为公司官方自述，也不得单独证明股权、财务、处罚、诉讼或可交易性结论。若股权/集团归属信息不足，应依靠工商、监管、上市公司或正式披露等基础来源继续定向搜索，不得用"未公开出售意愿"替代交易障碍证据。

输出原则：每个维度必须充分展开——结论、证据链、评分理由缺一不可；不确定处写"未见公开依据"而非推测。搜索工具只在确实缺失关键维度信息时使用，信息够用就停止搜索直接输出。

分析完成后，紧接在末尾输出结构化JSON。字段与维度的对应关系（严格遵守）：
- sec1 = 行业匹配度（身份核验+产品服务）
- sec2 = 商业模式（资产轻重、毛利、可复制性）
- sec3 = Recurring Revenue（重复收入判断）
- sec4 = 客户粘性（切换成本、客户集中度）
- sec5 = 市场地位分析（行业排名、KA覆盖）
- sec6 = 竞争位置与差异化（护城河、壁垒）
- sec7 = 角色定位（Platform/Add-on判断+理由）
- sec8 = 规模适配与交易可行性（规模推断+交易障碍）
- sec9 = 证据统计表
- sec10 = 最终建议
注意：维度顺序与上方分析顺序一致，但sec2=商业模式(非RR)，sec3=RR(非市场地位)。务必把每个维度的内容填入对应secN字段，不得串位。

---JSON---
{完整的JSON对象，包含 grade/total_score/industry/business_type/exec_summary/revenue_estimate/revenue_evidence/revenue_process_html/recurring_revenue/red_flags_count/red_flags_severity/red_flags/radar_scores/evidence_levels/sec1_summary/sec1_content/sec2_summary/sec2_content/sec3_summary/sec3_verdict/sec3_content/sec4_summary/sec4_content/sec5_summary/sec5_content/sec6_summary/sec6_content/sec7_summary/sec7_verdict/sec7_content/sec8_summary/revenue_range/revenue_method/sec8_table/maturity/role_recommendation/role_recommendation_reason/transaction_feasibility/transaction_feasibility_reason/has_transaction_obstacle/sec9_summary/sec9_table/sec10_summary/sec10_verdict/verdict_type/sec10_reason/key_question/approach/next_steps_html/references_html/red_flags_html 等全部字段。其中 role_recommendation 取值 platform/addon/both/neither；transaction_feasibility 取值 evidence_yes/inferred_yes/inferred_no/evidence_no；has_transaction_obstacle 为布尔值 true/false。}
---END_JSON---

注意：content 字段用HTML格式，角标用<sup class=\"cite\">N</sup>；red_flags 包含1-14全部项；references_html输出空字符串。"""

SYSTEM_PROMPT_JSON = """你是数据格式转换助手。将上面的深筛分析结果转为严格的JSON格式。

输出格式（不要输出其他内容）：
{"grade":"A/B/C/D初步判断","total_score":4.2,"industry":"行业","business_type":"业务类型","exec_summary":"150-300字执行摘要","revenue_estimate":"1.5-2亿","revenue_evidence":"A/B/C/D","revenue_process_html":"<ol><li>原始线索：...</li><li>中间假设：...</li><li>换算/校正：...</li><li>最终区间：...</li></ol>","recurring_revenue":"高/中/低","red_flags_count":3,"red_flags_severity":"严重/中等/轻微","red_flags":{"1":"🟢","2":"🟡","3":"🔴","4":"🟢","5":"🟢","6":"🟢","7":"🟢","8":"🟢","9":"🟢","10":"🟢","11":"🟢","12":"🟢","13":"🟢","14":"🟢"},"radar_scores":[4,3,3,3,4],"radar_scores_name_placeholder":null,"evidence_levels":[3,2,2,2,2,2,2],"sec1_summary":"行业匹配度摘要(身份核验+产品服务)","sec1_content":"行业匹配度HTML","sec2_summary":"商业模式摘要","sec2_content":"商业模式HTML","sec3_summary":"RecurringRevenue判断+理由","sec3_verdict":"高/中/低","sec3_content":"RecurringRevenue分析HTML","sec4_summary":"客户粘性摘要","sec4_content":"客户粘性HTML","sec5_summary":"市场地位摘要","sec5_content":"市场地位分析HTML","sec6_summary":"竞争位置摘要","sec6_content":"竞争位置HTML","sec7_summary":"Platform/Add-on判断+理由","sec7_verdict":"Platform/Add-on/Borderline/非目标","sec7_content":"Platform/Add-on分析HTML","sec8_summary":"营收区间+规模适配+交易障碍","revenue_range":"1-2亿","revenue_method":"推算方法","sec8_table":"<tr>...</tr>HTML行","maturity":"判断","role_recommendation":"platform","role_recommendation_reason":"理由","transaction_feasibility":"inferred_yes","transaction_feasibility_reason":"理由","transaction_feasibility_reason":"理由","sec9_summary":"证据统计","sec9_table":"<tr>...</tr>HTML行","sec10_summary":"初步判断","sec10_verdict":"优先深挖/值得继续看/保留观察/暂不优先初步建议","verdict_type":"pass/watch/reject初步建议","sec10_reason":"理由","key_question":"关键问题","approach":"切入点","next_steps_html":"HTML","references_html":"","red_flags_html":"<div>...</div>HTML","has_transaction_obstacle":false}

注意：
- content字段用HTML格式（<p><table><strong>等标签）
- 角标引用用<sup class="cite">N</sup>
- 红旗清单（red_flags_html）：每个红旗用 <div>🟢/🟡/🔴 序号.名称：结论。证据：[N]原文关键句</div> 表示。
- red_flags 字段必须和红旗清单逐项一致，不得省略1-14任一项。
- evidence_levels 依次对应行业归属、商业模式、Recurring Revenue、客户质量、切换成本、市场地位、规模推断，等级为4=明确直接证据、3=较强公开证据、2=合理推断、1=未见公开依据。
- 未见公开依据时，不得将对应维度写为中等水平；分数只能给1-2分。ABCD评级和推进建议会由程序根据维度分数、规模适配、证据和红旗重新计算。
  🟢示例：<div>🟢 1.一次性项目收入为主：未发现，长期代工模式收入稳定</div>
  🟡/🔴示例：<div>🔴 6.依赖单一大客户：存在。证据：[3]前五大客户占比超80%，最大客户占营收45%</div>
  不要用外层容器包裹，<div>标签直接拼接即可。必须包含全部14项PE红旗（按顺序），然后追加sec8补充核查项。
- 参考来源由程序根据实际检索结果生成；references_html 必须输出空字符串，不得自行新增、改写或补造来源条目。
  href先写#，服务器会自动替换为真实URL。N必须与搜索结果的[N]编号一致。
- 直接输出JSON，不要用代码块包裹"""


def _extract_json(text: str, task_id: str = "") -> dict:
    """从DeepSeek回复中鲁棒地提取JSON对象
    
    策略优先级：
    1. 提取 ```json ... ``` 代码块
    2. 直接解析整个文本
    3. 提取第一个 { 到最后一个 } 的子串
    4. 逐层剥离非法字符后重试
    5. 所有方法失败则抛出详细错误信息
    """
    import logging
    logger = logging.getLogger("pe-screening")
    
    original_text = text
    text = text.strip()
    
    # 策略1：提取 ```json ... ``` 代码块
    json_match = re.search(r'```json\s*([\s\S]*?)\s*```', text)
    if json_match:
        candidate = json_match.group(1).strip()
        try:
            return json.loads(candidate)
        except json.JSONDecodeError:
            pass  # 继续尝试其他策略
    
    # 策略2：直接解析
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass
    
    # 策略3：提取第一个 { 到最后一个 } 的子串
    start = text.find('{')
    end = text.rfind('}') + 1
    if start != -1 and end > start:
        candidate = text[start:end]
        try:
            return json.loads(candidate)
        except json.JSONDecodeError:
            pass
        
        # 策略4：修复常见JSON问题
        # 4a: 移除JS风格注释 // 和 /* */
        cleaned = re.sub(r'//.*?$', '', candidate, flags=re.MULTILINE)
        cleaned = re.sub(r'/\*[\s\S]*?\*/', '', cleaned)
        try:
            return json.loads(cleaned)
        except json.JSONDecodeError:
            pass
        
        # 4b: 修复尾部逗号 (trailing comma)
        cleaned2 = re.sub(r',\s*([}\]])', r'\1', candidate)
        try:
            return json.loads(cleaned2)
        except json.JSONDecodeError:
            pass
        
        # 4c: 修复单引号 → 双引号
        cleaned3 = candidate.replace("'", '"')
        try:
            return json.loads(cleaned3)
        except json.JSONDecodeError:
            pass
        
        # 4d: 修复HTML内容中的换行导致的非法JSON
        try:
            return _parse_json_lenient(candidate)
        except Exception:
            pass

        # 4e: 大语言模型经常在 JSON 字符串值内嵌入未转义的 ASCII 双引号
        # 如 "message": "他说"你好"" — 将字符串值内的 " 替换为中文引号“”
        try:
            return _fix_embedded_quotes(candidate)
        except Exception:
            pass

        # 4f: 移除所有控制字符（\x00-\x1f 除了 \n \r \t）
        cleaned5 = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', candidate)
        try:
            return json.loads(cleaned5)
        except json.JSONDecodeError:
            pass

        # 4g: 使用 json5 宽松解析（若可用）
        try:
            import json5  # type: ignore
            return json5.loads(candidate)
        except Exception:
            pass

    # 所有策略都失败，尝试最终兜底：逐字段正则提取核心字段
    try:
        salvage = _salvage_json_fields(original_text)
        if salvage:
            print(f"[JSON-ERROR] Task {task_id}: Salvaged {len(salvage)} fields via regex extraction")
            return salvage
    except Exception:
        pass
    
    # 所有策略都失败，输出详细调试信息
    print(f"[JSON-ERROR] Task {task_id}: All parsing strategies failed")
    print(f"[JSON-ERROR] Response length: {len(original_text)}")
    print(f"[JSON-ERROR] First 1000 chars: {original_text[:1000]}")
    print(f"[JSON-ERROR] Last 500 chars: {original_text[-500:]}")
    
    # 尝试定位JSON错误位置
    if start != -1 and end > start:
        candidate = text[start:end]
        try:
            json.loads(candidate)
        except json.JSONDecodeError as e:
            error_pos = e.pos if hasattr(e, 'pos') else -1
            context_start = max(0, error_pos - 50)
            context_end = min(len(candidate), error_pos + 50)
            print(f"[JSON-ERROR] Error at pos {error_pos}: ...{candidate[context_start:context_end]}...")
    
    raise ValueError(
        f"JSON解析失败（所有策略均已尝试）。"
        f"回复长度={len(original_text)}，"
        f"前300字：{original_text[:300]}"
    )


def _parse_json_lenient(text: str) -> dict:
    """宽松的JSON解析器，处理字符串值内的控制字符和换行"""
    # 简单策略：逐字符解析，在字符串值内转义控制字符
    result = []
    i = 0
    in_string = False
    escape_next = False
    
    while i < len(text):
        c = text[i]
        if escape_next:
            result.append(c)
            escape_next = False
            i += 1
            continue
        if c == '\\' and in_string:
            result.append(c)
            escape_next = True
            i += 1
            continue
        if c == '"':
            in_string = not in_string
            result.append(c)
            i += 1
            continue
        # 在字符串内，转义控制字符
        if in_string and ord(c) < 0x20:
            if c == '\n':
                result.append('\\n')
            elif c == '\r':
                result.append('\\r')
            elif c == '\t':
                result.append('\\t')
            else:
                result.append(f'\\u{ord(c):04x}')
            i += 1
            continue
        result.append(c)
        i += 1
    
    return json.loads(''.join(result))


def _fix_embedded_quotes(text: str) -> dict:
    """Handle AI-generated JSON where ASCII double quotes appear inside string values.

    Heuristic: after every key-value pair ``"key":``, scan the value string and
    replace bare ``"`` that are clearly not JSON delimiters with Chinese quotes.
    """
    # Mark JSON structural quotes vs content quotes
    # Walk through: if we see "key": then "value", the value delimiters are
    # structural; any " inside the value is content and should be replaced.
    result = []
    i = 0
    in_key_position = True  # true = expecting key or end of obj
    in_value = False
    depth = 0

    while i < len(text):
        c = text[i]
        if c == '"' and depth == 0 and not in_value:
            # Potentially opening a key or closing a value-quote
            result.append(c)
            if in_key_position:
                in_key_position = False
            else:
                # This is a value-closing quote — find the preceding colon+quote
                in_key_position = True
            i += 1
            continue
        if c == '"' and not in_key_position:
            # We're inside a string value — this " is content, not structure
            # Replace with Chinese left/right double quotes based on context
            prev_ch = result[-1] if result else ''
            next_ch = text[i + 1] if i + 1 < len(text) else ''
            if prev_ch and prev_ch not in (' ', ':', ',', '{', '['):
                result.append('”')  # right quote
            else:
                result.append('“')  # left quote
            i += 1
            continue
        if c in ('{', '['):
            result.append(c)
            depth += 1
            in_key_position = True
            i += 1
            continue
        if c in ('}', ']'):
            result.append(c)
            depth -= 1
            i += 1
            continue
        if c == ':' and depth == 0:
            in_key_position = False
            result.append(c)
            i += 1
            continue
        if c == ',' and depth == 0:
            in_key_position = True
            result.append(c)
            i += 1
            continue
        result.append(c)
        i += 1

    return json.loads(''.join(result))


def _salvage_json_fields(text: str) -> dict:
    """Ultimate fallback: regex-extract core fields when JSON is irreparably broken."""
    result: dict = {}
    patterns = {
        "grade": r'"grade"\s*:\s*"([ABCDEX])"',
        "total_score": r'"total_score"\s*:\s*([\d.]+)',
        "industry": r'"industry"\s*:\s*"([^"]+)"',
        "business_type": r'"business_type"\s*:\s*"([^"]*)"',
        "verdict_type": r'"verdict_type"\s*:\s*"(pass|watch|reject|PASS|WATCH|REJECT)"',
        "rr_label": r'"recurring_revenue"\s*:\s*"([^"]*)"',
        "summary": r'"exec_summary"\s*:\s*"([^"]*)"',
        "revenue_est": r'"revenue_estimate"\s*:\s*"([^"]*)"',
        "revenue_evidence": r'"revenue_evidence"\s*:\s*"([ABCD])"',
        "revenue_process_html": r'"revenue_process_html"\s*:\s*"(.*?)"\s*(?:,\s*"|}\s*$)',  # multi-line aware
        "red_flags_count": r'"red_flags_count"\s*:\s*(\d+)',
        "red_flags_severity": r'"red_flags_severity"\s*:\s*"([^"]*)"',
        "role_recommendation": r'"role_recommendation"\s*:\s*"([^"]*)"',
        "transaction_feasibility": r'"transaction_feasibility"\s*:\s*"([^"]*)"',
        "transaction_feasibility_reason": r'"transaction_feasibility_reason"\s*:\s*"([^"]*)"',
        "has_transaction_obstacle": r'"has_transaction_obstacle"\s*:\s*(true|false)',
        "scale_fit_label": r'"scale_fit_label"\s*:\s*"([^"]*)"',
        "transaction_obstacle_label": r'"transaction_obstacle_label"\s*:\s*"([^"]*)"',
        "evidence_confidence": r'"evidence_confidence"\s*:\s*"([ABCD])"',
        "key_question": r'"key_question"\s*:\s*"([^"]*)"',
        "approach": r'"approach"\s*:\s*"([^"]*)"',
        "advance_recommendation": r'"advance_recommendation"\s*:\s*"([^"]*)"',
        "advance_recommendation_reason": r'"advance_recommendation_reason"\s*:\s*"([^"]*)"',
    }
    for key, pat in patterns.items():
        m = re.search(pat, text, re.DOTALL)
        if m:
            val = m.group(1)
            try:
                result[key] = int(val)
            except ValueError:
                try:
                    result[key] = float(val)
                except ValueError:
                    result[key] = val

    # Section summary / content / verdict fields
    for n in range(1, 11):
        for field in ("summary", "content", "verdict"):
            key = f"sec{n}_{field}"
            if n < 10:
                # Content ends at next section or at comma-before-next-key
                if field == "content":
                    pat = rf'"{key}"\s*:\s*"(.*?)"\s*(?:,\s*"sec{n+1}_|,\s*"scales"|,\s*"maturity")'
                else:
                    pat = rf'"{key}"\s*:\s*"(.+?)"\s*[,}}]'
            else:
                pat = rf'"{key}"\s*:\s*"(.*?)"\s*[,}}]'
            m = re.search(pat, text, re.DOTALL if field == "content" else 0)
            if m:
                result[key] = m.group(1)

    # Tables
    for tbl in ("sec8_table", "sec9_table"):
        m = re.search(rf'"{tbl}"\s*:\s*"(.*?)"\s*[,}}]', text, re.DOTALL)
        if m:
            result[tbl] = m.group(1)

    # radar_scores
    scores_match = re.search(r'"radar_scores"\s*:\s*\[([^\]]*)\]', text)
    if scores_match:
        try:
            result["radar_scores"] = [int(x.strip()) for x in scores_match.group(1).split(",") if x.strip().isdigit()]
        except Exception:
            pass

    # evidence_levels
    ev_match = re.search(r'"evidence_levels"\s*:\s*\[([^\]]*)\]', text)
    if ev_match:
        try:
            result["evidence_levels"] = [int(x.strip()) for x in ev_match.group(1).split(",") if x.strip().isdigit()]
        except Exception:
            pass

    # red_flags dict
    flags = {}
    for i in range(1, 15):
        m_flag = re.search(rf'"{i}"\s*:\s*"([^"]*)"', text)
        if m_flag:
            flags[str(i)] = m_flag.group(1)
    if flags:
        result["red_flags"] = flags

    result.setdefault("grade", "B")
    result.setdefault("total_score", 2.5)
    result.setdefault("industry", "待确认")
    result.setdefault("business_type", "")
    result.setdefault("verdict_type", "watch")
    result.setdefault("recurring_revenue", "待确认")
    result.setdefault("summary", "（JSON解析失败，从原始文本中提取）")
    result.setdefault("red_flags_count", 0)
    result.setdefault("red_flags_severity", "轻微")
    result.setdefault("radar_scores", [2, 2, 2, 2, 2])
    result.setdefault("evidence_levels", [2, 2, 2, 2, 2, 2, 2])
    result.setdefault("role_recommendation", "")
    result.setdefault("transaction_feasibility", "")
    result.setdefault("has_transaction_obstacle", False)
    result.setdefault("revenue_evidence", "C")
    result.setdefault("scale_fit_label", "")
    result.setdefault("transaction_obstacle_label", "")
    result.setdefault("evidence_confidence", "C")
    result.setdefault("revenue_process_html", "")
    return result


_SCREENING_TIMEOUT = int(os.environ.get("SCREENING_TIMEOUT", "720"))
_KIMI_SCREENING_TIMEOUT = int(os.environ.get("KIMI_SCREENING_TIMEOUT", "1200"))
_BATCH_ANALYSIS_TIMEOUT = int(os.environ.get("BATCH_ANALYSIS_TIMEOUT", "600"))


async def _run_screening_core(task_id, company_name, website, industry, bp_text):
    """Internal: all screening work, returns nothing, updates task dict."""
    task = tasks[task_id]
    try:
        # 第1阶段：搜索
        task["status"] = "searching"
        task["progress"] = "正在搜索公开信息..."

        # 第一层：按固定尽调主题搜索，不随公司或模型输出临时改变框架。
        search_plan = _build_screening_search_plan(
            company_name, "", industry, website
        )
        search_topics = {item["query"]: item["topic"] for item in search_plan}
        raw_search_results = await multi_search(
            [item["query"] for item in search_plan],
            company_name=company_name,
            comprehensive=True,
        )
        seen_urls: set[str] = set()
        search_results = _govern_search_result_groups(
            raw_search_results, company_name, "", seen_urls
        )

        # 百度是中文基础来源层；微信不是基础检索，只在两轮基础检索后补经营事实缺口。
        task["progress"] = "正在百度搜索中文内容..."
        baidu_plan = _build_baidu_search_plan(company_name)
        baidu_topics = {item["query"]: item["topic"] for item in baidu_plan}
        raw_baidu_results = {
            item["query"]: await search_baidu(item["query"], max_results=6)
            for item in baidu_plan
        }
        baidu_results = _govern_search_result_groups(
            raw_baidu_results, company_name, "", seen_urls
        )

        # 第二层：只为第一轮尚未形成公司相关证据的主题执行针对性补搜。
        base_coverage_text = _coverage_text_from_results(search_results, baidu_results)
        coverage = _assess_search_coverage(base_coverage_text)
        # QCC MCP covers identity/ownership — don't gap-search these
        if QCC_MCP_KEY:
            coverage["ownership"] = True  # QCC covers 股东/实控人, skip gap search
        gap_plan = _build_gap_search_plan(company_name, coverage)
        if gap_plan:
            task["progress"] = "正在针对公开信息缺口补搜..."
            gap_topics = {item["query"]: item["topic"] for item in gap_plan}
            gap_search = await multi_search(
                [item["query"] for item in gap_plan],
                company_name=company_name,
                comprehensive=True,
            )
            gap_search = _govern_search_result_groups(
                gap_search, company_name, "", seen_urls
            )
            search_results.update(gap_search)
            search_topics.update(gap_topics)
            raw_baidu_gap = {
                item["query"]: await search_baidu(item["query"], max_results=4)
                for item in gap_plan
            }
            baidu_gap = _govern_search_result_groups(
                raw_baidu_gap, company_name, "", seen_urls
            )
            baidu_results.update(baidu_gap)
            baidu_topics.update(gap_topics)
            base_coverage_text = _coverage_text_from_results(search_results, baidu_results)
            coverage = _assess_search_coverage(base_coverage_text)
            if QCC_MCP_KEY:
                coverage["identity"] = True
                coverage["ownership"] = True

        wechat_queries = _wechat_supplement_queries(
            company_name, "", industry, base_coverage_text
        )
        wechat_results = []
        if wechat_queries:
            task["progress"] = "正在补充微信公众号经营事实与交叉验证..."
            for query in wechat_queries:
                wechat_results += await search_wechat_articles(query, max_results=5)
            candidates = _dedupe_wechat_results(wechat_results, limit=_WECHAT_MAX_ARTICLES * 2)
            verified_wechat = _filter_relevant_wechat_results(
                candidates, company_name, ""
            )
            wechat_results = []
            for article in verified_wechat:
                key = str(article.get("url", "")).strip().rstrip("/").lower()
                if key and key in seen_urls:
                    continue
                if key:
                    seen_urls.add(key)
                wechat_results.append(article)
            print(f"[WeChat] Kept {len(wechat_results)}/{len(candidates)} company-specific operating supplements")

        ref_counter = 1
        verified_refs_html = ""
        ref_provider_counts: dict[str, int] = {}
        qcc_info = ""
        if QCC_MCP_KEY:
            task["progress"] = "正在查询工商信息..."
            qcc_info, qcc_refs = await collect_qcc_evidence(company_name)
            for ref in qcc_refs:
                provider = ref.get("provider", "企查查MCP")
                if not _can_add_report_reference(provider, ref_counter, ref_provider_counts):
                    continue
                verified_refs_html += _reference_item_html(
                    ref_counter, provider, ref.get("title", ""), ref.get("url", "")
                )
                _mark_report_reference(provider, ref_provider_counts)
                ref_counter += 1

        # 基础来源先编号并提供给模型，避免补充来源主导结论。
        found_topics = [
            _SEARCH_TOPIC_LABELS[topic] for topic, found in coverage.items() if found
        ]
        missing_topics = [
            _SEARCH_TOPIC_LABELS[topic] for topic, found in coverage.items() if not found
        ]
        search_text = (
            "## 检索方法与覆盖度\n"
            "基础检索按六类尽调主题执行：主体与主营业务、产品/客户与持续经营、"
            "营收/产能与规模、股权与交易障碍、诉讼/处罚与合规、行业位置与协同价值。"
            "基础检索不足的主题已进行针对性补搜；微信用于补充及交叉验证经营事实。\n"
            f"已见公开线索的主题：{'、'.join(found_topics) if found_topics else '无'}。\n"
            f"仍待核实的主题：{'、'.join(missing_topics) if missing_topics else '无明显缺口'}。\n"
        )
        for query, results in search_results.items():
            providers = {r.get("provider", "") for r in results if r.get("type") == "result"}
            provider_label = "/".join(sorted(p for p in providers if p)) or "网页"
            topic_label = _SEARCH_TOPIC_LABELS.get(search_topics.get(query, ""), "补充检索")
            search_text += f"\n### [{topic_label}] {provider_label}搜索：{query}\n"
            for r in results:
                if r["type"] == "answer":
                    search_text += f"摘要：{r['content'][:500]}\n"
                elif r["type"] == "result":
                    provider = r.get("provider", "网页搜索")
                    if not _can_add_report_reference(provider, ref_counter, ref_provider_counts):
                        continue
                    search_text += f"[{ref_counter}] {r['title']}\nURL: {r['url']}\n{r['content'][:800]}\n\n"
                    verified_refs_html += _reference_item_html(
                        ref_counter, provider, r.get("title", ""), r.get("url", "")
                    )
                    _mark_report_reference(provider, ref_provider_counts)
                    ref_counter += 1
                elif r["type"] == "error":
                    search_text += f"搜索出错：{r['content']}\n"

        search_text += "\n\n## 百度搜索结果\n"
        for query, results in baidu_results.items():
            topic_label = _SEARCH_TOPIC_LABELS.get(baidu_topics.get(query, ""), "补充检索")
            search_text += f"\n### [{topic_label}] 百度搜索：{query}\n"
            for r in results:
                if not _can_add_report_reference("百度", ref_counter, ref_provider_counts):
                    continue
                search_text += f"[{ref_counter}] {r['title']}\nURL: {r['url']}\n{r['content'][:500]}\n\n"
                verified_refs_html += _reference_item_html(
                    ref_counter, "百度", r.get("title", ""), r.get("url", "")
                )
                _mark_report_reference("百度", ref_provider_counts)
                ref_counter += 1

        if wechat_results:
            search_text += (
                f"\n\n## 微信公众号补充来源（经营事实补充与交叉验证，去重后{len(wechat_results)}篇；"
                "仅辅助业务定位、产品、客户和持续经营线索，不单独证明财务、股权、处罚或交易障碍）\n"
            )
            for index, art in enumerate(wechat_results):
                if not _can_add_report_reference("微信补充", ref_counter, ref_provider_counts):
                    continue
                content_text = ""
                if index < _WECHAT_MAX_FULLTEXT and art.get("url"):
                    detail = await fetch_wechat_article_content(art["url"])
                    content_text = detail.get("content", "")[:2000] if detail else ""
                search_text += (
                    f"[{ref_counter}] {art.get('title','')}\nURL: {art.get('url','')}\n"
                    f"公众号: {art.get('account','')} | {art.get('time','')}\n"
                    f"{'正文: ' + content_text if content_text else '摘要: ' + art.get('digest','')}\n\n"
                )
                verified_refs_html += _reference_item_html(
                    ref_counter,
                    "微信补充",
                    f'{art.get("title", "")} - {art.get("account", "")}',
                    art.get("url", ""),
                )
                _mark_report_reference("微信补充", ref_provider_counts)
                ref_counter += 1

        task["progress"] = "正在生成并执行模型补搜计划..."
        supplement_text, supplement_refs_html, ref_counter = await _run_model_planned_supplemental_searches(
            company_name=company_name,
            current_search_text=search_text,
            ref_counter=ref_counter,
            seen_urls=seen_urls,
            provider_counts=ref_provider_counts,
            model=task.get("model", ""),
        )
        if supplement_text:
            search_text += supplement_text
        if supplement_refs_html:
            verified_refs_html += supplement_refs_html

        # 截断总搜索文本
        if len(search_text) > 25000:
            search_text = search_text[:35000] + "\n\n[搜索结果已截断，共" + str(ref_counter-1) + "条参考]"

        # 第1.5阶段：Playwright 抓取关键页面全文
        task["status"] = "crawling"
        task["progress"] = "正在抓取关键网页全文..."

        # 收集搜索结果中的关键URL，并做初步过滤（去重+域名限制）
        crawl_urls = []
        crawl_domain_counts = {}
        company_crawl_domain = ""
        if website:
            try:
                cw = website if "://" in website else f"https://{website}"
                company_crawl_domain = urlparse(cw).netloc.lower().replace("www.", "")
            except:
                pass

        for query, results in search_results.items():
            for r in results:
                if r["type"] == "result" and r.get("url"):
                    url = r["url"]
                    if not _is_url_worth_crawling(url, website):
                        continue
                    # 域名去重（同域名最多2个，官网不限）
                    try:
                        d = urlparse(url).netloc.lower().replace("www.", "")
                    except:
                        d = url
                    cnt = crawl_domain_counts.get(d, 0)
                    if cnt >= 2 and d != company_crawl_domain:
                        continue
                    crawl_domain_counts[d] = cnt + 1
                    crawl_urls.append(url)

        # 限制总数
        crawl_urls = crawl_urls[:10]

        crawled_text = ""
        crawled_count = 0

        # 1) 先抓公司官网（如果提供了）
        if website:
            task["progress"] = f"正在抓取公司官网 {website}..."
            website_results = await crawl_company_website(website, max_pages=5)
            for url, info in website_results.items():
                if info.get("content") and info["length"] > 100:
                    content_preview = info["content"][:3000]  # 每页最多3000字
                    crawled_text += f"\n### [官网抓取] {info.get('title', url)}\nURL: {url}\n{content_preview}\n\n"
                    crawled_count += 1

        # 2) 再抓搜索结果中的高价值URL
        if crawl_urls:
            task["progress"] = f"正在抓取搜索结果中的 {len(crawl_urls)} 个关键页面..."
            try:
                crawl_results = await asyncio.wait_for(
                    crawl_urls_with_playwright(crawl_urls, company_website=website, max_concurrent=3),
                    timeout=60,
                )
            except asyncio.TimeoutError:
                print("[Playwright] Crawling timed out after 60s, continuing with partial results")
                crawl_results = {}
            for url, info in crawl_results.items():
                if info.get("content") and info["length"] > 200:
                    # 搜索结果页面的内容可能较长，截取重点部分
                    content_preview = info["content"][:2000]  # 搜索结果页最多2000字
                    crawled_text += f"\n### [深度抓取] {info.get('title', url)}\nURL: {url}\n{content_preview}\n\n"
                    crawled_count += 1

        print(f"[Playwright] Total crawled pages: {crawled_count}, total text length: {len(crawled_text)}")

        # 如果抓取到了内容，追加到搜索文本中
        if crawled_text:
            # 重新计算总文本长度，可能需要调整截断
            full_text = search_text + "\n\n## 深度抓取内容（Playwright全文提取）\n" + crawled_text
            # 总上限（含网页搜索+微信文章+Playwright抓取）
            if len(full_text) > 45000:
                full_text = full_text[:45000] + f"\n\n[总内容已截断，含搜索{ref_counter-1}条+抓取{crawled_count}页]"
            search_text = full_text

        # 第1.8阶段：企查查工商信息已在上方查询，直接复用
        # qcc_info already populated from line ~3506, no second call needed

        # 第2阶段：DeepSeek 分析
        task["status"] = "analyzing"
        task["progress"] = "正在分析公开信息..."

        user_message = f"""请对以下公司进行PE并购深筛分析：

公司名称：{company_name}
{'官网：' + website if website else ''}
{'行业：' + industry if industry else ''}

{qcc_info}## 搜索结果

{search_text}

{'## BP（商业计划书）内容' + chr(10) + bp_text if bp_text else ''}

请完成5个维度+角色定位+交易可行性的深筛分析。"""

        # 第2步a：分析（输出 Markdown）
        messages_analysis = [
            {"role": "system", "content": SYSTEM_PROMPT_ANALYSIS},
            {"role": "user", "content": user_message},
        ]

        # 调用 DeepSeek - 第1步：分析（AI可自主搜索补充信息，结果纳入引用编号）
        task["status"] = "analyzing"
        task["progress"] = "正在分析公开信息（AI可自主搜索补充）..."

        pre_search_count = ref_counter  # next available ref number
        use_tools = not task.get("model", "").startswith("kimi-")  # Kimi too slow with tools
        task["live_text"] = ""
        def _on_stream_chunk(chunk: str):
            task["live_text"] = (task.get("live_text", "") + chunk)[-8000:]

        analysis_md, ai_refs_html, _ = await call_deepseek_with_tools(
            messages_analysis,
            tools=ANALYSIS_TOOLS if use_tools else None,
            max_rounds=3,
            ref_counter_start=pre_search_count,
            company_name=company_name,
            model=task.get("model", ""),
            progress_callback=lambda msg: task.update({"progress": msg}),
            stream_callback=_on_stream_chunk if not task.get("model", "").startswith("kimi-") else None,
        )
        task["live_text"] = ""  # done streaming
        if ai_refs_html:
            verified_refs_html += ai_refs_html

        if not analysis_md.strip() or analysis_md.startswith("[思维链"):
            raise ValueError(f"AI分析失败：{analysis_md[:200] if analysis_md else '返回为空'}")

        task["status"] = "generating"
        task["progress"] = "正在生成结构化报告..."

        # 第2步b：提取 JSON — 优先从分析响应中直接提取，省掉二次调用
        analysis_json = None
        json_match = re.search(r"---JSON---\s*(.*?)\s*---END_JSON---", analysis_md, re.DOTALL)
        if json_match:
            try:
                analysis_json = _extract_json(json_match.group(1), task_id)
                print("[Screening] JSON extracted directly from analysis response")
            except ValueError as e:
                print(f"[Screening] Direct JSON extraction failed: {str(e)[:200]}, falling back to two-step")

        # Fallback: traditional two-step JSON conversion
        if analysis_json is None:
            messages_json = [
                {"role": "system", "content": SYSTEM_PROMPT_JSON},
                {"role": "user", "content": f"将以下深筛分析转为JSON格式：\n\n{analysis_md}"},
            ]
            json_parse_error = None
            for json_attempt in range(3):
                full_response = await call_deepseek(messages_json, retry_count=1, max_tokens=8192, model=task.get("model", ""))
                if not full_response.strip():
                    if json_attempt < 2:
                        continue
                    raise ValueError("JSON转换步骤返回为空")
                try:
                    analysis_json = _extract_json(full_response, task_id)
                    json_parse_error = None
                    break
                except ValueError as e:
                    json_parse_error = str(e)
                    if json_attempt < 2:
                        messages_json = [
                            {"role": "system", "content": SYSTEM_PROMPT_JSON},
                            {"role": "user", "content": f"将以下深筛分析转为JSON格式：\n\n{analysis_md}"},
                            {"role": "assistant", "content": full_response[:3000]},
                            {"role": "user", "content": f"上一次输出的JSON解析失败，错误：{str(e)[:300]}\n请重新输出完整的正确JSON。"},
                        ]
                        continue
            if json_parse_error and analysis_json is None:
                raise ValueError(json_parse_error)

        # 第3阶段：程序按统一规则重算分数与结论，再生成报告
        task["status"] = "rendering"
        task["progress"] = "正在渲染报告..."

        analysis_json = _apply_screening_score(analysis_json)

        # ---- 自我审查：检查内部一致性 ----
        _self_critique(analysis_json, company_name)

        report_html = fill_report_template(company_name, analysis_json, search_text, verified_refs_html)

        # 保存报告
        report_filename = _report_filename(company_name, task_id)
        report_path = REPORT_DIR / report_filename
        report_path.write_text(report_html, encoding="utf-8")

        task["status"] = "completed"
        task["progress"] = "报告生成完成"
        task["report_url"] = f"/report/{task_id}"
        task["report_filename"] = report_filename
        task["analysis"] = analysis_json
        _save_task_history(task_id, task)

    except Exception as e:
        import traceback
        error_detail = traceback.format_exc()
        print(f"[ERROR] Task {task_id} failed:\n{error_detail}")
        task["status"] = "error"
        task["progress"] = f"错误：{str(e)}"
        task["error_detail"] = error_detail


async def run_screening(task_id: str, company_name: str, website: str, industry: str, bp_text: str, model: str = ""):
    """执行完整的深筛流程，带整体超时保护。model 参数可选，默认 DEEPSEEK_MODEL。"""
    effective_model = model or DEEPSEEK_MODEL
    provider = _get_model_config(effective_model)["provider"]
    timeout_sec = _KIMI_SCREENING_TIMEOUT if provider == "kimi" else _SCREENING_TIMEOUT
    task = tasks.get(task_id, {})
    task["model"] = effective_model
    try:
        await asyncio.wait_for(
            _run_screening_core(task_id, company_name, website, industry, bp_text),
            timeout=timeout_sec,
        )
    except asyncio.TimeoutError:
        task = tasks.get(task_id, {})
        task["status"] = "error"
        task["progress"] = f"任务超时（模型：{effective_model}，{timeout_sec}秒），请重试"






# ============ 批量完整分析 ============
batch_tasks: dict = {}

_BATCH_MAX_COMPANIES = 2000
_BATCH_COMPANY_HEADERS = {"企业名称", "公司名称", "企业名", "名称"}
_BATCH_IMPORT_ALIASES = {
    "registered_capital": ("注册资本", "注册资金"),
    "legal_representative": ("法定代表人", "法人"),
    "established_at": ("成立日期", "成立时间", "注册日期"),
    "registration_status": ("经营状态", "登记状态", "企业状态"),
    "credit_code": ("统一社会信用代码", "统一信用代码"),
    "province": ("省份", "所属地区", "地区", "region", "province"),
    "enterprise_type": ("企业类型", "公司类型", "类型"),
    "insured_count": ("参保人数", "社保人数", "员工人数"),
    "paid_in_capital": ("实缴资本", "实缴资金", "实缴出资"),
}

# ---- 批量预筛选常量 ----
_TARGET_REGIONS = frozenset({"上海", "江苏", "浙江", "安徽", "福建", "湖南", "湖北", "广东", "四川", "重庆"})
_ABNORMAL_STATUS_KW = frozenset({"注销", "吊销", "清算", "停业", "经营异常", "严重违法"})
_R12_ENTERPRISE_TYPE_KW = frozenset({"外商", "港澳台", "中外合资", "外资", "外国法人"})
_R5_ENTERPRISE_TYPE_KW = frozenset({"国有独资", "集体所有制", "全民所有制"})
_R3_ENTERPRISE_TYPE_KW = frozenset({"国有", "集体", "集团", "央企"})


def _apply_table_prescreen_rules(record: dict, *, exclude_foreign: bool = False) -> dict | None:
    """Apply table-based exclusion rules using imported Qichacha fields. Returns
    ``{"exclude_code": "R1", "exclude_reason": "..."}`` when a rule triggers,
    or ``None`` when the company passes all checks.  Zero API calls — pure data.

    ``exclude_foreign`` toggles R12 (外资控股排除), defaults to False."""
    company_name = record.get("company_name", "")

    # R8 经营状态异常
    status = record.get("registration_status", "")
    if status:
        status_lower = status.replace("（", "(").replace("）", ")")
        for kw in _ABNORMAL_STATUS_KW:
            if kw in status_lower:
                return {"exclude_code": "R8", "exclude_reason": f"经营状态异常（{status}）"}

    # R1 地域不符合
    province_raw = record.get("province", "")
    province = province_raw.strip().rstrip("省").rstrip("市")
    if province:
        target = {r.rstrip("省").rstrip("市") for r in _TARGET_REGIONS}
        if province not in target:
            display_parts = sorted(_TARGET_REGIONS, key=lambda x: x)
            return {"exclude_code": "R1",
                    "exclude_reason": f"注册地（{province_raw}）不在目标地域（仅保留{'/'.join(display_parts)}）"}

    # R7 成立时间过短
    established = record.get("established_at", "")
    if established:
        try:
            date_str = re.sub(r"[年月]", "-", str(established)).replace("日", "").replace("/", "-").strip()
            m = re.match(r"(\d{4})-(\d{1,2})(?:-(\d{1,2}))?", date_str)
            if m:
                y, mo, d = int(m.group(1)), int(m.group(2)), int(m.group(3) or 1)
                est = datetime(y, mo, d, tzinfo=timezone.utc)
                years = (datetime.now(timezone.utc) - est).days / 365.25
                if years < 3:
                    return {"exclude_code": "R7",
                            "exclude_reason": f"成立时间（{established}）距今仅{years:.1f}年，不足3年"}
        except (ValueError, TypeError):
            pass

    # R9 参保人数过低
    insured = record.get("insured_count", "")
    if insured:
        try:
            count = int(re.sub(r"[^\d]", "", str(insured)))
            if count <= 5:
                return {"exclude_code": "R9", "exclude_reason": f"参保人数（{insured}）≤5人"}
        except (ValueError, TypeError):
            pass

    # R11 注册资本 < 50万 且实缴为 0
    reg_cap = record.get("registered_capital", "")
    paid_in = record.get("paid_in_capital", "")
    if reg_cap:
        try:
            cap_val = float(re.sub(r"[^\d.]", "", str(reg_cap)))
            if cap_val < 50:
                paid_val = 0.0
                if paid_in:
                    try:
                        paid_val = float(re.sub(r"[^\d.]", "", str(paid_in)))
                    except (ValueError, TypeError):
                        pass
                if paid_val == 0:
                    return {"exclude_code": "R11",
                            "exclude_reason": f"注册资本{reg_cap}万<50万且实缴为0"}
        except (ValueError, TypeError):
            pass

    # R3 / R5 / R12 企业类型（R12 → R5 → R3 顺序，避免交叉命中）
    ent_type = record.get("enterprise_type", "")
    if ent_type:
        # R12 外资（含港澳台）控股 — 仅当用户主动开启时生效
        if exclude_foreign and any(kw in ent_type for kw in _R12_ENTERPRISE_TYPE_KW):
            return {"exclude_code": "R12",
                    "exclude_reason": f"企业类型含外资/港澳台成分（{ent_type}）"}
        # R5 国资/集体控股
        if any(kw in ent_type for kw in _R5_ENTERPRISE_TYPE_KW):
            return {"exclude_code": "R5",
                    "exclude_reason": f"企业类型为国资/集体控股（{ent_type}）"}
        # R3 集团/国企/央企体系
        if any(kw in ent_type for kw in _R3_ENTERPRISE_TYPE_KW):
            return {"exclude_code": "R3",
                    "exclude_reason": f"企业类型属于集团/国有体系（{ent_type}）"}

    return None  # 通过全部表格规则


def _cell_text(value) -> str:
    if value is None:
        return ""
    return str(value).strip()


def _parse_batch_excel(content: bytes) -> list[dict]:
    """Read a Qichacha-style .xlsx export while retaining imported registry fields."""
    try:
        import openpyxl
    except ImportError as exc:
        raise HTTPException(500, "服务器未安装 Excel 解析组件 openpyxl") from exc
    try:
        workbook = openpyxl.load_workbook(BytesIO(content), read_only=True, data_only=True)
    except Exception as exc:
        raise HTTPException(400, "无法读取 Excel，请上传 .xlsx 格式的企查查导出文件") from exc
    sheet = workbook.active
    values = list(sheet.iter_rows(values_only=True))
    header_index = -1
    headers: list[str] = []
    for index, row in enumerate(values[:20]):
        candidate = [_cell_text(cell) for cell in row]
        if any(header in _BATCH_COMPANY_HEADERS for header in candidate):
            header_index = index
            headers = candidate
            break
    if header_index < 0:
        raise HTTPException(400, "未识别到企业名称列，请确认文件包含“企业名称”或“公司名称”")
    name_col = next(index for index, header in enumerate(headers) if header in _BATCH_COMPANY_HEADERS)
    records = []
    seen_names = set()
    for row in values[header_index + 1:]:
        cells = [_cell_text(cell) for cell in row]
        company_name = cells[name_col] if name_col < len(cells) else ""
        if not company_name or company_name in seen_names:
            continue
        seen_names.add(company_name)
        imported = {
            header: cells[index]
            for index, header in enumerate(headers)
            if header and index < len(cells) and cells[index]
        }
        record = {"company_name": company_name, "imported_fields": imported}
        for key, aliases in _BATCH_IMPORT_ALIASES.items():
            record[key] = next((imported[alias] for alias in aliases if imported.get(alias)), "")
        records.append(record)
    if not records:
        raise HTTPException(400, "Excel 中没有可分析的企业记录")
    if len(records) > _BATCH_MAX_COMPANIES:
        raise HTTPException(400, f"单次最多 {_BATCH_MAX_COMPANIES} 家，请拆分文件后上传")
    return records


def _manual_batch_records(company_names: list[str]) -> list[dict]:
    records = []
    seen = set()
    for value in company_names:
        name = str(value or "").strip()
        if name and name not in seen:
            records.append({"company_name": name, "imported_fields": {}})
            seen.add(name)
    if len(records) > _BATCH_MAX_COMPANIES:
        raise HTTPException(400, f"批量完整分析单次最多 {_BATCH_MAX_COMPANIES} 家")
    return records


def _plain_fragment(value: str) -> str:
    return BeautifulSoup(str(value or ""), "html.parser").get_text(" ", strip=True)

SYSTEM_PROMPT_QUICK_SCORE = """你是PE并购深筛分析师。对以下公司进行快速评分。

5维度打分标准（1-5分）：
1.行业匹配度：冻品/饮料/调味料OEM或商用清洗消毒=5分，其他食品加工=4分（含鲜切水果/果切/净菜加工/果蔬分拣），食品供应链非OEM=3分，不相关=1-2分；速冻/冷冻调理食品或冷冻半成品的B端OEM/ODM必须按冻品OEM给5分
2.Recurring Revenue：长约/订阅/重复收入占比或多年稳定复购等强证据=5分；确认为盒马/奥乐齐/叮咚/美团/山姆/开市客等头部KA零售商正式供应商且持续供货=5分；多个具名客户且有持续供货、提前排产、周期性订单或连续产量等经营证据=4分（中小企业无需公开合同期限/续约率/排他条款）；仅由产品特性推测复购=3分；项目制为主=2分；纯一次性=1分
3.市场地位（隐形冠军识别）：全国性细分龙头/权威确认为第一/最大/龙头=5分；区域性龙头/省级第一/头部KA全覆盖=4分；有知名度但非领导地位=3分；跟随者=2分；新进入者/无地位=1分
4.客户粘性：切换成本极高+客户分散=5分，有粘性证据=4分，推测有粘性=3分，标准产品=2分，纯贸易=1分
5.商业模式：轻资产/高毛利/可复制=5分，中等=3分，重资产/低毛利=2分，无壁垒纯贸易=1分
角色定位（定性）：输出 role_recommendation 取值 platform/addon/both/neither + role_recommendation_reason
交易可行性（定性）：输出 transaction_feasibility 取值 evidence_yes/inferred_yes/inferred_no/evidence_no + transaction_feasibility_reason。非上市中小企业未公开出售意愿属正常，不得据此降为 evidence_no 或 inferred_no

证据原则：
- 每个评分维度同时输出 evidence_levels 的对应等级（依次：行业归属、商业模式、Recurring Revenue、客户质量、切换成本、市场地位、规模推断，共7项）：4=直接明确证据，3=较强公开证据，2=合理推断，1=未见公开依据。
- 不能因为没查到负面就给3分。某项关键事实未见公开依据时，该项只能给1-2分，并在 analysis 中说明待核实内容。
- 你提供评分观察和红旗判断；最终总分、评级、结论由程序重新计算。
- 对非上市中小企业，未公开合同期限、续约率、重复收入占比或排他条款不构成RR降分理由；若已有多个具名客户和持续经营行为的公开证据，RR可给4分。

红旗→维度软联动（🔴触发时强制上限）：
#1一次性项目为主🔴→维度2≤2 | #3纯贸易低壁垒🔴→维度3≤2,维度4≤2 | #6依赖大客户🔴→维度3≤2
#7依赖老板🔴→transaction_feasibility不能为evidence_yes | #8技术替代🔴→role_recommendation不能为platform | #10行业边界模糊🔴→维度1≤2
#11平台不足🔴→role_recommendation不能为platform | #12主体可疑🔴→transaction_feasibility不能为evidence_yes
#13客户质量不足🔴→维度3≤2 |

14项PE红旗（逐项标注🟢/🟡/🔴）：
1.一次性项目收入为主 2.必须持续开发新客户 3.纯贸易低壁垒 4.重资产/重CAPEX 5.原材料价格暴露
6.依赖单一大客户 7.过度依赖老板关系 8.技术替代风险高 9.监管合规风险高 10.行业边界模糊
11.平台能力不足 12.财务规范性/主体清晰度可疑 13.客户质量不足 14.公开信息极度稀缺
注意：(12)不公开财务数据属正常，仅缺少财务数据、财报或审计报表不得标红，只能写为待核实；只有主体不清、造假、处罚、冻结、失信等实质异常才可标红。(14)搜索结果少不等于公司有问题。

硬规则（最终由程序复算）：
- ABCD评级只衡量企业质量，不纳入角色定位或交易可行性；企业质量分为行业匹配度 10%、RR 30%、客户粘性 15%、商业模式 20%、市场地位 25%，合计 100%。
- 企业质量评级：≥3.8→A，3.0-3.7→B，2.0-2.9→C，＜2.0→D。
- 推进建议：规模明确/大概率过大→不建议推进；其余规模情形下A→推进、B→观察、C/D→不建议推进；transaction_feasibility为evidence_no→强制不建议推进；inferred_no→推进建议旁标⚠。未公开出售意愿或一般性的股权待核实不应把A降为观察。
- 严重红旗或行业匹配度≤2→不建议推进；可推进但证据可信度为C/D→降为观察并提示补证。
- 红旗严重度：#9或#12为🔴，或任意🔴达到3项→严重；其他存在🔴→中等；只有🟡→轻微

输出严格JSON（不要代码块）：
{"scores":[4,3,3,3,3],"evidence_levels":[3,2,2,2,2,2,2],"total":3.1,"grade":"C","verdict":"WATCH","rr":"中","role_recommendation":"addon","role_recommendation_reason":"","transaction_feasibility":"inferred_yes","transaction_feasibility_reason":"","red_flags":{"1":"🟢","2":"🟢","3":"🟢","4":"🟢","5":"🟢","6":"🟢","7":"🟢","8":"🟢","9":"🟢","10":"🟢","11":"🟢","12":"🟢","13":"🟢","14":"🟢"},"red_count":0,"red_severity":"轻微","summary":"一句话摘要（50字以内）","industry":"行业","revenue_est":"营收推断","revenue_evidence":"A/B/C/D","analysis":"每维度一句话分析逻辑（如：维度1-食品OEM核心赛道5分 维度2-有长期代工合同4分 维度3-市场地位3分 ...）","refs":["来源1标题","来源2标题"]}"""


async def quick_score_company(company_name: str) -> dict:
    """快速评分单家公司：按简化尽调主题检索 + AI一步输出含分析和来源数"""
    search_plan = _build_quick_search_plan(company_name)
    seen_urls: set[str] = set()

    search_text = "## 快速检索框架：主营业务、经营持续性、规模、股权、风险五个主题\n"
    ref_counter = 1
    for item in search_plan:
        q = item["query"]
        try:
            results = await public_web_search(
                q, company_name=company_name, max_results=8, search_depth="basic"
            )
        except Exception as e:
            print(f"[QuickScore] Search failed for {company_name}/{q}: {e}")
            continue
        results = _filter_verified_search_results(results, company_name, seen_urls=seen_urls)
        topic_label = _SEARCH_TOPIC_LABELS.get(item["topic"], "补充检索")
        search_text += f"\n### [{topic_label}] 搜索：{q}\n"
        for r in results:
            if r["type"] == "result":
                search_text += f"[{ref_counter}] {r['title']}\nURL: {r['url']}\n{r['content'][:600]}\n\n"
                ref_counter += 1
            elif r["type"] == "answer":
                search_text += f"摘要：{r['content'][:400]}\n"

    # 百度搜索补充
    baidu_results = await search_baidu(f"{company_name} 主营 工商 风险", max_results=5)
    baidu_results = _filter_verified_search_results(baidu_results, company_name, seen_urls=seen_urls)
    if baidu_results:
        search_text += "\n\n## 百度搜索结果\n"
        for r in baidu_results:
            search_text += f"[{ref_counter}] {r['title']}\nURL: {r['url']}\n{r['content'][:400]}\n\n"
            ref_counter += 1

    sources_count = ref_counter - 1
    if len(search_text) > 15000:
        search_text = search_text[:15000]

    user_msg = f"公司名称：{company_name}\n\n## 搜索结果（共{sources_count}条来源，含网页搜索和百度）\n{search_text}\n\n请输出该公司的快速评分JSON（必须包含analysis字段写分析逻辑和refs字段写来源URL列表）。"
    messages = [
        {"role": "system", "content": SYSTEM_PROMPT_QUICK_SCORE},
        {"role": "user", "content": user_msg},
    ]

    response_text = ""
    async with httpx.AsyncClient(timeout=httpx.Timeout(300.0, connect=30.0)) as client:
        resp = await client.post(
            f"{DEEPSEEK_BASE_URL}/v1/chat/completions",
            headers={"Authorization": f"Bearer {DEEPSEEK_API_KEY}", "Content-Type": "application/json"},
            json={"model": DEEPSEEK_MODEL, "messages": messages, "stream": False, "max_tokens": 8192, "temperature": 0},
        )
        if resp.status_code != 200:
            raise Exception(f"DeepSeek error {resp.status_code}: {resp.text[:300]}")
        data = resp.json()
        response_text = data["choices"][0]["message"].get("content", "") or ""

    # Parse JSON
    json_text = response_text.strip()
    m = re.search(r'\{[^{}]*"scores"\s*:\s*\[[^\]]*\][^{}]*\}', json_text, re.DOTALL)
    if m:
        json_text = m.group(0)
    else:
        start = json_text.find('{')
        end = json_text.rfind('}') + 1
        if start != -1 and end > start:
            json_text = json_text[start:end]

    result = json.loads(json_text)
    result["company_name"] = company_name
    result["sources_count"] = sources_count
    result = _apply_screening_score(result, quick=True)
    result.setdefault("rr", "中")
    result.setdefault("red_count", 0)
    result.setdefault("red_severity", "轻微")
    result.setdefault("summary", "")
    result.setdefault("industry", "待确认")
    result.setdefault("revenue_est", "待推算")
    result.setdefault("analysis", "")
    result.setdefault("refs", [])
    return result


async def _collect_batch_full_evidence(record: dict) -> tuple[str, list[str], int]:
    """Collect the same diligence layers used in a full screen without rendering HTML."""
    company_name = record["company_name"]
    seen_urls: set[str] = set()
    plan = _build_screening_search_plan(company_name)
    raw_results = await multi_search(
        [item["query"] for item in plan], company_name=company_name, comprehensive=True
    )
    search_results = _govern_search_result_groups(raw_results, company_name, "", seen_urls)

    baidu_plan = _build_baidu_search_plan(company_name)
    raw_baidu = {
        item["query"]: await search_baidu(item["query"], max_results=6)
        for item in baidu_plan
    }
    baidu_results = _govern_search_result_groups(raw_baidu, company_name, "", seen_urls)
    coverage_text = _coverage_text_from_results(search_results, baidu_results)
    coverage = _assess_search_coverage(coverage_text)
    gap_plan = _build_gap_search_plan(company_name, coverage)
    if gap_plan:
        raw_gap = await multi_search(
            [item["query"] for item in gap_plan], company_name=company_name, comprehensive=True
        )
        gap_results = _govern_search_result_groups(raw_gap, company_name, "", seen_urls)
        search_results.update(gap_results)
        raw_baidu_gap = {
            item["query"]: await search_baidu(item["query"], max_results=4)
            for item in gap_plan
        }
        baidu_results.update(
            _govern_search_result_groups(raw_baidu_gap, company_name, "", seen_urls)
        )
        coverage_text = _coverage_text_from_results(search_results, baidu_results)

    wechat_results = []
    for query in _wechat_supplement_queries(company_name, "", "", coverage_text):
        wechat_results.extend(await search_wechat_articles(query, max_results=5))
    wechat_results = _filter_relevant_wechat_results(
        _dedupe_wechat_results(wechat_results, limit=_WECHAT_MAX_ARTICLES * 2),
        company_name,
        "",
    )

    search_text = "## 检索方法与覆盖度\n采用单家深筛同等的六主题网页检索、缺口补搜、百度中文补充及微信公众号经营事实交叉验证。\n"
    imported_fields = record.get("imported_fields", {})
    if imported_fields:
        search_text += "\n## 导入的工商基础信息（企查查导出，作为主体与规模线索，不等于营收）\n"
        for field, value in imported_fields.items():
            search_text += f"- {field}: {value}\n"
    qcc_text = ""
    qcc_refs: list[dict] = []
    has_basic_registry = imported_fields and all(
        imported_fields.get(k) for k in ("注册资本", "法定代表人", "成立时间")
    )
    if QCC_MCP_KEY and not has_basic_registry:
        qcc_text, qcc_refs = await collect_qcc_evidence(company_name)
        if qcc_text:
            search_text += "\n" + qcc_text + "\n"

    refs: list[str] = []
    source_urls: list[str] = []
    ref_counter = 1
    for ref in qcc_refs:
        title = ref.get("title", "")
        url = ref.get("url", "")
        provider = ref.get("provider", "企查查MCP")
        refs.append(f"[{provider}] {title} | {url}")
        if url:
            source_urls.append(url)
    for query, results in search_results.items():
        search_text += f"\n### 网页搜索：{query}\n"
        for result in results:
            if result.get("type") == "answer":
                search_text += f"摘要：{result.get('content', '')[:500]}\n"
            elif result.get("type") == "result":
                title = result.get("title", "")
                url = result.get("url", "")
                provider = result.get("provider", "网页搜索")
                search_text += f"[{ref_counter}] [{provider}] {title}\nURL: {url}\n{result.get('content', '')[:800]}\n\n"
                refs.append(f"[{provider}] {title} | {url}")
                if url:
                    source_urls.append(url)
                ref_counter += 1
    for query, results in baidu_results.items():
        search_text += f"\n### 百度搜索：{query}\n"
        for result in results:
            title = result.get("title", "")
            url = result.get("url", "")
            search_text += f"[{ref_counter}] [百度] {title}\nURL: {url}\n{result.get('content', '')[:500]}\n\n"
            refs.append(f"[百度] {title} | {url}")
            ref_counter += 1
    if wechat_results:
        search_text += "\n## 微信公众号补充来源（仅用于经营事实补充和交叉验证）\n"
        for index, article in enumerate(wechat_results):
            content = ""
            if index < _WECHAT_MAX_FULLTEXT and article.get("url"):
                detail = await fetch_wechat_article_content(article["url"])
                content = detail.get("content", "")[:2000] if detail else ""
            title = article.get("title", "")
            url = article.get("url", "")
            account = article.get("account", "")
            search_text += f"[{ref_counter}] [微信补充] {title} - {account}\nURL: {url}\n{content or article.get('digest', '')}\n\n"
            refs.append(f"[微信补充] {_plain_fragment(title)} - {account} | {url}")
            ref_counter += 1

    crawl_urls = list(dict.fromkeys(source_urls))[:6]
    if crawl_urls:
        try:
            crawled = await asyncio.wait_for(
                crawl_urls_with_playwright(crawl_urls, max_concurrent=3),
                timeout=60,
            )
            deep_text = ""
            for url, info in crawled.items():
                if info.get("content") and info.get("length", 0) > 200:
                    deep_text += f"\n### [深度抓取] {info.get('title', url)}\nURL: {url}\n{info['content'][:2000]}\n"
            if deep_text:
                search_text += "\n## 深度抓取内容\n" + deep_text
        except Exception as exc:
            print(f"[BatchFull] Crawl skipped for {company_name}: {type(exc).__name__}: {exc}")
    return search_text[:45000], refs, ref_counter - 1


async def full_score_company(record: dict) -> dict:
    """Perform full diligence analysis for a batch row and retain structured fields only."""
    company_name = record["company_name"]
    search_text, refs, source_count = await _collect_batch_full_evidence(record)
    user_message = (
        f"请对以下公司进行PE并购深筛分析：\n\n公司名称：{company_name}\n\n"
        f"## 已核验检索与导入信息\n\n{search_text}\n\n请完成5个维度+角色定位+交易可行性的深筛分析。"
    )
    analysis_md, ai_refs_html, _ = await call_deepseek_with_tools(
        [{"role": "system", "content": SYSTEM_PROMPT_ANALYSIS}, {"role": "user", "content": user_message}],
        tools=ANALYSIS_TOOLS,
        max_rounds=2,
        ref_counter_start=source_count + 1,
        company_name=company_name,
    )
    if not analysis_md.strip() or analysis_md.startswith("[思维链"):
        raise ValueError("AI完整分析返回为空")
    # Try direct JSON extraction first
    json_match = re.search(r"---JSON---\s*(.*?)\s*---END_JSON---", analysis_md, re.DOTALL)
    if json_match:
        try:
            response = json_match.group(1)
            print("[Batch] JSON extracted directly from analysis response")
        except Exception:
            json_match = None
    if not json_match:
        response = await call_deepseek(
            [{"role": "system", "content": SYSTEM_PROMPT_JSON},
             {"role": "user", "content": f"将以下深筛分析转为JSON格式：\n\n{analysis_md}"}],
            retry_count=1,
        )
    analysis = _apply_screening_score(_extract_json(response))
    _self_critique(analysis, company_name)
    verdict = {"pass": "PASS", "watch": "WATCH", "reject": "REJECT"}.get(
        analysis.get("verdict_type", "watch"), "WATCH"
    )
    analysis.update({
        "company_name": company_name,
        "imported_fields": record.get("imported_fields", {}),
        "registered_capital": record.get("registered_capital", ""),
        "legal_representative": record.get("legal_representative", ""),
        "established_at": record.get("established_at", ""),
        "registration_status": record.get("registration_status", ""),
        "credit_code": record.get("credit_code", ""),
        "sources_count": source_count,
        "refs": refs,
        "scores": analysis.get("radar_scores", []),
        "total": analysis.get("total_score", 0),
        "verdict": verdict,
        "rr": analysis.get("recurring_revenue", "待确认"),
        "red_count": analysis.get("red_flags_count", 0),
        "red_severity": analysis.get("red_flags_severity", "轻微"),
        "summary": analysis.get("exec_summary", ""),
        "revenue_est": analysis.get("revenue_estimate", ""),
        "analysis": "\n".join(
            _plain_fragment(analysis.get(field, ""))
            for field in ("sec1_summary", "sec2_summary", "sec3_summary", "sec4_summary", "sec5_summary", "sec6_summary", "sec7_summary")
            if analysis.get(field)
        ),
    })
    return analysis


# ---- 排除结果工厂 ----
def _make_excluded_result(record: dict, exclude_code: str, exclude_reason: str) -> dict:
    name = record.get("company_name", "")
    return {
        "company_name": name, "excluded": True, "exclude_code": exclude_code,
        "exclude_reason": exclude_reason, "scores": [0] * 5, "total": 0,
        "grade": "EXCLUDED", "verdict": "EXCLUDED",
        "advance_recommendation": "被排除", "advance_recommendation_reason": exclude_reason,
        "rr": "-", "red_flags": {}, "red_count": 0, "red_severity": "-",
        "summary": f"被排除（{exclude_code}）：{exclude_reason}",
        "industry": "?", "revenue_est": "?", "evidence_confidence": "-",
        "role_recommendation": "", "transaction_feasibility": "",
        "imported_fields": record.get("imported_fields", {}),
        "registered_capital": record.get("registered_capital", ""),
        "legal_representative": record.get("legal_representative", ""),
        "established_at": record.get("established_at", ""),
        "registration_status": record.get("registration_status", ""),
        "credit_code": record.get("credit_code", ""),
        "sources_count": 0, "analysis": "", "refs": [],
    }


def _make_error_result(record: dict, error_msg: str) -> dict:
    name = record.get("company_name", "")
    return {
        "company_name": name, "scores": [0] * 5, "total": 0, "grade": "N/A",
        "verdict": "ERROR", "rr": "?", "red_flags": {}, "red_count": 0,
        "red_severity": "?", "summary": f"错误：{error_msg[:100]}",
        "industry": "?", "revenue_est": "?",
        "advance_recommendation": "错误", "advance_recommendation_reason": error_msg[:100],
        "imported_fields": record.get("imported_fields", {}),
        "registered_capital": record.get("registered_capital", ""),
        "legal_representative": record.get("legal_representative", ""),
        "established_at": record.get("established_at", ""),
        "registration_status": record.get("registration_status", ""),
        "credit_code": record.get("credit_code", ""),
    }


async def run_batch_scoring(batch_id: str, company_records: list[dict]):
    """两阶段批量分析：表格预筛选 → 仅通过者进入完整 AI 深筛。"""
    batch = batch_tasks[batch_id]
    batch["_full_list"] = company_records
    total = len(company_records)
    progress_lock = asyncio.Lock()

    try:
        await _run_batch_scoring_impl(batch_id, batch, company_records, total, progress_lock)
    except Exception as e:
        batch["status"] = "error"
        batch["progress"] = f"批量分析崩溃：{e}"
        _persist_batch_task(batch_id, batch)
        print(f"[Batch] CRASH: {e}")
        import traceback
        traceback.print_exc()


async def _run_batch_scoring_impl(batch_id: str, batch: dict, company_records: list[dict],
                                  total: int, progress_lock: asyncio.Lock):
    """实际的批量分析流程（与 run_batch_scoring 分离，便于顶层 try-except 保护）"""

    # ====== Phase 1: 表格预筛选（纯同步，零 API） ======
    prescreen_enabled = batch.get("prescreen_enabled", True)
    if prescreen_enabled:
        batch["phase"] = "pre_screen"
        batch["progress"] = "正在进行表格预筛选（第1/2阶段）…"
        _persist_batch_task(batch_id, batch)

        excluded_results: list[dict] = []
        passed_records: list[dict] = []

        try:
            for i, record in enumerate(company_records):
                name = record.get("company_name", "").strip()
                if not name:
                    continue
                ruled = _apply_table_prescreen_rules(
                    record, exclude_foreign=batch.get("exclude_foreign", False)
                )
                if ruled:
                    excluded_results.append(
                        _make_excluded_result(record, ruled["exclude_code"], ruled["exclude_reason"])
                    )
                    print(f"[PreScreen] {name} → EXCLUDED {ruled['exclude_code']}: {ruled['exclude_reason']}")
                else:
                    passed_records.append(record)
        except Exception as e:
            batch["status"] = "error"
            batch["progress"] = f"预筛选阶段崩溃：{e}"
            _persist_batch_task(batch_id, batch)
            print(f"[PreScreen] CRASH: {e}")
            import traceback
            traceback.print_exc()
            return

        excluded_count = len(excluded_results)
        passed_count = len(passed_records)
        batch["excluded_count"] = excluded_count
        batch["passed_count"] = passed_count
        batch["_passed_records"] = passed_records  # for resume
        batch["completed"] = excluded_count
        batch["results"] = list(excluded_results)
        _persist_batch_task(batch_id, batch)

        if not passed_records:
            batch["phase"] = "completed"
            batch["status"] = "generating"
            batch["progress"] = "正在生成 Excel…"
            excel_path = REPORT_DIR / f"batch_{batch_id}.xlsx"
            _generate_batch_excel(excluded_results, excel_path)
            batch["status"] = "completed"
            batch["progress"] = f"全部 {total} 家公司均被预筛排除，无公司进入完整分析"
            batch["completed"] = total
            batch["excel_path"] = str(excel_path)
            _persist_batch_task(batch_id, batch)
            return
    else:
        # 跳过预筛选，所有公司直接进入完整分析
        passed_records = list(company_records)
        batch["excluded_count"] = 0
        batch["passed_count"] = len(passed_records)
        batch["_passed_records"] = passed_records
        batch["results"] = []
        batch["completed"] = 0
        _persist_batch_task(batch_id, batch)

    # ====== Phase 1 done ======
    if prescreen_enabled:
        # Stop here — user must explicitly trigger Phase 2
        batch["phase"] = "prescreen_done"
        batch["status"] = "prescreen_done"
        batch["progress"] = f"预筛选完成：{batch.get('passed_count',0)}家通过，{batch.get('excluded_count',0)}家排除。点击「开始深筛分析」继续。"
        _persist_batch_task(batch_id, batch)
        return

    # else: prescreen disabled → go straight to analysis
    await _batch_analysis_phase(batch_id, batch)

async def _batch_analysis_phase(batch_id: str, batch: dict | None = None):
    """Execute Phase 2: full AI analysis on passed companies."""
    if batch is None:
        batch = batch_tasks[batch_id]
        if batch.get("status") in ("completed",):
            return
        passed_records = batch.get("_passed_records", [])
        if not passed_records:
            batch["status"] = "completed"
            batch["phase"] = "completed"
            batch["progress"] = "没有通过预筛的公司"
            _persist_batch_task(batch_id, batch)
            return
        batch["status"] = "running"
        batch["phase"] = "analyzing"
        batch["progress"] = f"预筛完成：{batch.get('passed_count',0)}家进入完整分析…"
        _persist_batch_task(batch_id, batch)

    passed_records = batch.get("_passed_records", [])
    if not passed_records:
        batch["status"] = "completed"
        batch["phase"] = "completed"
        batch["progress"] = "没有可分析的公司"
        _persist_batch_task(batch_id, batch)
        return
    if len(passed_records) > 100:
        batch["status"] = "error"
        batch["progress"] = f"通过预筛的公司（{len(passed_records)}家）超过深筛上限（100家），请拆分后重试"
        _persist_batch_task(batch_id, batch)
        return

    excluded_results = list(batch.get("results", []))
    excluded_count = batch.get("excluded_count", 0)
    total = batch.get("total", 0)

    # ====== Phase 2: 完整 AI 深筛（仅通过预筛的公司） ======
    excluded_count = batch.get("excluded_count", 0)
    passed_count = batch.get("passed_count", 0)
    if "excluded_results" not in dir():
        excluded_results = []
    batch["phase"] = "analyzing"
    prescreen_label = "预筛完成：" if batch.get("prescreen_enabled", True) else "跳过预筛选："
    batch["progress"] = f"{prescreen_label}{passed_count}家，开始完整分析…"
    _persist_batch_task(batch_id, batch)

    progress_lock = asyncio.Lock()
    analysis_slots: dict[int, dict] = {}
    analysis_completed = 0
    analysis_total = len(passed_records)
    sem = asyncio.Semaphore(3)

    async def analyze_one(idx: int, record: dict) -> tuple[int, dict]:
        async with sem:
            name = record.get("company_name", "").strip()
            nonlocal analysis_completed
            async with progress_lock:
                analysis_completed += 1
                total_completed = excluded_count + analysis_completed
                batch["current"] = name
                batch["completed"] = total_completed
                batch["progress"] = (
                    f"正在完整分析 [{analysis_completed}/{analysis_total}] {name}"
                    f"（已排除{excluded_count}家）"
                )
                _persist_batch_task(batch_id, batch)

            try:
                result = await asyncio.wait_for(
                    full_score_company(record), timeout=_BATCH_ANALYSIS_TIMEOUT
                )
                print(f"[Batch] [{analysis_completed}/{analysis_total}] {name}: "
                      f"{result.get('grade', '?')} {result.get('verdict', '?')}")
                return (idx, result)
            except asyncio.TimeoutError:
                print(f"[Batch] [{analysis_completed}/{analysis_total}] {name} TIMEOUT after {_BATCH_ANALYSIS_TIMEOUT}s")
                return (idx, _make_error_result(record, f"分析超时（{_BATCH_ANALYSIS_TIMEOUT}秒），跳过"))
            except Exception as e:
                print(f"[Batch] [{analysis_completed}/{analysis_total}] {name} FAILED: {e}")
                return (idx, _make_error_result(record, str(e)))

    pending = [
        analyze_one(idx, rec)
        for idx, rec in enumerate(passed_records)
        if rec.get("company_name", "").strip()
    ]

    for coro in asyncio.as_completed(pending):
        idx, result = await coro
        analysis_slots[idx] = result
        # Rebuild merged results list: excluded + analyzed (in original passed order)
        merged: list[dict] = list(excluded_results)
        for i in range(len(passed_records)):
            if i in analysis_slots:
                merged.append(analysis_slots[i])
        batch["results"] = merged
        batch["completed"] = excluded_count + len(analysis_slots)
        _persist_batch_task(batch_id, batch)

        # Quota check (Tavily-only mode)
        if SEARCH_PROVIDER not in {"exa", "hybrid"}:
            d = _load_keys()
            ki = d["current_index"]
            current_key = d["keys"][ki]
            used = current_key.get("credits_used", 0)
            rotated = _rotate_key_if_needed(reserve=20)
            current_limit = int(current_key.get("credits_limit") or TAVILY_MONTHLY_LIMIT)
            if not rotated and used >= current_limit - 20:
                all_full = all(
                    k.get("credits_used", 0) >= int(k.get("credits_limit") or TAVILY_MONTHLY_LIMIT) - 20
                    for k in d["keys"]
                )
                if all_full:
                    batch["status"] = "paused_quota"
                    batch["progress"] = f"所有Tavily Key额度耗尽，已处理 {len(merged)}/{total}"
                    batch["paused_at"] = len(merged)
                    _persist_batch_task(batch_id, batch)
                    return

    # Final merge
    final_results: list[dict] = list(excluded_results)
    for i in range(len(passed_records)):
        if i in analysis_slots:
            final_results.append(analysis_slots[i])
        else:
            # record that somehow didn't get analyzed — mark as error
            rec = passed_records[i]
            final_results.append(_make_error_result(rec, "分析未完成"))

    batch["status"] = "generating"
    batch["phase"] = "completed"
    batch["progress"] = "正在生成 Excel…"

    excel_path = REPORT_DIR / f"batch_{batch_id}.xlsx"
    _generate_batch_excel(final_results, excel_path)

    batch["status"] = "completed"
    batch["progress"] = f"完成 {len(final_results)}/{total} 家（{excluded_count}家被预筛排除）"
    batch["completed"] = len(final_results)
    batch["results"] = final_results
    batch["excel_path"] = str(excel_path)
    _persist_batch_task(batch_id, batch)


async def _run_analysis_phase(batch_id: str, passed_records: list[dict]):
    """独立的分析阶段任务 — 仅对已通过预筛的公司运行 full_score_company。供 resume 使用。"""
    batch = batch_tasks[batch_id]
    excluded_count = batch.get("excluded_count", 0)
    total = batch.get("total", 0)
    progress_lock = asyncio.Lock()
    analysis_slots: dict[int, dict] = {}
    analysis_completed = 0
    analysis_total = len(passed_records)
    sem = asyncio.Semaphore(3)

    async def analyze_one(idx: int, record: dict) -> tuple[int, dict]:
        async with sem:
            name = record.get("company_name", "").strip()
            nonlocal analysis_completed
            async with progress_lock:
                analysis_completed += 1
                total_completed = excluded_count + analysis_completed
                batch["current"] = name
                batch["completed"] = total_completed
                batch["progress"] = (
                    f"正在完整分析 [{analysis_completed}/{analysis_total}] {name}"
                    f"（已排除{excluded_count}家）"
                )
                _persist_batch_task(batch_id, batch)
            try:
                result = await asyncio.wait_for(
                    full_score_company(record), timeout=_BATCH_ANALYSIS_TIMEOUT
                )
                print(f"[Batch-Resume] [{analysis_completed}/{analysis_total}] {name}: "
                      f"{result.get('grade', '?')} {result.get('verdict', '?')}")
                return (idx, result)
            except asyncio.TimeoutError:
                print(f"[Batch-Resume] [{analysis_completed}/{analysis_total}] {name} TIMEOUT after {_BATCH_ANALYSIS_TIMEOUT}s")
                return (idx, _make_error_result(record, f"分析超时（{_BATCH_ANALYSIS_TIMEOUT}秒），跳过"))
            except Exception as e:
                print(f"[Batch-Resume] [{analysis_completed}/{analysis_total}] {name} FAILED: {e}")
                return (idx, _make_error_result(record, str(e)))

    pending = [analyze_one(idx, rec) for idx, rec in enumerate(passed_records) if rec.get("company_name", "").strip()]
    for coro in asyncio.as_completed(pending):
        idx, result = await coro
        analysis_slots[idx] = result
        excluded_results = [r for r in batch.get("results", []) if r.get("excluded")]
        merged = list(excluded_results)
        for i in range(len(passed_records)):
            if i in analysis_slots:
                merged.append(analysis_slots[i])
        batch["results"] = merged
        batch["completed"] = excluded_count + len(analysis_slots)
        _persist_batch_task(batch_id, batch)

    excluded_results = [r for r in batch.get("results", []) if r.get("excluded")]
    final_results = list(excluded_results)
    for i in range(len(passed_records)):
        if i in analysis_slots:
            final_results.append(analysis_slots[i])
        else:
            final_results.append(_make_error_result(passed_records[i], "分析未完成"))

    batch["status"] = "generating"
    batch["phase"] = "completed"
    batch["progress"] = "正在生成 Excel…"
    excel_path = REPORT_DIR / f"batch_{batch_id}.xlsx"
    _generate_batch_excel(final_results, excel_path)
    batch["status"] = "completed"
    batch["progress"] = f"完成 {len(final_results)}/{total} 家（{excluded_count}家被预筛排除）"
    batch["completed"] = len(final_results)
    batch["results"] = final_results
    batch["excel_path"] = str(excel_path)
    _persist_batch_task(batch_id, batch)


def _excel_role_label(val):
    _map = {"platform": "Platform", "addon": "Add-on", "both": "两者皆可", "neither": "皆不适合"}
    return _map.get(val, val or "")

def _excel_feasibility_label(val):
    _map = {"evidence_yes": "有明确证据可推进", "inferred_yes": "推断可推进", "inferred_no": "推断存在障碍", "evidence_no": "有明确证据不可推进"}
    return _map.get(val, val or "")

def _role_display_label(role, fallback):
    _map = {"platform": "Platform", "addon": "Add-on", "both": "两者皆可", "neither": "皆不适合"}
    if role in _map:
        return _map[role]
    return fallback or "待判断"

def _generate_batch_excel(results: list[dict], path: Path):
    """Generate a structured, human-readable workbook with logical column grouping."""

    # ---- Column definitions: (label, width, center?, wrap?) ----
    COLUMNS = [
        # Section 1: 筛选 + 公司
        ("预筛结果", 26, True,  False),
        ("公司名称", 22, False, False),
        # Section 2: 工商信息
        ("注册资本", 12, True,  False),
        ("法定代表人", 12, False, False),
        ("成立日期", 12, True,  False),
        ("经营状态", 10, True,  False),
        ("行业",       14, False, False),
        # Section 3: 评分总览
        ("企业质量分", 10, True,  False),
        ("评级",        8, True,  False),
        ("行业匹配度\n10%",   10, True, False),
        ("RR\n30%",          10, True, False),
        ("市场地位\n25%",    10, True, False),
        ("商业模式\n20%",    10, True, False),
        ("客户粘性\n15%",    10, True, False),
        ("RR判断",      8, True,  False),
        ("证据可信度", 10, True,  False),
        ("搜索来源",   8, True,  False),
        # Section 4: 规模与推进
        ("规模推断",     16, False, False),
        ("规模适配",     24, False, True),
        ("角色定位",     14, True,  False),
        ("交易可行性",   20, False, True),
        ("交易障碍",     24, False, True),
        ("推进建议",     12, True,  False),
        ("推进理由",     32, False, True),
        # Section 5: 红旗与摘要
        ("红旗",     6, True,  False),
        ("严重度",   8, True,  False),
        ("执行摘要", 44, False, True),
        # Section 6: 细项分析
        ("身份核验",       28, False, True),
        ("产品与服务",     28, False, True),
        ("RR分析",         28, False, True),
        ("客户粘性分析",   28, False, True),
        ("竞争位置",       28, False, True),
        ("Platform/Add-on", 28, False, True),
        ("规模与交易障碍", 28, False, True),
        # Section 7: 其他
        ("关键问题",   24, False, True),
        ("接触切入点", 24, False, True),
        ("参考来源",   40, False, True),
        ("导入工商字段", 30, False, True),
    ]

    headers = [c[0] for c in COLUMNS]
    center_cols = frozenset(i + 1 for i, c in enumerate(COLUMNS) if c[2])
    wrap_cols = frozenset(i + 1 for i, c in enumerate(COLUMNS) if c[3])

    def export_row(result: dict) -> list:
        scores = (result.get("scores", []) + [0] * 5)[:5]
        imported = result.get("imported_fields", {})
        imported_text = "；".join(f"{k}: {v}" for k, v in imported.items())
        # Strip credit codes / biz type — they're already in imported_fields
        return [
            f"{result.get('exclude_code')}: {result.get('exclude_reason')}"
            if result.get("excluded") else "通过",
            result.get("company_name", ""),
            result.get("registered_capital", ""),
            result.get("legal_representative", ""),
            result.get("established_at", ""),
            result.get("registration_status", ""),
            result.get("industry", ""),
            result.get("total", 0),
            result.get("grade", ""),
            *scores,
            result.get("rr", ""),
            result.get("evidence_confidence", ""),
            result.get("sources_count", 0),
            result.get("revenue_est", ""),
            result.get("scale_fit_label", "") + "：" + (result.get("scale_fit_reason") or ""),
            _excel_role_label(result.get("role_recommendation", "")),
            _excel_feasibility_label(result.get("transaction_feasibility", "")),
            result.get("transaction_obstacle_label", "") + "：" + (result.get("transaction_obstacle_reason") or ""),
            result.get("advance_recommendation", ""),
            result.get("advance_recommendation_reason", ""),
            result.get("red_count", 0),
            result.get("red_severity", ""),
            result.get("summary", ""),
            _plain_fragment(result.get("sec1_content", "") or result.get("sec1_summary", "")),
            _plain_fragment(result.get("sec2_content", "") or result.get("sec2_summary", "")),
            _plain_fragment(result.get("sec3_content", "") or result.get("sec3_summary", "")),
            _plain_fragment(result.get("sec4_content", "") or result.get("sec4_summary", "")),
            _plain_fragment(result.get("sec6_content", "") or result.get("sec6_summary", "")),
            _plain_fragment(result.get("sec7_content", "") or result.get("sec7_summary", "")),
            _plain_fragment(result.get("sec8_summary", "")),
            result.get("key_question", ""),
            result.get("approach", ""),
            "\n".join(result.get("refs", [])),
            imported_text,
        ]

    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    except ImportError:
        import csv
        csv_path = path.with_suffix('.csv')
        with open(csv_path, 'w', newline='', encoding='utf-8-sig') as f:
            writer = csv.writer(f)
            writer.writerow(headers)
            for r in results:
                writer.writerow(export_row(r))
        return

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "批量完整分析结果"

    header_font = Font(bold=True, color="FFFFFF", size=10)
    header_fill = PatternFill(start_color="1B2A4A", end_color="1B2A4A", fill_type="solid")
    pass_fill   = PatternFill(start_color="E8F5E9", end_color="E8F5E9", fill_type="solid")
    reject_fill = PatternFill(start_color="FFEBEE", end_color="FFEBEE", fill_type="solid")
    watch_fill  = PatternFill(start_color="FFF8E1", end_color="FFF8E1", fill_type="solid")
    excl_fill   = PatternFill(start_color="F5F5F5", end_color="F5F5F5", fill_type="solid")
    error_fill  = PatternFill(start_color="FFF3E0", end_color="FFF3E0", fill_type="solid")
    thin_border = Border(
        left=Side(style='thin'), right=Side(style='thin'),
        top=Side(style='thin'), bottom=Side(style='thin'),
    )

    # Header row
    for col, label in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=label)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
        cell.border = thin_border

    # Data rows
    for row_idx, r in enumerate(results, 2):
        row_data = export_row(r)
        for col, val in enumerate(row_data, 1):
            cell = ws.cell(row=row_idx, column=col, value=val)
            cell.border = thin_border
            halign = 'center' if col in center_cols else 'left'
            cell.alignment = Alignment(horizontal=halign, vertical='top',
                                       wrap_text=col in wrap_cols)

        verdict = r.get("verdict", "WATCH")
        if r.get("excluded"):
            row_fill = excl_fill
        elif verdict == "PASS":
            row_fill = pass_fill
        elif verdict == "REJECT":
            row_fill = reject_fill
        elif verdict == "ERROR":
            row_fill = error_fill
        else:
            row_fill = watch_fill
        for c in range(1, len(headers) + 1):
            ws.cell(row=row_idx, column=c).fill = row_fill

    # Column widths
    for i, (_, width, _, _) in enumerate(COLUMNS, 1):
        ws.column_dimensions[openpyxl.utils.get_column_letter(i)].width = width

    # Freeze header row
    ws.freeze_panes = 'A2'
    ws.auto_filter.ref = ws.dimensions
    wb.save(str(path))
    print(f"[Batch] Excel saved: {path}")


# ============ API 路由 ============

@app.get("/", response_class=HTMLResponse)
async def index():
    """首页"""
    return (BASE_DIR / "static" / "index.html").read_text(encoding="utf-8")


@app.post("/api/screen")
async def start_screening(
    company_name: str = Form(...),
    website: str = Form(""),
    industry: str = Form(""),
    bp_file: UploadFile = File(None),
    model: str = Form(""),
):
    """启动深筛任务"""
    provider = "kimi" if model and model.startswith(("kimi-", "moonshot-")) else "deepseek"
    if provider == "deepseek" and not DEEPSEEK_API_KEY:
        raise HTTPException(500, "未配置 DEEPSEEK_API_KEY 环境变量")
    if provider == "kimi" and not KIMI_API_KEY:
        raise HTTPException(500, "未配置 KIMI_API_KEY 环境变量")
    if website:
        normalized_website = website if "://" in website else f"https://{website}"
        if not await _is_safe_public_url(normalized_website):
            raise HTTPException(400, "官网地址必须是可公开访问的 HTTP/HTTPS 网站")

    task_id = str(uuid.uuid4())[:8]
    tasks[task_id] = {
        "status": "created",
        "progress": "任务已创建",
        "company_name": company_name,
        "website": website,
        "industry": industry,
        "created_at": datetime.now().isoformat(),
    }

    # 处理 BP 文件上传
    bp_text = ""
    if bp_file and bp_file.filename:
        suffix = Path(bp_file.filename).suffix
        save_path = UPLOAD_DIR / f"{task_id}{suffix}"
        content = await bp_file.read()
        save_path.write_bytes(content)
        bp_text = await extract_bp_text(save_path)

    # 异步启动分析
    asyncio.create_task(run_screening(task_id, company_name, website, industry, bp_text, model=model))

    return {"task_id": task_id, "company_name": company_name}


@app.get("/api/status/{task_id}")
async def get_status(task_id: str):
    """查询任务状态"""
    if task_id not in tasks:
        raise HTTPException(404, "任务不存在")
    task = tasks[task_id]
    return {
        "task_id": task_id,
        "status": task["status"],
        "progress": task["progress"],
        "company_name": task.get("company_name", ""),
        "live_text": task.get("live_text", ""),
        "grade": (task.get("analysis", {}).get("grade", "") or task.get("grade", "")) if isinstance(task.get("analysis"), dict) else task.get("grade", ""),
        "verdict": (task.get("analysis", {}).get("verdict_type", "") or task.get("verdict", "")) if isinstance(task.get("analysis"), dict) else task.get("verdict", ""),
        "advance_recommendation": (task.get("analysis", {}).get("advance_recommendation", "") or task.get("advance_recommendation", "")) if isinstance(task.get("analysis"), dict) else task.get("advance_recommendation", ""),
        "advance_reason": (task.get("analysis", {}).get("advance_recommendation_reason", "") or task.get("advance_reason", "")) if isinstance(task.get("analysis"), dict) else task.get("advance_reason", ""),
        "report_url": task.get("report_url"),
        "report_filename": task.get("report_filename"),
        "error_detail": task.get("error_detail", ""),
    }


@app.get("/report/{task_id}", response_class=HTMLResponse)
async def view_report(task_id: str):
    """查看报告"""
    if task_id not in tasks:
        raise HTTPException(404, "任务不存在")
    task = tasks[task_id]
    if task["status"] != "completed":
        raise HTTPException(400, "报告尚未完成")
    report_path = REPORT_DIR / task["report_filename"]
    if not report_path.exists():
        raise HTTPException(404, "报告文件不存在")
    return report_path.read_text(encoding="utf-8")


@app.get("/api/download/{task_id}")
async def download_report(task_id: str):
    """下载报告"""
    if task_id not in tasks:
        raise HTTPException(404, "任务不存在")
    task = tasks[task_id]
    if task["status"] != "completed":
        raise HTTPException(400, "报告尚未完成")
    report_path = REPORT_DIR / task["report_filename"]
    if not report_path.exists():
        raise HTTPException(404, "报告文件不存在")
    return FileResponse(
        str(report_path),
        media_type="text/html",
        filename=task["report_filename"],
    )


# ============ 批量完整分析 API ============
def _create_batch_task(records: list[dict], prescreen: bool = True, exclude_foreign: bool = False) -> dict:
    if not records:
        raise HTTPException(400, "请提供公司名单或上传 Excel")
    batch_id = uuid.uuid4().hex[:8]
    batch_tasks[batch_id] = {
        "status": "running",
        "phase": "pre_screen" if prescreen else "analyzing",
        "prescreen_enabled": prescreen,
        "exclude_foreign": exclude_foreign,
        "progress": "准备完整检索与分析…",
        "total": len(records),
        "completed": 0,
        "current": "",
        "results": [],
        "excluded_count": 0,
        "passed_count": 0,
        "created_at": datetime.now().isoformat(),
        "input_type": "excel" if any(record.get("imported_fields") for record in records) else "manual",
    }
    asyncio.create_task(run_batch_scoring(batch_id, records))
    return {"batch_id": batch_id, "total": len(records)}


@app.post("/api/batch-score")
async def start_batch_scoring(request: dict):
    """Start full batch analysis from manually entered legal entity names."""
    if not DEEPSEEK_API_KEY:
        raise HTTPException(500, "未配置 DEEPSEEK_API_KEY 环境变量")
    company_names = request.get("company_names", [])
    if isinstance(company_names, str):
        company_names = [n.strip() for n in company_names.split("\n") if n.strip()]
    prescreen = request.get("prescreen", True)
    exclude_foreign = request.get("exclude_foreign", False)
    return _create_batch_task(_manual_batch_records(company_names), prescreen=prescreen, exclude_foreign=exclude_foreign)


@app.post("/api/batch-upload")
async def start_batch_from_excel(batch_file: UploadFile = File(...), prescreen: bool = Form(True),
                                  exclude_foreign: bool = Form(False)):
    """Start full batch analysis from a Qichacha exported .xlsx workbook."""
    if not DEEPSEEK_API_KEY:
        raise HTTPException(500, "未配置 DEEPSEEK_API_KEY 环境变量")
    filename = batch_file.filename or ""
    if Path(filename).suffix.lower() != ".xlsx":
        raise HTTPException(400, "请上传 .xlsx 格式的企查查导出文件")
    content = await batch_file.read()
    if len(content) > 20 * 1024 * 1024:
        raise HTTPException(400, "Excel 文件过大，请上传不超过 20MB 的文件")
    return _create_batch_task(_parse_batch_excel(content), prescreen=prescreen, exclude_foreign=exclude_foreign)


@app.get("/api/batch-status/{batch_id}")
async def get_batch_status(batch_id: str):
    """查询批量完整分析任务状态"""
    if batch_id not in batch_tasks:
        raise HTTPException(404, "任务不存在")
    b = batch_tasks[batch_id]
    return {
        "batch_id": batch_id,
        "status": b["status"],
        "phase": b.get("phase", "pre_screen"),
        "progress": b["progress"],
        "total": b["total"],
        "completed": b["completed"],
        "current": b.get("current", ""),
        "excluded_count": b.get("excluded_count", 0),
        "passed_count": b.get("passed_count", 0),
        "results": [
            {"company_name": r["company_name"], "total": r["total"], "grade": r["grade"],
             "verdict": r["verdict"], "advance_recommendation": r.get("advance_recommendation", "不建议推进"),
             "advance_reason": r.get("advance_recommendation_reason", ""),
             "rr": r["rr"], "red_count": r["red_count"],
             "red_severity": r["red_severity"], "summary": r["summary"],
             "evidence_confidence": r.get("evidence_confidence", "C"),
             "analysis": r.get("analysis", "")[:200],
             "refs_count": len(r.get("refs", [])),
             "sources_count": r.get("sources_count", 0),
             "registered_capital": r.get("registered_capital", ""),
             "excluded": r.get("excluded", False),
             "exclude_code": r.get("exclude_code", ""),
             "exclude_reason": r.get("exclude_reason", ""),
             }
            for r in b.get("results", [])
        ],
    }


@app.get("/api/batch-detail/{batch_id}/{index}")
async def get_batch_detail(batch_id: str, index: int):
    """查看批量完整分析中某家公司的结构化结果"""
    if batch_id not in batch_tasks:
        raise HTTPException(404, "任务不存在")
    results = batch_tasks[batch_id].get("results", [])
    if index < 0 or index >= len(results):
        raise HTTPException(404, "序号超出范围")
    r = results[index]
    if r.get("excluded"):
        return {
            "company_name": r["company_name"], "excluded": True,
            "exclude_code": r.get("exclude_code", ""),
            "exclude_reason": r.get("exclude_reason", ""),
            "total": 0, "grade": "EXCLUDED",
            "registered_capital": r.get("registered_capital", ""),
            "legal_representative": r.get("legal_representative", ""),
            "established_at": r.get("established_at", ""),
            "registration_status": r.get("registration_status", ""),
            "sources_count": 0, "scores": [], "red_flags": {},
            "red_count": 0, "red_severity": "", "summary": r.get("summary", ""),
            "rr": "", "verdict": "EXCLUDED",
            "advance_recommendation": "被排除",
            "advance_reason": r.get("exclude_reason", ""),
        }
    refs = r.get("refs", [])
    return {
        "company_name": r["company_name"],
        "scores": r.get("scores", []),
        "total": r["total"],
        "grade": r["grade"],
        "verdict": r["verdict"],
        "advance_recommendation": r.get("advance_recommendation", "不建议推进"),
        "advance_reason": r.get("advance_recommendation_reason", ""),
        "scale_fit_label": r.get("scale_fit_label", "未评估"),
        "scale_fit_reason": r.get("scale_fit_reason", ""),
        "transaction_obstacle_label": r.get("transaction_obstacle_label", "未评估"),
        "transaction_obstacle_reason": r.get("transaction_obstacle_reason", ""),
        "rr": r["rr"],
        "red_flags": r.get("red_flags", {}),
        "red_count": r["red_count"],
        "red_severity": r["red_severity"],
        "evidence_confidence": r.get("evidence_confidence", "C"),
        "decision_note": r.get("decision_note", ""),
        "summary": r["summary"],
        "industry": r.get("industry", ""),
        "revenue_est": r.get("revenue_est", ""),
        "analysis": r.get("analysis", ""),
        "refs": refs,
        "sources_count": r.get("sources_count", 0),
        "registered_capital": r.get("registered_capital", ""),
        "legal_representative": r.get("legal_representative", ""),
        "established_at": r.get("established_at", ""),
        "registration_status": r.get("registration_status", ""),
        "business_type": r.get("business_type", ""),
        "sections": {
            "身份核验": _plain_fragment(r.get("sec1_content", "") or r.get("sec1_summary", "")),
            "产品与服务": _plain_fragment(r.get("sec2_content", "") or r.get("sec2_summary", "")),
            "Recurring Revenue": _plain_fragment(r.get("sec3_content", "") or r.get("sec3_summary", "")),
            "客户结构与粘性": _plain_fragment(r.get("sec4_content", "") or r.get("sec4_summary", "")),
            "竞争位置": _plain_fragment(r.get("sec6_content", "") or r.get("sec6_summary", "")),
            "Platform/Add-on": _plain_fragment(r.get("sec7_content", "") or r.get("sec7_summary", "")),
            "规模与交易障碍": _plain_fragment(r.get("sec8_summary", "")),
        },
    }


@app.get("/api/batch-download/{batch_id}")
async def download_batch_excel(batch_id: str):
    """下载批量完整分析Excel"""
    if batch_id not in batch_tasks:
        raise HTTPException(404, "任务不存在")
    b = batch_tasks[batch_id]
    if b["status"] != "completed":
        raise HTTPException(400, "任务未完成")
    path = b.get("excel_path")
    if not path or not Path(path).exists():
        raise HTTPException(404, "Excel文件不存在")
    return FileResponse(str(path), media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                       filename=f"PE深筛_批量完整分析_{batch_id}.xlsx")


@app.post("/api/batch-analyze/{batch_id}")
async def start_batch_analysis(batch_id: str):
    """从 prescreen_done 状态启动完整 AI 深筛分析"""
    if batch_id not in batch_tasks:
        raise HTTPException(404, "任务不存在")
    b = batch_tasks[batch_id]
    if b.get("status") not in ("prescreen_done",) and b.get("phase") not in ("prescreen_done",):
        raise HTTPException(400, f"任务状态为 {b.get('status')}，无法启动分析。请先完成预筛选。")
    asyncio.create_task(_batch_analysis_phase(batch_id))
    return {"status": "ok", "batch_id": batch_id}


@app.get("/api/batch-download-interim/{batch_id}")
async def download_batch_interim(batch_id: str):
    """下载当前批量的已有结果（即使分析尚未完成）"""
    if batch_id not in batch_tasks:
        raise HTTPException(404, "任务不存在")
    b = batch_tasks[batch_id]
    results = b.get("results", [])
    if not results:
        raise HTTPException(400, "暂无已完成的结果")
    # Generate a temp Excel from current results
    import tempfile
    tmp = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
    try:
        _generate_batch_excel(results, Path(tmp.name))
        return FileResponse(
            tmp.name,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=f"PE深筛_批量分析_进行中_{batch_id}.xlsx",
        )
    except Exception:
        if Path(tmp.name).exists():
            Path(tmp.name).unlink(missing_ok=True)
        raise HTTPException(500, "生成 Excel 失败")


@app.get("/api/batch-prescreen-excel/{batch_id}")
async def download_prescreen_excel(batch_id: str):
    """下载预筛选结果Excel — 每家公司的通过/排除状态及原因"""
    if batch_id not in batch_tasks:
        raise HTTPException(404, "任务不存在")
    b = batch_tasks[batch_id]
    passed_records = b.get("_passed_records", [])
    excluded_results = [r for r in b.get("results", []) if r.get("excluded")]

    passed_names = {r.get("company_name", "") for r in passed_records}
    all_companies: list[dict] = []

    # Full list from _full_list preserves original order
    for rec in b.get("_full_list", []):
        name = rec.get("company_name", "").strip()
        if not name:
            continue
        if name in passed_names:
            all_companies.append({
                "company_name": name, "result": "通过", "exclude_code": "", "exclude_reason": "",
                "registered_capital": rec.get("registered_capital", ""),
                "established_at": rec.get("established_at", ""),
                "registration_status": rec.get("registration_status", ""),
                "province": rec.get("province", ""),
                "enterprise_type": rec.get("enterprise_type", ""),
            })
        else:
            # Find the excluded result detail
            excl = next((e for e in excluded_results if e.get("company_name") == name), None)
            all_companies.append({
                "company_name": name, "result": "排除",
                "exclude_code": excl.get("exclude_code", "") if excl else "",
                "exclude_reason": excl.get("exclude_reason", "") if excl else "",
                "registered_capital": rec.get("registered_capital", ""),
                "established_at": rec.get("established_at", ""),
                "registration_status": rec.get("registration_status", ""),
                "province": rec.get("province", ""),
                "enterprise_type": rec.get("enterprise_type", ""),
            })

    import tempfile
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    except ImportError:
        raise HTTPException(500, "服务器未安装 openpyxl")

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "预筛选结果"

    headers = ["序号", "公司名称", "预筛结果", "排除规则", "排除原因", "注册资本", "成立日期", "经营状态", "省份", "企业类型"]
    header_font = Font(bold=True, color="FFFFFF", size=10)
    header_fill = PatternFill(start_color="1B2A4A", end_color="1B2A4A", fill_type="solid")
    pass_fill = PatternFill(start_color="E8F5E9", end_color="E8F5E9", fill_type="solid")
    excl_fill = PatternFill(start_color="FFEBEE", end_color="FFEBEE", fill_type="solid")
    thin_border = Border(
        left=Side(style='thin'), right=Side(style='thin'),
        top=Side(style='thin'), bottom=Side(style='thin'),
    )

    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal='center', vertical='center')
        cell.border = thin_border

    for idx, c in enumerate(all_companies):
        row_idx = idx + 2
        vals = [idx + 1, c["company_name"], c["result"], c["exclude_code"], c["exclude_reason"],
                c["registered_capital"], c["established_at"], c["registration_status"],
                c["province"], c["enterprise_type"]]
        row_fill = pass_fill if c["result"] == "通过" else excl_fill
        for col, val in enumerate(vals, 1):
            cell = ws.cell(row=row_idx, column=col, value=val)
            cell.font = Font(bold=True) if col == 3 else Font()
            cell.fill = row_fill if col == 3 else PatternFill()
            cell.border = thin_border
            cell.alignment = Alignment(horizontal='center' if col in (1, 3, 4) else 'left', vertical='top')

    widths = [6, 24, 10, 10, 40, 14, 14, 12, 10, 22]
    for i, w in enumerate(widths):
        ws.column_dimensions[openpyxl.utils.get_column_letter(i + 1)].width = w

    ws.freeze_panes = 'A2'
    ws.auto_filter.ref = ws.dimensions

    tmp = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
    wb.save(tmp.name)
    return FileResponse(
        tmp.name,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        filename=f"PE深筛_预筛选_{batch_id}.xlsx",
    )


@app.post("/api/batch-resume/{batch_id}")
async def resume_batch(batch_id: str):
    """恢复暂停的批量任务（更换Key后调用）"""
    if batch_id not in batch_tasks:
        raise HTTPException(404, "任务不存在")
    b = batch_tasks[batch_id]
    if b["status"] not in ("paused_quota",):
        raise HTTPException(400, f"任务状态为 {b['status']}，无法恢复")

    # 检查是否已完成预筛选 — 如果是则只恢复分析阶段
    passed_records = b.get("_passed_records")
    if passed_records is not None:
        b["status"] = "running"
        b["phase"] = "analyzing"
        b["progress"] = "已恢复，从分析阶段继续…"
        asyncio.create_task(_run_analysis_phase(batch_id, passed_records))
        return {"status": "resumed", "phase": "analyzing"}

    full_list = b.get("_full_list") or [
        {"company_name": result["company_name"], "imported_fields": result.get("imported_fields", {})}
        for result in b.get("results", [])
    ]
    b["status"] = "running"
    b["progress"] = "已恢复，重新开始…"
    asyncio.create_task(run_batch_scoring(batch_id, full_list))
    return {"status": "resumed"}  # fallback: run everything


@app.post("/api/tavily-key")
async def update_tavily_key(request: dict, x_admin_token: str = Header("")):
    """添加新的 Tavily API Key 到轮换池"""
    _require_admin_token(x_admin_token)
    new_key = request.get("key", "").strip()
    label = request.get("label", "").strip() or f"Key{int(time.time())}"
    if not new_key:
        raise HTTPException(400, "请提供新的 API Key")
    d = _load_keys()
    d["keys"].append({"key": new_key, "credits_used": 0, "label": label})
    _save_keys(d)
    return {"status": "ok", "total_keys": len(d["keys"])}


# ============ 历史报告持久化 ============
HISTORY_FILE = REPORT_DIR / "history.json"

def _load_history():
    if HISTORY_FILE.exists():
        try:
            return json.loads(HISTORY_FILE.read_text(encoding="utf-8"))
        except (json.JSONDecodeError, KeyError):
            pass
    return {}


def _persist_batch_task(batch_id: str, batch: dict) -> None:
    results_path = REPORT_DIR / f"batch_{batch_id}.json"
    persisted = {key: value for key, value in batch.items() if key != "_full_list"}
    results_path.write_text(json.dumps(persisted, ensure_ascii=False), encoding="utf-8")
    history = _load_history()
    history["batch_" + batch_id] = {
        "task_id": "batch_" + batch_id,
        "batch_id": batch_id,
        "company_name": f"批量完整分析 ({batch.get('completed', 0)}/{batch.get('total', 0)})",
        "status": batch.get("status", ""),
        "grade": "-",
        "verdict": "-",
        "advance_recommendation": "-",
        "advance_reason": "-",
        "report_url": "",
        "report_filename": batch.get("excel_path", ""),
        "results_path": str(results_path),
        "created_at": batch.get("created_at", ""),
    }
    HISTORY_FILE.write_text(json.dumps(history, ensure_ascii=False), encoding="utf-8")

def _save_task_history(task_id: str, task: dict):
    h = _load_history()
    analysis = task.get("analysis")
    h[task_id] = {
        "task_id": task_id,
        "company_name": task.get("company_name", ""),
        "status": task.get("status", ""),
        "grade": analysis.get("grade", "-") if isinstance(analysis, dict) else "-",
        "verdict": analysis.get("verdict_type", "-") if isinstance(analysis, dict) else "-",
        "advance_recommendation": analysis.get("advance_recommendation", "-") if isinstance(analysis, dict) else "-",
        "advance_reason": analysis.get("advance_recommendation_reason", "-") if isinstance(analysis, dict) else "-",
        "total_score": analysis.get("total_score", 0) if isinstance(analysis, dict) else 0,
        "analysis": analysis,
        "report_url": task.get("report_url", ""),
        "report_filename": task.get("report_filename", ""),
        "created_at": task.get("created_at", ""),
    }
    # Also persist batch tasks
    for bid, b in batch_tasks.items():
        h["batch_" + bid] = {
            "task_id": "batch_" + bid,
            "batch_id": bid,
            "company_name": f"批量完整分析 ({b.get('completed',0)}/{b.get('total',0)})",
            "status": b.get("status", ""),
            "grade": "-",
            "verdict": "-",
            "advance_recommendation": "-",
            "advance_reason": "-",
            "report_url": "",
            "report_filename": b.get("excel_path", ""),
            "results_path": str(REPORT_DIR / f"batch_{bid}.json"),
            "created_at": b.get("created_at", ""),
        }
    HISTORY_FILE.write_text(json.dumps(h, ensure_ascii=False), encoding="utf-8")

# Load history on startup
_history = _load_history()
for tid, t in _history.items():
    if tid.startswith("batch_"):
        batch_id = t.get("batch_id", tid.removeprefix("batch_"))
        persisted_path = Path(t.get("results_path", REPORT_DIR / f"batch_{batch_id}.json"))
        if persisted_path.exists():
            try:
                batch_tasks[batch_id] = json.loads(persisted_path.read_text(encoding="utf-8"))
            except json.JSONDecodeError:
                pass
    else:
        tasks[tid] = t
        tasks[tid]["analysis"] = {}  # placeholder, actual analysis lost on restart
print(f"[History] Loaded {len(_history)} records")


@app.get("/api/reports")
async def list_reports(page: int = 1, limit: int = 15):
    """列出报告，支持分页 + 统计"""
    items = []
    h = _load_history()
    for tid, t in tasks.items():
        h[tid] = {
            "task_id": tid,
            "company_name": t.get("company_name", ""),
            "status": t.get("status", "?"),
            "grade": (t.get("analysis", {}).get("grade") or t.get("grade", "-")) if isinstance(t.get("analysis"), dict) else t.get("grade", "-"),
            "verdict": (t.get("analysis", {}).get("verdict_type") or t.get("verdict", "-")) if isinstance(t.get("analysis"), dict) else t.get("verdict", "-"),
            "advance_recommendation": (t.get("analysis", {}).get("advance_recommendation") or t.get("advance_recommendation", "-")) if isinstance(t.get("analysis"), dict) else t.get("advance_recommendation", "-"),
            "advance_reason": (t.get("analysis", {}).get("advance_recommendation_reason") or t.get("advance_reason", "-")) if isinstance(t.get("analysis"), dict) else t.get("advance_reason", "-"),
            "total_score": (t.get("analysis", {}).get("total_score") or 0) if isinstance(t.get("analysis"), dict) else 0,
            "report_url": t.get("report_url", ""),
            "report_filename": t.get("report_filename", ""),
            "created_at": t.get("created_at", ""),
        }
    for tid, t in h.items():
        if not tid.startswith("batch_"):
            items.append(t)
    items.sort(key=lambda x: x.get("created_at", ""), reverse=True)

    # Stats
    grades = {}
    verdicts = {}
    scores = []
    for it in items:
        g = str(it.get("grade", "")).strip()
        if g: grades[g] = grades.get(g, 0) + 1
        v = str(it.get("verdict", "")).strip()
        if v: verdicts[v] = verdicts.get(v, 0) + 1
        s = it.get("total_score", 0)
        if isinstance(s, (int, float)) and s > 0: scores.append(s)
    avg_score = round(sum(scores) / len(scores), 1) if scores else 0

    # Paginate
    total = len(items)
    total_pages = max(1, (total + limit - 1) // limit)
    page = max(1, min(page, total_pages))
    start = (page - 1) * limit
    paged = items[start:start + limit]

    return {
        "reports": paged,
        "total": total,
        "page": page,
        "total_pages": total_pages,
        "stats": {
            "total": total,
            "avg_score": avg_score,
            "by_grade": grades,
            "by_verdict": verdicts,
        }
    }


@app.delete("/api/reports/{task_id}")
async def delete_report(task_id: str):
    """删除单个报告（HTML文件 + 历史记录）"""
    if task_id in tasks:
        t = tasks[task_id]
        rp = REPORT_DIR / (t.get("report_filename", "") or f"{task_id}.html")
        if rp.exists():
            rp.unlink()
        del tasks[task_id]
    h = _load_history()
    if task_id in h:
        del h[task_id]
        HISTORY_FILE.write_text(json.dumps(h, ensure_ascii=False), encoding="utf-8")
    return {"status": "deleted", "task_id": task_id}


@app.get("/api/batches")
async def list_batches():
    """列出所有批量完整分析任务"""
    items = []
    for bid, b in batch_tasks.items():
        items.append({
            "batch_id": bid,
            "status": b["status"],
            "total": b["total"],
            "completed": b["completed"],
            "created_at": b.get("created_at", ""),
            "excel_path": b.get("excel_path", ""),
        })
    items.sort(key=lambda x: x["created_at"], reverse=True)
    return {"batches": items}


# ============ 用量查询 ============
@app.get("/api/usage")
async def get_usage():
    """查询所有 Tavily Key 的官方用量，并提前切走接近耗尽的 Key。"""
    d = await _refresh_tavily_usage()
    _rotate_key_if_needed(reserve=20)
    d = _load_keys()
    keys_info = []
    for i, k in enumerate(d["keys"]):
        limit = int(k.get("credits_limit") or TAVILY_MONTHLY_LIMIT)
        used = int(k.get("credits_used", 0))
        pct = used / limit * 100 if limit else 0
        keys_info.append({
            "index": i,
            "label": k.get("label", ""),
            "credits_used": used,
            "limit": limit,
            "percent": round(pct, 1),
            "remaining": max(limit - used, 0),
            "is_current": i == d.get("current_index", 0),
            "status": k.get("usage_status", ""),
        })
    return {
        "total_keys": len(d["keys"]),
        "current_index": d.get("current_index", 0),
        "keys": keys_info,
    }


@app.delete("/api/tavily-keys/{index}")
async def delete_tavily_key(index: int, x_admin_token: str = Header("")):
    """删除指定 Key"""
    _require_admin_token(x_admin_token)
    d = _load_keys()
    if index < 0 or index >= len(d["keys"]):
        raise HTTPException(404, "Key不存在")
    if len(d["keys"]) <= 1:
        raise HTTPException(400, "至少保留一个 Key")
    removed = d["keys"].pop(index)
    if d["current_index"] >= len(d["keys"]):
        d["current_index"] = 0
    _save_keys(d)
    return {"status": "deleted", "removed": removed["label"]}


# ============ 启动 ============
def _sync_abcd_threshold_text(text: str) -> str:
    replacements = {
        "≥4.5→A，3.5-4.4→B，2.5-3.4→C，≤2.4→D": "≥3.8→A，3.0-3.7→B，2.0-2.9→C，＜2.0→D",
        ">=4.5->A,3.5-4.4->B,2.5-3.4->C,<=2.4->D": ">=3.8->A,3.0-3.7->B,2.0-2.9->C,<2.0->D",
        "4.5→A，3.5-4.4→B，2.5-3.4→C，≤2.4→D": "3.8→A，3.0-3.7→B，2.0-2.9→C，＜2.0→D",
        "4.5->A,3.5-4.4->B,2.5-3.4->C,2.4->D": "3.8->A,3.0-3.7->B,2.0-2.9->C,<2.0->D",
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    return text


if "SYSTEM_PROMPT_ANALYSIS" in globals():
    SYSTEM_PROMPT_ANALYSIS = _sync_abcd_threshold_text(SYSTEM_PROMPT_ANALYSIS)
if "SYSTEM_PROMPT_QUICK_SCORE" in globals():
    SYSTEM_PROMPT_QUICK_SCORE = _sync_abcd_threshold_text(SYSTEM_PROMPT_QUICK_SCORE)


if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8000))
    print(f"""
╔══════════════════════════════════════════╗
║   PE Deep Screening Web App - MVP        ║
║   http://localhost:{port}                  ║
╚══════════════════════════════════════════╝
    """)
    uvicorn.run(app, host="0.0.0.0", port=port)
